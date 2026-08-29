"""
Tools de acessibilidade expostas ao chat agentico via o registry local de tools.

O agente conversacional chama estas tools durante o tool-loop; cada uma embrulha
o pipeline deterministico existente (orchestrate). Registrado sob o toolset
'a11y_chat'.

O tool-loop multi-turn roda nos providers nativos suportados
(OpenAI/Anthropic/Gemini/xAI/Ollama Cloud), todos com suporte ao round-trip.
"""

import asyncio
import base64
import contextlib
import io
import json
import logging
import os
import re
import tempfile
import uuid
import zipfile
from collections.abc import Callable
from typing import Any

from backend.src.services import chat_progress

logger = logging.getLogger(__name__)

A11Y_CHAT_TOOLSET = "a11y_chat"
# Toolset separado porque o chat_runtime habilita 'clarify' explicitamente:
# é a pergunta interativa que as regras 12/13 do system prompt exigem.
CLARIFY_TOOLSET = "clarify"

_SEVERITY_DEDUCTION = {"critical": 20, "high": 10, "medium": 5, "low": 2}


def _safe_async_run(coro: Any) -> Any:
    """Executa uma corotina com segurança, criando thread isolada se já houver um loop ativo.

    Preserva o conversation_id corrente (session_context) quando precisa rodar a
    corotina numa thread worker. Sem isso, stores por sessão (last_fix_store,
    last_analysis_store, fix_checkpoint_store) leem/escritam no slot 'default' em
    vez da conversa atual, e ferramentas como open_live_preview nao encontram o
    preview gerado por fix_and_zip_files.
    """
    try:
        loop = asyncio.get_running_loop()
    except RuntimeError:
        loop = None

    if loop and loop.is_running():
        import concurrent.futures

        from backend.src.services.session_context import (
            reset_current_session,
            resolve_session,
            set_current_session,
        )

        session_id = resolve_session()

        def _run_in_thread() -> Any:
            token = set_current_session(session_id)
            try:
                return asyncio.run(coro)
            finally:
                reset_current_session(token)

        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            return pool.submit(_run_in_thread).result()
    return asyncio.run(coro)


def _summarize_issues(issues: list[dict[str, Any]]) -> dict[str, Any]:
    """Resumo compacto para o modelo narrar (sem estourar contexto)."""
    counts: dict[str, int] = {"critical": 0, "high": 0, "medium": 0, "low": 0}
    score = 100
    for issue in issues:
        # NAO usar str(...) aqui: severity pode vir como Severity(str, Enum) (model_dump
        # em modo python preserva o membro do enum, nao o valor). Enum.__str__ e sobrescrito
        # e retorna "Severity.LOW", quebrando o match contra `counts`. Como Severity extends
        # str, .lower() direto opera no conteudo real da string ("low"), ignorando __str__.
        sev = (issue.get("severity") or "").lower()
        if sev in counts:
            counts[sev] += 1
        score -= _SEVERITY_DEDUCTION.get(sev, 0)
    top = [
        {
            "criterion": i.get("criterion_pt") or i.get("criterion", ""),
            "severity": i.get("severity", ""),
            "element": (i.get("element", "") or "")[:120],
            "description": (i.get("description", "") or "")[:200],
        }
        for i in issues[:10]
    ]
    return {
        "total_issues": len(issues),
        "score": max(0, score),
        "counts_by_severity": counts,
        "top_issues": top,
    }


def _normalize_issue_for_fixer(issue: dict[str, Any], path: str) -> dict[str, Any]:
    """Completa campos mínimos de issues legados antes do fixer.

    Alguns agentes/integrações antigas podem devolver apenas ``criterion``,
    ``severity`` e ``description``. O pipeline de correção precisa de um
    AccessibilityIssue completo para gerar também a página corrigida do Live
    Preview; um issue parcial não deve fazer o arquivo inteiro cair no caminho
    de erro e desaparecer do preview.
    """
    normalized = dict(issue)
    normalized["url"] = path
    normalized.setdefault("guideline", "WCAG 2.2")
    normalized.setdefault(
        "suggestion",
        normalized.get("suggestion_technical")
        or normalized.get("description")
        or "Corrigir o elemento conforme o critério WCAG informado.",
    )
    normalized.setdefault("element", "documento HTML")
    normalized.setdefault("description", "Problema de acessibilidade identificado.")
    normalized.setdefault("severity", "medium")
    normalized.setdefault("criterion", "Critério WCAG não especificado")
    normalized.setdefault("id", f"issue-{uuid.uuid4().hex[:8]}")
    return normalized


async def _resolve_and_analyze(html: str, url: str, only_agents: list[str] | None = None) -> tuple[Any, str]:
    """Resolve o conteúdo (HTML direto ou via fetch de URL) e roda o pipeline.

    Devolve (resultado, content) -- o `content` resolvido (HTML semântico
    extraído, no caso de URL) é cacheado por `analyze_page` para uma correção
    subsequente na mesma conversa não ficar sem HTML para trabalhar (ver
    last_analyzed_content_store.py).
    """
    from backend.src.agents.orchestrator.orchestrator import orchestrate
    from backend.src.shared.models import TaskType

    content = html
    screenshot_base64 = None
    if url and not html:
        # Reaproveita o fetch com JS rendering + extracao semântica das rotas.
        from backend.src.routes.analyze import _extract_semantic_html
        from backend.src.services.browser import (
            fetch_accessibility_tree_snapshot,
            fetch_rendered_html_and_screenshot,
        )

        chat_progress.emit_tool_progress(None, "analyze_page", f"Buscando e renderizando {url}...")
        raw, screenshot_base64 = await fetch_rendered_html_and_screenshot(url)
        chat_progress.emit_tool_progress(None, "analyze_page", "Extraindo estrutura semântica...")
        # Best-effort: árvore de acessibilidade REAL (motor do navegador), nunca
        # derruba a análise se indisponível (ver fetch_accessibility_tree_snapshot).
        accessibility_tree = await fetch_accessibility_tree_snapshot(url)
        content = _extract_semantic_html(raw, accessibility_tree)

    chat_progress.emit_tool_progress(None, "analyze_page", "Executando especialistas em acessibilidade...")
    result = await orchestrate(content, TaskType.ANALYZE, screenshot_base64=screenshot_base64, only_agents=only_agents)
    chat_progress.emit_tool_progress(None, "analyze_page", "Consolidando resultados da análise...")
    return result, content


def analyze_page(args: dict[str, Any], **_kw: Any) -> str:
    """
    Handler da tool `analyze_page`. Aceita `html` (conteúdo direto) OU `url`
    (busca com JS rendering + extracao semântica). Roda o pipeline seletivo de
    especialistas, expert review e verificacao deterministica de contraste, e
    retorna um resumo JSON compacto para o agente narrar ao usuário.
    """
    html = str(args.get("html", "")).strip()
    url = str(args.get("url", "")).strip()
    append = bool(args.get("append", False))
    only_agents = args.get("only_agents")
    if not isinstance(only_agents, list) and only_agents:
        only_agents = [str(only_agents)]

    if not html and not url:
        return json.dumps({"error": "Forneca 'html' (conteúdo) ou 'url' da página a analisar."}, ensure_ascii=True)

    try:
        chat_progress.emit_tool_progress(None, "analyze_page", "Obtendo conteúdo da página...")
        result, resolved_content = _safe_async_run(_resolve_and_analyze(html, url, only_agents=only_agents))
    except Exception as exc:  # pragma: no cover - caminho de erro defensivo
        logger.error("[a11y_chat] analyze_page falhou: %s", exc)
        return json.dumps({"error": f"Falha na análise: {exc}"}, ensure_ascii=True)

    if not result.success:
        return json.dumps({"error": result.error or "Análise sem resultado."}, ensure_ascii=True)

    issues = result.data.get("issues", [])
    # Garante que cada issue tenha o campo 'url'
    for issue in issues:
        if not issue.get("url"):
            issue["url"] = url or "Conteúdo enviado por arquivo"

    from backend.src.services.last_analysis_store import set_last_analysis

    set_last_analysis(issues, url or "Conteúdo enviado por arquivo", append=append)

    # Cacheia o HTML resolvido (bruto se `html` foi passado direto, extraído
    # semanticamente se veio de `url`) -- sem isso, um pedido de correção
    # subsequente na mesma conversa não tinha HTML disponível para repassar a
    # `fix_and_zip_files` e acabava reanalisando a página do zero (achado real).
    from backend.src.services.last_analyzed_content_store import set_last_analyzed_content

    set_last_analyzed_content(resolved_content, url)

    summary = _summarize_issues(issues)

    # "Shift-right" sob demanda (decisão do usuário 2026-08-11): se essa MESMA
    # URL já foi analisada antes neste ambiente (qualquer sessão, qualquer
    # dia), compara com o snapshot anterior real e reporta regressão real
    # (issues novos) e melhoria real (issues resolvidos) -- nunca cria issue
    # no GitHub sozinho, só sinaliza pro modelo oferecer via create_github_issue.
    if url:
        from backend.src.services import url_scan_history_store

        previous = url_scan_history_store.get_previous_scan(url)
        if previous and previous.get("issues"):
            diff = url_scan_history_store.diff_scans(previous["issues"], issues)
            if diff["new_issues_count"] or diff["resolved_issues_count"]:
                summary["regression_vs_previous_scan"] = {
                    "previous_scan_at_unix": previous.get("scanned_at"),
                    "new_issues_count": diff["new_issues_count"],
                    "resolved_issues_count": diff["resolved_issues_count"],
                    "new_issues_preview": [
                        {
                            "criterion": i.get("criterion"),
                            "severity": i.get("severity"),
                            "description": i.get("description"),
                        }
                        for i in diff["new_issues"][:5]
                    ],
                }
        url_scan_history_store.save_scan(url, issues)

    return json.dumps(summary, ensure_ascii=True)


_ANALYZE_PAGE_SCHEMA: dict[str, Any] = {
    "description": (
        "Analyze the accessibility of a web page against WCAG 2.2, WAI-ARIA and "
        "Section 508. Runs a selective specialist pipeline plus deterministic contrast "
        "verification. Returns a compact summary: total issues, accessibility score "
        "(0-100), counts by severity, and the top issues. Call this whenever the user "
        "asks to analyze, audit, or check accessibility. Provide EITHER `url` (the page "
        "is fetched with JS rendering) OR `html` (raw content, e.g. an attached file or "
        "a whole project concatenated). After this tool succeeds, do not call web extraction "
        "for the same URL; use the returned summary and last-analysis cache. "
        "If this same URL was analyzed before (any past session), the response may include "
        "`regression_vs_previous_scan` -- real new issues that were not there last time, and "
        "real issues that got fixed since then. When present, narrate this to the user in "
        "natural language and offer to file a GitHub issue for the new regressions via "
        "`create_github_issue` (never create it without the user's approval, same as any "
        "effectful tool)."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Vou iniciar a análise de acessibilidade da página...').",
            },
            "url": {
                "type": "string",
                "description": "URL of the page to fetch and analyze (JS rendering).",
            },
            "html": {
                "type": "string",
                "description": "Raw HTML/code content to analyze (attached file or project).",
            },
            "append": {
                "type": "boolean",
                "description": "Set to true if you are analyzing multiple files in sequence and want to accumulate the results in the spreadsheet.",
            },
            "only_agents": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Optional list of specific accessibility specialist agent names to run (e.g., ['forms_a11y', 'css_analyzer', 'widgets_a11y', 'visual_a11y']). If omitted, the orchestrator automatically selects the relevant specialists.",
            },
        },
        "required": ["pre_exec_msg"],
    },
}

_ANALYZE_SITE_SCHEMA: dict[str, Any] = {
    "description": (
        "Analyze the accessibility of multiple pages or an entire website. "
        "Provide EITHER `url` (to crawl internal links starting from the root URL) "
        "OR `urls` (a list of specific URLs to analyze directly). "
        "Optional `max_pages` limits the crawler (default 10, max 50). "
        "Returns a consolidated summary: total issues, accessibility score (0-100), "
        "counts by severity, top issues, and audited page counts. "
        "Call this whenever the user wants to audit multiple URLs, a domain, or crawl a site."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Vou iniciar a varredura e análise de acessibilidade do site...').",
            },
            "url": {
                "type": "string",
                "description": "Root URL of the website to crawl and analyze.",
            },
            "urls": {
                "type": "array",
                "items": {
                    "type": "string",
                },
                "description": "A list of specific URLs to analyze directly (no crawling).",
            },
            "max_pages": {
                "type": "integer",
                "description": "Maximum number of pages to crawl. Default is 10, max is 50.",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


async def _run_site_crawl_and_analyze(
    url: str | None,
    urls: list[str] | None,
    max_pages: int,
) -> dict[str, Any]:
    """Crawl ou lista de URLs especificas, executa o pipeline em cada uma e consolida.

    Renderização das páginas e execução do pipeline de analise são paralelizadas
    (com semaforos) para reduzir a latencia total do site. Falhas isoladas nao
    interrompem as demais paginas.
    """
    from backend.src.agents.orchestrator.orchestrator import orchestrate
    from backend.src.config.settings import get_settings
    from backend.src.routes.analyze import _extract_semantic_html
    from backend.src.services.browser import fetch_rendered_html_and_screenshot
    from backend.src.services.crawler import crawl_site
    from backend.src.shared.models import TaskType

    # Limites de concorrencia: renderizacao de paginas (I/O de rede) e
    # analise (LLM/custo de API). Usamos a mesma configuracao max_concurrent
    # do orchestrator para analise, mas reservamos um semaforo separado.
    max_render = 5
    max_analysis = max(1, get_settings().a11y_max_concurrent_agents)
    render_sem = asyncio.Semaphore(max_render)
    analysis_sem = asyncio.Semaphore(max_analysis)

    target_pages: list[tuple[str, str]] = []

    chat_progress.emit_tool_progress(None, "analyze_site", "Coletando páginas para análise...")
    if urls:

        async def _render_one(u_clean: str) -> tuple[str, str] | None:
            async with render_sem:
                try:
                    chat_progress.emit_tool_progress(None, "analyze_site", f"Renderizando {u_clean}...")
                    raw, _ = await fetch_rendered_html_and_screenshot(u_clean)
                    return (u_clean, raw)
                except Exception as exc:
                    logger.error("[a11y_chat] Falha ao renderizar URL %s: %s", u_clean, exc)
                    return None

        render_tasks = [_render_one(u.strip()) for u in urls if u.strip()]
        target_pages = [t for t in await asyncio.gather(*render_tasks) if t is not None]
    elif url:
        # Crawling a partir do dominio raiz
        chat_progress.emit_tool_progress(None, "analyze_site", f"Varrendo links a partir de {url}...")
        page_results = await crawl_site(url, max_pages=max_pages)
        for page in page_results:
            if page.success and page.html:
                target_pages.append((page.url, page.html))

    if not target_pages:
        return {"error": "Não foi possivel acessar nenhuma das URLs fornecidas."}

    all_issues: list[dict[str, Any]] = []
    pages_ok = 0
    pages_failed = 0

    chat_progress.emit_tool_progress(None, "analyze_site", f"Auditando {len(target_pages)} página(s) em paralelo...")

    async def _analyze_one(page_url: str, html: str) -> tuple[bool, list[dict[str, Any]]]:
        async with analysis_sem:
            try:
                chat_progress.emit_tool_progress(None, "analyze_site", f"Analisando {page_url}...")
                semantic_html = _extract_semantic_html(html)
                res = await orchestrate(semantic_html, TaskType.ANALYZE)
                if res.success:
                    page_issues = res.data.get("issues", [])
                    for issue in page_issues:
                        issue["element"] = f"[{page_url}] {issue.get('element', '')}"
                    return (True, page_issues)
                return (False, [])
            except Exception as exc:
                logger.error("[a11y_chat] Erro ao analisar página %s: %s", page_url, exc)
                return (False, [])

    analysis_results = await asyncio.gather(*[_analyze_one(page_url, html) for page_url, html in target_pages])
    for ok, page_issues in analysis_results:
        if ok:
            pages_ok += 1
            all_issues.extend(page_issues)
        else:
            pages_failed += 1

    chat_progress.emit_tool_progress(None, "analyze_site", "Consolidando resultados do site...")

    # Cache na store
    from backend.src.services.last_analysis_store import set_last_analysis

    set_last_analysis(all_issues, f"Crawl/Lista de URLs: {url or ', '.join(urls or [])}")

    summary = _summarize_issues(all_issues)
    summary["total_pages"] = len(target_pages)
    summary["pages_ok"] = pages_ok
    summary["pages_failed"] = pages_failed
    return summary


def analyze_site(args: dict[str, Any], **_kw: Any) -> str:
    """
    Handler da tool `analyze_site`. Executa crawl ou analisa lista especifica
    de URLs, retornando o consolidado em formato JSON para o chat.
    """
    url = str(args.get("url", "")).strip()
    urls_raw = args.get("urls")
    urls = [str(u).strip() for u in urls_raw] if isinstance(urls_raw, list) else None
    max_pages = args.get("max_pages", 10)

    if not url and not urls:
        return json.dumps(
            {"error": "Forneca 'url' para crawl ou 'urls' para lista direta de páginas."}, ensure_ascii=True
        )

    try:
        max_pages = int(max_pages)
    except (ValueError, TypeError):
        max_pages = 10
    max_pages = max(1, min(max_pages, 50))

    try:
        result = _safe_async_run(_run_site_crawl_and_analyze(url or None, urls, max_pages))
        return json.dumps(result, ensure_ascii=True)
    except Exception as exc:
        logger.error("[a11y_chat] analyze_site falhou: %s", exc)
        return json.dumps({"error": f"Falha na análise de site: {exc}"}, ensure_ascii=True)


_UNZIP_AND_LIST_FILES_SCHEMA = {
    "description": (
        "Unzip a base64-encoded ZIP file containing local files/folders and return their paths and contents. "
        "Call this tool when the user uploads a zipped project/folder to audit. "
        "Returns a list of objects with `path` and `content` for each text file."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Vou extrair e listar os arquivos do projeto...').",
            },
            "zip_base64": {
                "type": "string",
                "description": "Base64 encoded string of the ZIP file.",
            },
        },
        "required": ["pre_exec_msg", "zip_base64"],
    },
}


def _read_docx_text(file_like) -> str:
    try:
        import docx

        doc = docx.Document(io.BytesIO(file_like.read()))
        lines = []
        for p in doc.paragraphs:
            if p.text.strip():
                style = p.style.name if p.style else "Normal"
                lines.append(f"[Style: {style}] {p.text}")
        for t in doc.tables:
            lines.append("[Table]")
            for row in t.rows:
                row_txt = " | ".join(cell.text.strip() for cell in row.cells)
                lines.append(f"| {row_txt} |")
        return "\n".join(lines)
    except Exception as e:
        logger.error("[a11y_chat] Erro ao ler DOCX: %s", e)
        return f"[Erro ao extrair texto do DOCX: {e}]"


def _run_verapdf_audit(pdf_path: str) -> list[str]:
    """Roda o validador oficial veraPDF CLI se estiver no PATH e retorna erros de PDF/UA."""
    import re
    import shutil
    import subprocess
    import xml.etree.ElementTree as ET

    if not shutil.which("verapdf"):
        return []

    try:
        # Executa o validador veraPDF no modo de validação PDF/UA (processor: rmd)
        res = subprocess.run(
            ["verapdf", "--format", "xml", "--processor", "rmd", pdf_path], capture_output=True, text=True, timeout=30
        )
        xml_content = res.stdout
        if not xml_content or "<report>" not in xml_content:
            return []

        # Limpa namespaces para facilitar o parsing com ET
        xml_content = re.sub(r'\sxmlns="[^"]+"', "", xml_content, count=1)
        xml_content = re.sub(r'\sxmlns:[^=]+="[^"]+"', "", xml_content)

        root = ET.fromstring(xml_content.encode("utf-8"))
        violations = []

        # Percorre as regras falhas (FAILED)
        for rule in root.findall(".//rule"):
            status_elem = rule.find("status")
            if status_elem is not None and status_elem.text == "FAILED":
                spec_elem = rule.find("specification")
                clause_elem = rule.find("clause")
                desc_elem = rule.find("description")
                test_num_elem = rule.find("testNumber")
                spec = spec_elem.text if spec_elem is not None else ""
                clause = clause_elem.text if clause_elem is not None else ""
                desc = desc_elem.text if desc_elem is not None else ""
                test_num = test_num_elem.text if test_num_elem is not None else ""

                rule_id = f"{spec} Cl. {clause}"
                if test_num:
                    rule_id += f" (Teste {test_num})"
                violations.append(f"- {rule_id}: {desc}")

        return violations
    except subprocess.TimeoutExpired:
        logger.warning("[a11y_chat] Tempo limite de execucao do veraPDF esgotado.")
        return []
    except Exception as e:
        logger.error("[a11y_chat] Erro ao rodar veraPDF: %s", e)
        return []


def _read_pdf_text(file_like) -> str:
    try:
        import os
        import tempfile

        import pypdf

        # Leitura dos bytes originais para pypdf e arquivo temporário (veraPDF)
        pdf_bytes = file_like.read()

        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))

        # Extração de campos de formulário nativos
        form_info = []
        try:
            fields = reader.get_fields()
            if fields:
                for field_name, field_value in fields.items():
                    field_type = field_value.get("/FT")
                    ft_map = {
                        "/Tx": "Campo de Texto",
                        "/Btn": "Caixa de Selecao/Opcao",
                        "/Ch": "Selecao (Dropdown)",
                        "/Sig": "Assinatura",
                    }
                    type_name = ft_map.get(field_type, "Campo")
                    form_info.append(f"- {field_name} ({type_name})")
        except Exception as fe:
            logger.warning("[a11y_chat] Não foi possivel extrair campos de formulário do PDF: %s", fe)

        # Execução do veraPDF via arquivo temporário
        verapdf_violations = []
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(pdf_bytes)
                tmp_path = tmp.name
            try:
                verapdf_violations = _run_verapdf_audit(tmp_path)
            finally:
                with contextlib.suppress(Exception):
                    os.remove(tmp_path)
        except Exception as ve:
            logger.warning("[a11y_chat] Não foi possivel rodar veraPDF: %s", ve)

        text_pages = []
        for idx, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            text_pages.append(f"--- Page {idx + 1} ---\n{page_text}")

        prefix = ""
        if form_info:
            prefix += "[Campos de Formulário Interativos Nativos Detectados no PDF]\n" + "\n".join(form_info) + "\n\n"
        if verapdf_violations:
            prefix += (
                "[Violacoes de Conformidade PDF/UA Detectadas pelo veraPDF]\n" + "\n".join(verapdf_violations) + "\n\n"
            )

        return prefix + "\n".join(text_pages)
    except Exception as e:
        logger.error("[a11y_chat] Erro ao ler PDF: %s", e)
        return f"[Erro ao extrair texto do PDF: {e}]"


def _read_xlsx_structure(file_like) -> str:
    """Extrai um resumo ESTRUTURAL (não o conteúdo célula a célula) de uma
    planilha XLSX real via openpyxl -- nome das abas, dimensões, painel
    congelado, faixas mescladas, cabeçalho detectado e imagens/gráficos
    embutidos -- para alimentar `run_excel_accessibility` (agente novo,
    ver backend/src/agents/excel_accessibility)."""
    try:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(file_like.read()), data_only=True)
        lines = [f"Planilha com {len(wb.sheetnames)} aba(s): {', '.join(wb.sheetnames)}\n"]

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"--- Aba '{sheet_name}' ---")
            lines.append(f"Dimensão: {ws.dimensions} ({ws.max_row} linhas x {ws.max_column} colunas)")
            lines.append(f"Painel congelado: {ws.freeze_panes or 'nenhum'}")

            merged = [str(r) for r in ws.merged_cells.ranges]
            lines.append(f"Células mescladas: {', '.join(merged) if merged else 'nenhuma'}")

            header_row = [c.value for c in next(ws.iter_rows(min_row=1, max_row=1), [])]
            has_header = any(isinstance(v, str) and v.strip() for v in header_row)
            lines.append(f"Linha 1 parece cabeçalho de texto: {has_header} (valores: {header_row[:10]})")

            try:
                # openpyxl não expõe API pública pra contar imagens embutidas nem
                # seu alt-text -- `_images` é atributo interno não tipado no stub
                # oficial, por isso o type: ignore; acesso já protegido por
                # try/except (biblioteca pode mudar esse detalhe interno sem aviso).
                image_count = len(ws._images)  # type: ignore[attr-defined]
            except Exception:
                image_count = 0
            lines.append(
                f"Imagens/gráficos embutidos nesta aba: {image_count} (openpyxl não expõe o alt-text definido no Excel -- reporte isso como incerteza, não como ausência confirmada)"
            )
            lines.append("")

        return "\n".join(lines)
    except Exception as e:
        logger.error("[a11y_chat] Erro ao ler estrutura do XLSX: %s", e)
        return f"[Erro ao extrair estrutura do XLSX: {e}]"


_ANALYZE_DOCUMENT_SCHEMA: dict[str, Any] = {
    "description": (
        "Analyze the REAL accessibility of an uploaded PDF or XLSX document (PDF/UA and "
        "Excel accessibility rules) via a dedicated specialist agent -- returns structured "
        "issues, same shape as analyze_page/analyze_site, NOT a blind fix. For PDF this "
        "also runs the real veraPDF PDF/UA validator when available on the server, plus "
        "native form field detection. Call this when the user uploads a PDF or XLSX and "
        "wants it analyzed/audited, before offering to fix it. "
        "IMPORTANT: do NOT pass the raw file bytes/base64 -- when a PDF or XLSX is "
        "uploaded, its content is already extracted into your own context as text right "
        "after the '=== filename ===' marker (labeled '[Texto extraído do PDF]' or "
        "'[Estrutura extraída do XLSX]'). Copy that exact extracted text into "
        "`document_text` here; do not re-derive, summarize, or guess it."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Analisando a acessibilidade do documento...').",
            },
            "filename": {
                "type": "string",
                "description": "Nome do arquivo enviado, usado para decidir o tipo (extensão .pdf ou .xlsx).",
            },
            "document_text": {
                "type": "string",
                "description": "O texto/estrutura JÁ EXTRAÍDO do documento, exatamente como apareceu no seu contexto após o upload (não os bytes brutos do arquivo).",
            },
        },
        "required": ["pre_exec_msg", "filename", "document_text"],
    },
}


def analyze_document(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `analyze_document`. Despacha PDF/XLSX pro agente
    especialista certo e cacheia o resultado em last_analysis_store, no mesmo
    formato de analyze_page -- export_xlsx/generate_checklist/generate_vpat
    funcionam em cima do resultado normalmente.

    Recebe o texto JÁ EXTRAÍDO (não bytes) -- mesma lição do Fix #1
    (fix_and_zip_files/last_analyzed_content_store): por essa altura da
    conversa o modelo não tem mais os bytes originais do upload, só o texto
    que o pré-processamento (chat_runtime.py::_preprocess_base64_attachments_impl)
    já colocou no contexto dele."""
    filename = str(args.get("filename", "")).strip()
    document_text = str(args.get("document_text", "")).strip()
    if not filename or not document_text:
        return json.dumps({"error": "Parâmetros 'filename' e 'document_text' são obrigatórios."}, ensure_ascii=True)

    ext = os.path.splitext(filename.lower())[1]
    if ext not in (".pdf", ".xlsx"):
        return json.dumps(
            {"error": f"Extensão '{ext}' não suportada por analyze_document (use .pdf ou .xlsx)."}, ensure_ascii=True
        )

    try:
        chat_progress.emit_tool_progress(None, "analyze_document", f"Analisando {filename}...")
        if ext == ".pdf":
            from backend.src.agents.pdf_accessibility.pdf_accessibility import run_pdf_accessibility

            result = _safe_async_run(run_pdf_accessibility(document_text))
        else:
            from backend.src.agents.excel_accessibility.excel_accessibility import run_excel_accessibility

            result = _safe_async_run(run_excel_accessibility(document_text))
    except Exception as exc:
        logger.error("[a11y_chat] analyze_document falhou: %s", exc)
        return json.dumps({"error": f"Falha ao analisar documento: {exc}"}, ensure_ascii=True)

    if not result.success:
        return json.dumps({"error": f"Falha ao analisar documento: {result.error}"}, ensure_ascii=True)

    issues = result.data.get("issues", [])
    from backend.src.services.last_analysis_store import set_last_analysis

    set_last_analysis(issues, filename)

    return json.dumps(
        _summarize_issues(issues) | {"filename": filename, "document_type": ext.lstrip(".")}, ensure_ascii=True
    )


def _read_epub_text(file_like) -> str:
    """Extrai texto estruturado de arquivos EPUB seguindo a ordem de leitura do spine."""
    try:
        import html
        import io
        import os
        import re
        import xml.etree.ElementTree as ET
        import zipfile
        from urllib.parse import unquote

        data = file_like.read() if hasattr(file_like, "read") else file_like

        with zipfile.ZipFile(io.BytesIO(data)) as epub:
            # Encontra o arquivo de manifesto (.opf)
            opf_name = None
            for name in epub.namelist():
                if name.endswith(".opf"):
                    opf_name = name
                    break

            if not opf_name:
                text_parts = []
                for name in sorted(epub.namelist()):
                    if name.endswith((".xhtml", ".html", ".htm")):
                        with epub.open(name) as f:
                            html_content = f.read().decode("utf-8", errors="ignore")
                            text_parts.append(re.sub(r"<[^>]+>", " ", html_content))
                return "\n\n".join(text_parts)

            # Lê o XML do OPF
            opf_dir = os.path.dirname(opf_name)
            with epub.open(opf_name) as f:
                opf_xml = f.read()

            root = ET.fromstring(opf_xml)

            # Mapeia os itens do manifesto por ID (independente de namespace)
            manifest = {}
            manifest_elem = None
            for child in root:
                if child.tag.endswith("manifest"):
                    manifest_elem = child
                    break

            if manifest_elem is not None:
                for item in manifest_elem:
                    if item.tag.endswith("item"):
                        item_id = item.get("id")
                        item_href = item.get("href")
                        if item_id and item_href:
                            manifest[item_id] = item_href

            # Lê os arquivos na ordem física do spine (independente de namespace)
            text_parts = []
            spine_elem = None
            for child in root:
                if child.tag.endswith("spine"):
                    spine_elem = child
                    break

            if spine_elem is not None:
                for itemref in spine_elem:
                    if itemref.tag.endswith("itemref"):
                        idref = itemref.get("idref")
                        href = manifest.get(idref) if idref is not None else None
                        if href:
                            rel_path = unquote(href)
                            full_path = os.path.normpath(os.path.join(opf_dir, rel_path)).replace("\\", "/")
                            if full_path in epub.namelist():
                                with epub.open(full_path) as f:
                                    html_content = f.read().decode("utf-8", errors="ignore")
                                    clean_text = re.sub(r"<[^>]+>", " ", html_content)
                                    clean_text = html.unescape(clean_text)
                                    clean_text = re.sub(r"\s+", " ", clean_text).strip()
                                    if clean_text:
                                        text_parts.append(clean_text)

            return "\n\n".join(text_parts)
    except Exception as e:
        logger.error("[a11y_chat] Erro ao ler EPUB: %s", e)
        return f"[Erro ao extrair texto do EPUB: {e}]"


def _read_pptx_text(file_like) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_like.read()))
        lines = []
        for idx, slide in enumerate(prs.slides):
            lines.append(f"--- Slide {idx + 1} ---")
            if slide.shapes.title:
                lines.append(f"[Title] {slide.shapes.title.text}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            lines.append(paragraph.text)
                if shape.has_table:
                    lines.append("[Table]")
                    for row in shape.table.rows:
                        row_txt = " | ".join(cell.text.strip() for cell in row.cells)
                        lines.append(f"| {row_txt} |")
        return "\n".join(lines)
    except Exception as e:
        logger.error("[a11y_chat] Erro ao ler PPTX: %s", e)
        return f"[Erro ao extrair texto do PPTX: {e}]"


async def _fix_docx_document(path: str, content_str: str, custom_instruction: str | None) -> bytes:
    import docx

    from backend.src.services.llm_client import call_llm, extract_json_object

    system_prompt = """
You are an expert accessibility document engineer. You will receive the text content and structure of a DOCX file.
Your task is to organize this content into a fully accessible semantic document structure, detecting its language and generating a descriptive title.
Fix any accessibility issues (e.g. ensure hierarchical headings, correct list structures, descriptive tables with headers).
Ensure the output text is written in correct Portuguese with proper grammar, spelling, and all standard accents (use á, é, í, ó, ú, â, ê, ô, ã, õ, ç). Do NOT strip or miss standard accents from words, as this breaks screen reader pronunciation.

Return a JSON object containing "metadata" (with "language" code like "pt-BR" or "en" and a descriptive document "title") and "blocks" (the document blocks in order):
{
  "metadata": {
    "language": "pt-BR",
    "title": "Document Title"
  },
  "blocks": [
    {"type": "heading", "level": 1, "text": "Document Title"},
    {"type": "paragraph", "text": "Paragraph content..."},
    {"type": "heading", "level": 2, "text": "Sub-section Title"},
    {"type": "list_item", "text": "Bullet item..."},
    {"type": "table", "headers": ["Header A", "Header B"], "rows": [["Cell A1", "Cell B1"]]}
  ]
}
Return ONLY the JSON object. Do NOT output markdown formatting outside the JSON block.
"""
    user_prompt = f"File: {path}\n\nContent:\n{content_str}"
    if custom_instruction:
        user_prompt += f"\n\nInstructions: {custom_instruction}"

    try:
        raw = await call_llm(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.1,
            model_tier="alto",
            agent_label="docx_fixer",
        )
        res = extract_json_object(raw)
        metadata = res.get("metadata", {}) if isinstance(res, dict) else {}
        blocks = res.get("blocks", []) if isinstance(res, dict) else []
        if not isinstance(blocks, list):
            blocks = []

        doc = docx.Document()

        # Configura propriedades de acessibilidade e metadados
        if metadata:
            doc.core_properties.title = metadata.get("title", "")
            doc.core_properties.language = metadata.get("language", "pt-BR")

        for b in blocks:
            b_type = b.get("type", "")
            if b_type == "heading":
                level = int(b.get("level", 1))
                doc.add_heading(b.get("text", ""), level=min(9, max(1, level)))
            elif b_type == "paragraph":
                doc.add_paragraph(b.get("text", ""))
            elif b_type == "list_item":
                doc.add_paragraph(b.get("text", ""), style="List Bullet")
            elif b_type == "table":
                headers = b.get("headers", [])
                rows = b.get("rows", [])
                if headers or rows:
                    cols_count = max(len(headers), max((len(r) for r in rows), default=0))
                    t = doc.add_table(rows=0, cols=cols_count)
                    t.style = "Table Grid"

                    from docx.oxml import parse_xml
                    from docx.oxml.ns import nsdecls

                    if headers:
                        row = t.add_row()
                        trPr = row._tr.get_or_add_trPr()
                        # Repete o cabecalho no topo de cada página (tblHeader)
                        trPr.append(parse_xml(r"<w:tblHeader {}/>".format(nsdecls("w"))))
                        # Impede que a linha do cabecalho se divida entre páginas (cantSplit)
                        trPr.append(parse_xml(r"<w:cantSplit {}/>".format(nsdecls("w"))))

                        hdr_cells = row.cells
                        for idx, text in enumerate(headers):
                            if idx < len(hdr_cells):
                                hdr_cells[idx].text = text
                    for r in rows:
                        row = t.add_row()
                        trPr = row._tr.get_or_add_trPr()
                        # Impede que as linhas de conteúdo se dividam entre páginas (cantSplit)
                        trPr.append(parse_xml(r"<w:cantSplit {}/>".format(nsdecls("w"))))

                        row_cells = row.cells
                        for idx, text in enumerate(r):
                            if idx < len(row_cells):
                                row_cells[idx].text = text

        out_stream = io.BytesIO()
        doc.save(out_stream)
        return out_stream.getvalue()
    except Exception as exc:
        logger.error("[a11y_chat] Erro ao corrigir DOCX %s: %s", path, exc)
        doc = docx.Document()
        doc.add_paragraph(content_str)
        out_stream = io.BytesIO()
        doc.save(out_stream)
        return out_stream.getvalue()


async def _fix_pptx_document(path: str, content_str: str, custom_instruction: str | None) -> bytes:
    from backend.src.services.llm_client import call_llm, extract_json_object

    system_prompt = """
You are an expert accessibility presentation engineer. You will receive the text content and structure of a PowerPoint (PPTX) file.
Your task is to remediate the presentation content to be highly accessible and compliant with Section 508 and WCAG 2.1/2.2 AA.
Ensure the output text is written in correct Portuguese with proper grammar, spelling, and all standard accents (use á, é, í, ó, ú, â, ê, ô, ã, õ, ç). Do NOT strip or miss standard accents from words, as this breaks screen reader pronunciation.
You must output a structured JSON containing slide-by-slide remediation:
1. Detect and declare the main presentation language (language: e.g. "pt-BR") and title (title: descriptive presentation title).
2. For each slide, define:
   - "title": A unique, descriptive slide title (mandatory under WCAG 2.4.2). If a slide repeats a topic, suffix it with " (Continued)" or " (1/2)".
   - "layout": Must be either "title_content" (standard for bullets/tables) or "title_only" (for graphs/images description).
   - "bullets": An array of strings representing bullet points. Keep lists structured, clear, and concise to avoid cognitive overload.
   - "table": An optional object containing "headers" (array of strings) and "rows" (array of arrays of strings).

Return a JSON object structured exactly like this:
{
  "metadata": {
    "title": "<presentation title>",
    "language": "<language code>"
  },
  "slides": [
    {
      "title": "<slide title>",
      "layout": "title_content",
      "bullets": ["bullet 1", "bullet 2"],
      "table": {
        "headers": ["Col 1", "Col 2"],
        "rows": [["cell 1", "cell 2"]]
      }
    }
  ]
}
Return ONLY valid JSON.
"""
    user_prompt = f"File: {path}\n\nContent:\n{content_str}"
    if custom_instruction:
        user_prompt += f"\n\nInstructions: {custom_instruction}"

    try:
        raw = await call_llm(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.1,
            model_tier="alto",
            agent_label="pptx_fixer",
        )
        res = extract_json_object(raw)
        metadata = res.get("metadata", {}) if isinstance(res, dict) else {}
        slides_data = res.get("slides", []) if isinstance(res, dict) else []
        if not isinstance(slides_data, list):
            slides_data = []

        from pptx import Presentation

        prs = Presentation()

        if metadata:
            prs.core_properties.title = metadata.get("title", "")

        for s_idx, s in enumerate(slides_data):
            title_text = s.get("title", f"Slide {s_idx + 1}")
            layout_type = s.get("layout", "title_content")

            if s_idx == 0:
                layout = prs.slide_layouts[0]
            elif layout_type == "title_only":
                layout = prs.slide_layouts[5]
            else:
                layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = title_text

            bullets = s.get("bullets", [])
            if bullets and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for b_idx, bullet in enumerate(bullets):
                    if b_idx == 0:
                        p = tf.paragraphs[0]
                        p.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                    p.level = 0

            table_data = s.get("table", {})
            if table_data and isinstance(table_data, dict):
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers or rows:
                    from pptx.util import Inches

                    cols_count = max(len(headers), max((len(r) for r in rows), default=0))
                    rows_count = len(rows) + (1 if headers else 0)

                    left = Inches(1.0)
                    top = Inches(2.0)
                    width = Inches(8.0)
                    height = Inches(0.8 + (rows_count * 0.4))

                    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
                    table = table_shape.table

                    current_row = 0
                    if headers:
                        for col_idx, text in enumerate(headers):
                            if col_idx < len(table.columns):
                                cell = table.cell(current_row, col_idx)
                                cell.text = text
                        current_row += 1
                    for r in rows:
                        for col_idx, text in enumerate(r):
                            if col_idx < len(table.columns):
                                cell = table.cell(current_row, col_idx)
                                cell.text = text
                        current_row += 1

        out_stream = io.BytesIO()
        prs.save(out_stream)
        return out_stream.getvalue()
    except Exception as exc:
        logger.error("[a11y_chat] Erro ao corrigir PPTX %s: %s", path, exc)
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Presentation Recovery"
        body = slide.placeholders[1]
        body.text = content_str
        out_stream = io.BytesIO()
        prs.save(out_stream)
        return out_stream.getvalue()


async def _fix_pdf_document_to_html(path: str, content_str: str, custom_instruction: str | None) -> str:
    from backend.src.services.llm_client import call_llm, extract_json_object

    system_prompt = """
You are an expert accessibility document engineer. You will receive the extracted text content of a PDF file.
Your task is to convert this content into a highly accessible, beautiful semantic HTML5 document equivalent to PDF/UA.
Ensure the output text is written in correct Portuguese with proper grammar, spelling, and all standard accents (use á, é, í, ó, ú, â, ê, ô, ã, õ, ç). Do NOT strip or miss standard accents from words, as this breaks screen reader pronunciation.
Ensure to:
1. Set the correct language code on the html element (e.g. <html lang="pt-BR">) and provide a descriptive <title> in the <head>.
2. Ensure proper heading hierarchy (no skipped heading levels) and structural landmarks (header, main, footer).
3. Ensure any tables are semantic: use <table>, <thead>, <tbody>, and <th> elements with explicit scope (scope="col" or scope="row").
4. Provide alt attributes for all images, and mark decorative icons/shapes with aria-hidden="true".
5. Add an attractive, modern, accessible stylesheet inline in the <style> block (dark mode compatible, minimum contrast 4.5:1, 16px font-size, visible focus outlines of at least 2px solid).

Return a JSON object:
{
  "fixed_html": "<corrected HTML string>",
  "changes_summary": ["<description of each change made>"]
}
Return ONLY valid JSON.
"""
    user_prompt = f"File: {path}\n\nContent:\n{content_str}"
    if custom_instruction:
        user_prompt += f"\n\nInstructions: {custom_instruction}"

    try:
        raw = await call_llm(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.1,
            model_tier="alto",
            agent_label="pdf_fixer",
        )
        res = extract_json_object(raw)
        return res.get("fixed_html") or f"<html><body><pre>{content_str}</pre></body></html>"
    except Exception as exc:
        logger.error("[a11y_chat] Erro ao converter PDF %s: %s", path, exc)
        return f"<html><body><pre>{content_str}</pre></body></html>"


async def _generate_alt_text_for_image(img_name: str, img_b64: str) -> str:
    from backend.src.services.llm_client import call_llm

    system_prompt = """
You are an expert web accessibility visual auditor. You will receive an image from a web page.
Generate a concise, descriptive alt text (alternative text) for this image to be used in an HTML <img> tag.
The alt text must describe the informative content of the image. Do NOT start with "image of" or "picture of".
If the image contains text, transcribe the text. If the image is purely decorative, return "decorative".
Return ONLY the alt text string. No formatting, no quotes, no markdown.
"""
    multimodal_prompt = [
        {"type": "text", "text": f"Generate alt text for this image named '{img_name}'."},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
    ]
    try:
        raw = await call_llm(
            system_prompt=system_prompt,
            user_prompt=multimodal_prompt,  # type: ignore[arg-type]  # call_llm repassa direto ao AIAgent
            temperature=0.1,
            model_tier="alto",
            agent_label="image_alt_generator",
        )
        return raw.strip().replace('"', "")
    except Exception as e:
        logger.error("[a11y_chat] Erro ao gerar alt text para imagem %s: %s", img_name, e)
        return "Descrição da imagem"


def unzip_and_list_files(args: dict[str, Any], **_kw: Any) -> str:
    zip_base64 = args.get("zip_base64", "").strip()
    if not zip_base64:
        return json.dumps({"error": "Parâmetro zip_base64 é obrigatório."}, ensure_ascii=True)

    try:
        chat_progress.emit_tool_progress(None, "unzip_and_list_files", "Decodificando o arquivo ZIP...")
        zip_data = base64.b64decode(zip_base64)
        extracted_files: list[dict[str, Any]] = []
        total_chars = 0
        max_chars = 100_000
        max_files = 30

        # Extensões válidas de arquivos de código/texto e documentos (incluindo linguagens mobile Swift, Kotlin, Dart, Java)
        valid_extensions = {
            ".html",
            ".htm",
            ".js",
            ".jsx",
            ".ts",
            ".tsx",
            ".css",
            ".json",
            ".md",
            ".docx",
            ".pdf",
            ".pptx",
            ".epub",
            ".swift",
            ".kt",
            ".dart",
            ".java",
        }
        image_extensions = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".svg"}

        with zipfile.ZipFile(io.BytesIO(zip_data)) as z:
            namelist = [n for n in z.namelist() if not n.endswith("/") and not n.endswith("\\") and "__MACOSX" not in n]
            for idx, name in enumerate(namelist):
                _, ext = os.path.splitext(name.lower())
                if ext not in valid_extensions and ext not in image_extensions:
                    continue

                if len(extracted_files) >= max_files:
                    logger.warning("[a11y_chat] Limite de %d arquivos atingido no unzip.", max_files)
                    break

                chat_progress.emit_tool_progress(
                    None, "unzip_and_list_files", f"Extraindo {name} ({idx + 1}/{len(namelist)})..."
                )
                try:
                    if ext in image_extensions:
                        img_bytes = z.read(name)
                        # Limite de 2MB por imagem
                        if len(img_bytes) <= 2 * 1024 * 1024:
                            img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                            extracted_files.append({"path": name, "content": img_b64, "is_image": True})
                    elif ext == ".docx":
                        content = _read_docx_text(z.open(name))
                        if total_chars + len(content) > max_chars:
                            content = content[: (max_chars - total_chars)]
                        total_chars += len(content)
                        if content:
                            extracted_files.append({"path": name, "content": content})
                    elif ext == ".pdf":
                        content = _read_pdf_text(z.open(name))
                        if total_chars + len(content) > max_chars:
                            content = content[: (max_chars - total_chars)]
                        total_chars += len(content)
                        if content:
                            extracted_files.append({"path": name, "content": content})
                    elif ext == ".pptx":
                        content = _read_pptx_text(z.open(name))
                        if total_chars + len(content) > max_chars:
                            content = content[: (max_chars - total_chars)]
                        total_chars += len(content)
                        if content:
                            extracted_files.append({"path": name, "content": content})
                    elif ext == ".epub":
                        content = _read_epub_text(z.open(name))
                        if total_chars + len(content) > max_chars:
                            content = content[: (max_chars - total_chars)]
                        total_chars += len(content)
                        if content:
                            extracted_files.append({"path": name, "content": content})
                    else:
                        with z.open(name) as f:
                            content = f.read().decode("utf-8", errors="ignore")

                        if total_chars + len(content) > max_chars:
                            logger.warning(
                                "[a11y_chat] Limite de caracteres (%d) excedido no unzip. Cortando arquivo %s",
                                max_chars,
                                name,
                            )
                            content = content[: (max_chars - total_chars)]
                            if content:
                                extracted_files.append({"path": name, "content": content})
                            break

                        total_chars += len(content)
                        extracted_files.append({"path": name, "content": content})
                except Exception as e:
                    logger.error("[a11y_chat] Falha ao extrair arquivo %s: %s", name, e)

        return json.dumps(extracted_files, ensure_ascii=True)
    except Exception as exc:
        logger.error("[a11y_chat] unzip_and_list_files falhou: %s", exc)
        return json.dumps({"error": f"Falha ao descompactar arquivo: {exc}"}, ensure_ascii=True)


_LOCAL_PROJECT_JUNK_DIRS = {
    ".git",
    "node_modules",
    "dist",
    "build",
    ".next",
    "__pycache__",
    "venv",
    ".venv",
    "target",
    ".idea",
    ".vscode",
    ".turbo",
    "coverage",
    ".cache",
    ".pytest_cache",
    "out",
}
_LOCAL_PROJECT_VALID_EXTENSIONS = {
    ".html",
    ".htm",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".vue",
    ".svelte",
    ".css",
    ".json",
    ".md",
    ".swift",
    ".kt",
    ".dart",
    ".java",
}
_LOCAL_PROJECT_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".svg"}
_LOCAL_PROJECT_MAX_FILES = 60
_LOCAL_PROJECT_MAX_CHARS = 150_000


def _resolve_safe_project_path(project_dir: str, relative_path: str) -> str | None:
    """Resolve `relative_path` dentro de `project_dir`, recusando qualquer
    caminho que escape do diretório (traversal via `../`, symlink, caminho
    absoluto embutido). Devolve None quando inseguro."""
    base = os.path.realpath(project_dir)
    candidate = os.path.realpath(os.path.join(base, relative_path))
    if candidate != base and not candidate.startswith(base + os.sep):
        return None
    return candidate


def _walk_local_project_files(project_dir: str) -> list[dict[str, Any]]:
    """Varre um diretório real no disco e devolve os mesmos shapes de
    `unzip_and_list_files` ({"path", "content", "is_image"?}), mas lendo do
    filesystem local em vez de um ZIP em memória. Caminhos retornados são
    relativos a `project_dir`, em estilo posix (compatível com o resto do
    pipeline de correção, que já assume `/` como separador)."""
    extracted_files: list[dict[str, Any]] = []
    total_chars = 0
    base = os.path.realpath(project_dir)
    file_count = 0

    for root, dirs, filenames in os.walk(base):
        dirs[:] = [d for d in dirs if d not in _LOCAL_PROJECT_JUNK_DIRS and not d.startswith(".")]
        for filename in sorted(filenames):
            if len(extracted_files) >= _LOCAL_PROJECT_MAX_FILES:
                logger.warning(
                    "[a11y_chat] Limite de %d arquivos atingido na leitura de projeto local.", _LOCAL_PROJECT_MAX_FILES
                )
                return extracted_files

            ext = os.path.splitext(filename.lower())[1]
            if ext not in _LOCAL_PROJECT_VALID_EXTENSIONS and ext not in _LOCAL_PROJECT_IMAGE_EXTENSIONS:
                continue

            file_count += 1
            abs_path = os.path.join(root, filename)
            rel_path = os.path.relpath(abs_path, base).replace(os.sep, "/")
            chat_progress.emit_tool_progress(None, "read_local_project_files", f"Lendo {rel_path}...")

            try:
                if ext in _LOCAL_PROJECT_IMAGE_EXTENSIONS:
                    with open(abs_path, "rb") as f:
                        img_bytes = f.read()
                    if len(img_bytes) <= 2 * 1024 * 1024:
                        extracted_files.append(
                            {
                                "path": rel_path,
                                "content": base64.b64encode(img_bytes).decode("utf-8"),
                                "is_image": True,
                            }
                        )
                else:
                    with open(abs_path, encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    if total_chars + len(content) > _LOCAL_PROJECT_MAX_CHARS:
                        content = content[: (_LOCAL_PROJECT_MAX_CHARS - total_chars)]
                        if content:
                            extracted_files.append({"path": rel_path, "content": content})
                        return extracted_files
                    total_chars += len(content)
                    extracted_files.append({"path": rel_path, "content": content})
            except Exception as exc:
                logger.error("[a11y_chat] Falha ao ler arquivo local %s: %s", abs_path, exc)

    return extracted_files


_READ_LOCAL_PROJECT_FILES_SCHEMA: dict[str, Any] = {
    "description": (
        "Read a real local project directory on the backend machine's disk (HTML, CSS, JS, TS, "
        "JSX, TSX, Vue, Svelte, Swift, Kotlin, Dart, Java, plus images) into your context, so you "
        "can analyze/chat about it the same way as an uploaded ZIP -- but without requiring the "
        "user to zip and upload it first. Call this when the user tells you a local path where "
        "their project lives (e.g. 'o projeto está em C:\\meuprojeto') and wants you to work with "
        "it directly. Never guess or invent a path -- only use a path the user explicitly gave you."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português avisando o usuário que a IA vai ler os arquivos desse diretório local antes de prosseguir.",
            },
            "project_dir": {
                "type": "string",
                "description": "Caminho absoluto de um diretório real no disco, exatamente como o usuário informou.",
            },
        },
        "required": ["pre_exec_msg", "project_dir"],
    },
}


def read_local_project_files(args: dict[str, Any], **_kw: Any) -> str:
    from backend.src.services.local_project_guard import (
        accessibility_scope_denial_message,
        is_accessibility_project_dir,
    )

    project_dir = str(args.get("project_dir") or "").strip()
    if not project_dir:
        return json.dumps({"error": "Parâmetro 'project_dir' é obrigatório."}, ensure_ascii=True)
    if not os.path.isdir(project_dir):
        return json.dumps({"error": f"Diretório não encontrado: {project_dir}"}, ensure_ascii=True)
    if not is_accessibility_project_dir(project_dir):
        return json.dumps({"error": accessibility_scope_denial_message(project_dir)}, ensure_ascii=True)

    try:
        chat_progress.emit_tool_progress(None, "read_local_project_files", f"Lendo arquivos de {project_dir}...")
        files = _walk_local_project_files(project_dir)
        if not files:
            return json.dumps(
                {
                    "error": "Nenhum arquivo com extensão reconhecida foi encontrado nesse diretório.",
                    "project_dir": project_dir,
                },
                ensure_ascii=True,
            )
        return json.dumps({"project_dir": project_dir, "files": files}, ensure_ascii=True)
    except Exception as exc:
        logger.error("[a11y_chat] read_local_project_files falhou: %s", exc)
        return json.dumps({"error": f"Falha ao ler o projeto local: {exc}"}, ensure_ascii=True)


_FIX_AND_ZIP_FILES_SCHEMA = {
    "description": (
        "Apply accessibility fixes to multiple local files (HTML, CSS, JS, TS, Swift, Kotlin, Dart, Java, DOCX, PDF) and pack them into a ZIP file for the user to download. "
        "Provide the list of files (with their paths and contents) and optional custom_instruction. "
        "IMPORTANT: if the user wants to fix the page most recently analyzed via `analyze_page` with a `url` (not an uploaded file/project), you do NOT have its HTML content in your own context -- do NOT re-call analyze_page to get it, and do NOT guess/reconstruct the HTML yourself. Simply omit `files` entirely (or pass an empty list): the page fetched by the last analyze_page(url=...) call in this conversation is used automatically as the single file to fix. "
        "Returns a download link for the generated ZIP and a summary of all changes made."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Iniciando correções automáticas de acessibilidade nos arquivos do projeto...').",
            },
            "files": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "content": {"type": "string"},
                    },
                    "required": ["path", "content"],
                },
                "description": "List of files with their paths and contents. May be omitted when fixing the page most recently analyzed via analyze_page(url=...) -- see tool description.",
            },
            "custom_instruction": {
                "type": "string",
                "description": "Optional custom instructions or guidelines for the fixes.",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


async def _run_non_html_code_fixer(path: str, content: str, custom_instruction: str | None) -> dict[str, Any]:
    from backend.src.services.llm_client import call_llm, extract_json_object

    system_prompt = """
You are an expert web and mobile accessibility engineer. You will receive a code file (CSS, JS, TS, JSX, TSX, Swift, Kotlin, Dart, Java, JSON, or Markdown) and an optional custom instruction.
Review the file for accessibility issues, for example:
- Web (HTML/CSS/JS): Color contrast variables, focus-visible outlines, keyboard event listeners alongside click handlers, ARIA states.
- iOS (Swift/SwiftUI): Missing accessibilityLabel, accessibilityHint, accessibilityValue, scaling dynamic type, scale metrics, or accessibilityElement(hidden: true).
- Android (Kotlin/Java/Jetpack Compose): Missing contentDescription, mergeDescendants, clickable elements without role/contentDescription, hardcoded text sizes instead of sp.
- Flutter (Dart): Missing Semantics widget, mergeSemantics, focusNode, or label/hint attributes.

Apply minimal, surgical fixes to solve accessibility issues while keeping all other logic, features, and styling intact. Do NOT add hardcoded keywords.
Ensure any output text, labels, or comments inside the corrected content are written in correct Portuguese with proper grammar, spelling, and all standard accents (use á, é, í, ó, ú, â, ê, ô, ã, õ, ç). Do NOT strip or miss standard accents from words, as this breaks screen reader pronunciation.
Return a JSON object:
{
  "fixed_content": "<corrected content string>",
  "changes_summary": ["<description of each change made>"]
}
Return ONLY valid JSON. No markdown, no preamble.
""".strip()

    user_prompt = f"File Path: {path}\n\nFile Content:\n{content}"
    if custom_instruction:
        user_prompt += f"\n\nCRITICAL CUSTOM INSTRUCTIONS FROM USER: {custom_instruction}"

    try:
        raw = await call_llm(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.1,
            model_tier="alto",
            agent_label="code_fixer",
        )
        res = extract_json_object(raw)
        return {
            "fixed_content": res.get("fixed_content") or content,
            "changes_summary": res.get("changes_summary") or [],
        }
    except Exception as exc:
        logger.error("[a11y_chat] Erro ao corrigir código não-HTML %s: %s", path, exc)
        return {"fixed_content": content, "changes_summary": []}


async def _render_html_to_screenshot(html_content: str) -> str:
    """Renderiza a string HTML e gera um screenshot base64 usando Playwright."""
    from playwright.async_api import async_playwright

    from backend.src.config.settings import get_settings

    settings = get_settings()
    ws_url = getattr(settings, "browserless_ws_url", None)
    user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"

    async with async_playwright() as pw:
        if ws_url:
            browser = await pw.chromium.connect_over_cdp(ws_url)
        else:
            browser = await pw.chromium.launch(headless=True)
        context = await browser.new_context(
            user_agent=user_agent,
            viewport={"width": 1280, "height": 800},
            java_script_enabled=True,
            ignore_https_errors=True,
        )
        page = await context.new_page()
        try:
            await page.set_content(html_content, timeout=15000)
            await page.wait_for_timeout(1000)
            screenshot_bytes = await page.screenshot(type="png", full_page=False)
            return base64.b64encode(screenshot_bytes).decode("utf-8")
        finally:
            await context.close()
            await browser.close()


async def _verify_layout_visually(original_html: str, fixed_html: str, screenshot_base64: str) -> dict[str, Any]:
    """Chama a LLM com o screenshot e o HTML para verificar se o layout foi quebrado."""
    from backend.src.services.llm_client import call_llm, extract_json_object

    system_prompt = """
You are an expert visual QA accessibility engineer. You will receive a screenshot and HTML content of a remediated web page.
Your task is to analyze the screenshot and check if the remediation (applied to fix accessibility issues) has broken, distorted, or deformed the visual design or layout.
Verify if:
1. Text is overlapping, clipped, or hidden.
2. Layout elements (like headers, sidebars, grids, forms) are severely misaligned or broken.
3. Contrast edits or styling overrides have made elements invisible or visually unappealing.

Return a JSON object:
{
  "layout_ok": true|false,
  "reasons": ["<brief description of each design bug found if layout_ok is false>"]
}
Return ONLY valid JSON.
""".strip()

    text_prompt = (
        "Check this screenshot of the fixed page. Check if the layout is intact or broken.\n\n"
        f"Remediated HTML (truncated):\n{fixed_html[:20000]}"
    )

    multimodal_prompt = [
        {"type": "text", "text": text_prompt},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{screenshot_base64}"}},
    ]

    try:
        raw = await call_llm(
            system_prompt=system_prompt,
            user_prompt=multimodal_prompt,  # type: ignore
            temperature=0.1,
            agent_label="layout_verifier",
            model_tier="alto",
        )
        res = extract_json_object(raw)
        return {"layout_ok": bool(res.get("layout_ok", True)), "reasons": res.get("reasons") or []}
    except Exception as exc:
        logger.error("[a11y_chat] Erro no visual layout verifier: %s", exc)
        return {"layout_ok": True, "reasons": []}


def _sanitize_accessible_links_and_labels(html_content: str) -> tuple[str, list[str]]:
    """Garante que nenhum link contenha URLs cruas como aria-label ou texto e
    aplica rotulagem semântica humanizada para redes sociais, parceiros e controles."""
    if not html_content or "<a" not in html_content:
        return html_content, []

    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html_content, "html.parser")
    modified = False
    notes: list[str] = []

    for a in soup.find_all("a"):
        href = str(a.get("href") or "").strip()
        target = str(a.get("target") or "").strip()
        aria_label = str(a.get("aria-label") or "").strip()
        text = a.get_text(strip=True)
        new_window_notice = " (abre em nova janela)" if target == "_blank" else ""
        href_lower = href.lower()

        # Remove aria-labels que contêm URLs cruas
        if aria_label.startswith(("http://", "https://", "Ir para http://", "Ir para https://")):
            # 1. Redes Sociais
            if "facebook.com" in href_lower:
                a["aria-label"] = f"Página do Facebook{new_window_notice}"
            elif "whatsapp.com" in href_lower:
                a["aria-label"] = f"Canal do WhatsApp{new_window_notice}"
            elif "twitter.com" in href_lower or "x.com" in href_lower:
                a["aria-label"] = f"Perfil do Twitter / X{new_window_notice}"
            elif "instagram.com" in href_lower:
                a["aria-label"] = f"Perfil do Instagram{new_window_notice}"
            elif "linkedin.com" in href_lower:
                a["aria-label"] = f"Página do LinkedIn{new_window_notice}"
            elif "t.me/" in href_lower or "telegram" in href_lower:
                a["aria-label"] = f"Canal do Telegram{new_window_notice}"
            elif "youtube.com" in href_lower:
                a["aria-label"] = f"Canal do YouTube{new_window_notice}"
            elif "github.com" in href_lower:
                a["aria-label"] = f"Repositório no GitHub{new_window_notice}"
            else:
                inner_img = a.find("img")
                if inner_img and inner_img.get("alt", "").strip():
                    a["aria-label"] = f"{inner_img.get('alt').strip()}{new_window_notice}"
                elif a.get("title", "").strip():
                    a["aria-label"] = f"{a.get('title').strip()}{new_window_notice}"
                else:
                    del a["aria-label"]
            modified = True
            notes.append(f"Rótulo acessível do link '{href[:30]}' humanizado para '{a.get('aria-label', '')}'.")

        # Links vazios sem texto nem rótulo
        elif not text and not aria_label:
            inner_img = a.find("img")
            if inner_img and inner_img.get("alt", "").strip():
                a["aria-label"] = f"{inner_img.get('alt').strip()}{new_window_notice}"
                modified = True
            elif a.get("title", "").strip():
                a["aria-label"] = f"{a.get('title').strip()}{new_window_notice}"
                modified = True
            elif "facebook.com" in href_lower:
                a["aria-label"] = f"Página do Facebook{new_window_notice}"
                modified = True
            elif "whatsapp.com" in href_lower:
                a["aria-label"] = f"Canal do WhatsApp{new_window_notice}"
                modified = True
            elif "twitter.com" in href_lower or "x.com" in href_lower:
                a["aria-label"] = f"Perfil do Twitter / X{new_window_notice}"
                modified = True
            elif "instagram.com" in href_lower:
                a["aria-label"] = f"Perfil do Instagram{new_window_notice}"
                modified = True
            elif "linkedin.com" in href_lower:
                a["aria-label"] = f"Página do LinkedIn{new_window_notice}"
                modified = True
            elif "t.me/" in href_lower or "telegram" in href_lower:
                a["aria-label"] = f"Canal do Telegram{new_window_notice}"
                modified = True
            elif "youtube.com" in href_lower:
                a["aria-label"] = f"Canal do YouTube{new_window_notice}"
                modified = True

    return (str(soup) if modified else html_content), notes


def _strip_injected_script_vectors(original_html: str, fixed_html: str) -> tuple[str, list[str]]:
    """Remove <script>/<iframe>/<object>/<embed> e atributos on*/javascript:/data:text/html
    que aparecem no HTML corrigido pela IA mas NÃO existiam no HTML original.

    Diferente de `_sanitize_preview_html` (que remove TODO script indiscriminadamente,
    porque o preview é só visual), aqui o objetivo é preservar scripts legítimos que já
    faziam parte do site original -- o ZIP é o entregável de produção, e removê-los
    quebraria a funcionalidade do site. Só bloqueia o que o fixer (LLM) introduziu de
    novo, que é exatamente o vetor de um prompt injection vindo do HTML analisado
    (ver SECURITY block em fixer.SYSTEM_PROMPT) conseguindo sobreviver até o download."""
    if not fixed_html:
        return fixed_html, []

    from bs4 import BeautifulSoup

    original_soup = BeautifulSoup(original_html or "", "html.parser")
    existing_tags = {str(tag) for tag in original_soup.find_all(["script", "iframe", "object", "embed"])}
    existing_attrs: set[tuple[str, str, str]] = set()
    for tag in original_soup.find_all(True):
        for attr, value in tag.attrs.items():
            value_str = value if isinstance(value, str) else " ".join(value or [])
            existing_attrs.add((tag.name, attr.lower(), value_str))

    fixed_soup = BeautifulSoup(fixed_html, "html.parser")
    notes: list[str] = []

    for tag in fixed_soup.find_all(["script", "iframe", "object", "embed"]):
        if str(tag) not in existing_tags:
            notes.append(f"Removida tag <{tag.name}> introduzida pela correcao (nao existia no HTML original).")
            tag.decompose()

    for tag in fixed_soup.find_all(True):
        for attr in list(tag.attrs):
            attr_lower = attr.lower()
            value = tag.attrs.get(attr)
            value_str = value if isinstance(value, str) else " ".join(value or [])
            is_event_handler = attr_lower.startswith("on")
            is_dangerous_uri = attr_lower in ("href", "src", "action") and value_str.strip().lower().startswith(
                ("javascript:", "data:text/html")
            )
            if (is_event_handler or is_dangerous_uri) and (tag.name, attr_lower, value_str) not in existing_attrs:
                notes.append(
                    f"Removido atributo '{attr}' introduzido pela correcao em <{tag.name}> (nao existia no HTML original)."
                )
                del tag.attrs[attr]

    return (str(fixed_soup) if notes else fixed_html), notes


async def _run_fixes_and_generate_zip(
    files: list[dict[str, Any]],
    custom_instruction: str | None,
    existing_issues: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    import posixpath

    from bs4 import BeautifulSoup

    from backend.src.agents.fixer.fixer import run_fixer
    from backend.src.agents.orchestrator.orchestrator import orchestrate
    from backend.src.services.session_context import resolve_session
    from backend.src.shared.models import AccessibilityIssue, TaskType

    session_id = resolve_session()
    logger.info(
        "[a11y_chat] _run_fixes_and_generate_zip iniciado (sessao %s, %d arquivo(s), %d issues pre-existentes)",
        session_id,
        len(files),
        len(existing_issues or []),
    )

    fixed_files = []
    changes_summary_all = []
    all_issues_accumulated = []
    preview_pages = []  # HTML pages {"title","original_html","fixed_html"} for open_live_preview

    # Mapeia todas as imagens do projeto
    project_images = {f["path"]: f["content"] for f in files if f.get("is_image")}

    def resolve_image_path(html_path: str, img_src: str) -> str:
        if img_src.startswith(("http:", "https:", "data:")):
            return ""
        html_dir = posixpath.dirname(html_path)
        resolved = posixpath.normpath(posixpath.join(html_dir, img_src))
        if resolved.startswith("./"):
            resolved = resolved[2:]
        if resolved.startswith("/"):
            resolved = resolved[1:]
        return resolved

    for idx, file_info in enumerate(files):
        path = file_info.get("path", "")
        content = file_info.get("content", "")

        chat_progress.emit_tool_progress(None, "fix_and_zip_files", f"Processando {path} ({idx + 1}/{len(files)})...")

        if file_info.get("is_image"):
            try:
                fixed_files.append({"path": path, "content": base64.b64decode(content)})
            except Exception as e:
                logger.error("[a11y_chat] Erro ao decodificar imagem %s: %s", path, e)
            continue

        ext = os.path.splitext(path.lower())[1]
        is_html = ext in (".html", ".htm")
        is_code = ext in (".css", ".js", ".jsx", ".ts", ".tsx", ".json", ".md", ".swift", ".kt", ".dart", ".java")
        is_docx = ext == ".docx"
        is_pdf = ext == ".pdf"
        is_pptx = ext == ".pptx"

        if is_html and content.strip():
            try:
                # Processamento Visual de Imagens para preencher alt text
                soup = BeautifulSoup(content, "html.parser")
                imgs = soup.find_all("img")
                # Fallback de URL simples: HTML pequeno, sem imagens, vindo do cache de analise
                # nao precisa do lento loop de validacao visual (Playwright + LLM multimodal).
                is_small_fallback = len(content) < 5000 and not imgs and existing_issues is not None
                html_modified = False
                for img in imgs:
                    if not img.has_attr("alt") or not str(img["alt"]).strip():
                        src = img.get("src", "").strip()
                        if src:
                            resolved_src = resolve_image_path(path, src)
                            matched_b64 = project_images.get(resolved_src)
                            if not matched_b64:
                                base_src = posixpath.basename(resolved_src)
                                for p, b64 in project_images.items():
                                    if posixpath.basename(p) == base_src:
                                        matched_b64 = b64
                                        break
                            if matched_b64:
                                alt_text = await _generate_alt_text_for_image(src, matched_b64)
                                img["alt"] = alt_text
                                html_modified = True
                                changes_summary_all.append(
                                    f"[{path}] Injetado alt='{alt_text}' para a imagem '{src}' gerado automaticamente pelo agente de visão computacional."
                                )

                if html_modified:
                    content = str(soup)

                # 1. Resolve issues: usa pre-existentes quando o HTML vem do fallback
                #    da análise (mesma conversa, mesma URL), evitando re-auditoria
                #    inconsistente que pode falhar em detectar os problemas e deixa
                #    o live preview sem páginas. Uploads de arquivo continuam sendo
                #    re-auditados pois nao ha cache de issues.
                file_issues_raw: list[dict[str, Any]] = []
                if existing_issues:
                    file_issues_raw = existing_issues
                    logger.info(
                        "[a11y_chat] Usando %d issue(s) pre-existente(s) do cache de analise para %s",
                        len(file_issues_raw),
                        path,
                    )
                else:
                    audit_res = await orchestrate(content, TaskType.ANALYZE)
                    if audit_res.success and audit_res.data.get("issues"):
                        file_issues_raw = audit_res.data["issues"]
                        logger.info(
                            "[a11y_chat] Re-auditoria encontrou %d issue(s) para %s",
                            len(file_issues_raw),
                            path,
                        )
                    else:
                        logger.info("[a11y_chat] Re-auditoria nao encontrou issues para %s", path)

                if file_issues_raw:
                    normalized_issues = [_normalize_issue_for_fixer(issue, path) for issue in file_issues_raw]
                    all_issues_accumulated.extend(normalized_issues)
                    file_issues = [AccessibilityIssue(**issue) for issue in normalized_issues]
                    # 2. Roda o fixer (fallback pequeno usa tier mais rapido)
                    fixer_model_tier = "fast" if is_small_fallback else "alto"
                    fix_res = await run_fixer(
                        content,
                        file_issues,
                        custom_instruction=custom_instruction,
                        model_tier=fixer_model_tier,
                    )
                    fixed_content = fix_res.data.get("fixed_html") or content
                    summary = fix_res.data.get("changes_summary", [])

                    # 3. Validação Visual (Visual Feedback Loop)
                    #    Pula para fallback pequeno para evitar latencia de Playwright + LLM multimodal.
                    if is_small_fallback:
                        chat_progress.emit(
                            {
                                "type": "phase",
                                "text": f"[{path}] Fallback simples detectado; pulando validação visual para resposta rápida.",
                            }
                        )
                    else:
                        try:
                            chat_progress.emit(
                                {
                                    "type": "phase",
                                    "text": f"[{path}] Renderizando página corrigida para validação visual...",
                                }
                            )
                            screenshot_b64 = await _render_html_to_screenshot(fixed_content)

                            chat_progress.emit(
                                {
                                    "type": "phase",
                                    "text": f"[{path}] Analisando layout da página com o validador de design...",
                                }
                            )
                            layout_res = await _verify_layout_visually(content, fixed_content, screenshot_b64)
                            if not layout_res.get("layout_ok", True):
                                reasons_str = "; ".join(layout_res.get("reasons", []))
                                chat_progress.emit(
                                    {
                                        "type": "phase",
                                        "text": f"[{path}] Quebra de layout detectada: {reasons_str}. Solicitando re-correção visual...",
                                    }
                                )
                                logger.warning(
                                    "[a11y_chat] Layout quebrado detectado pos-remediacao de %s: %s. Re-executando fixer com correcoes de design.",
                                    path,
                                    layout_res.get("reasons"),
                                )
                                extra_inst = f"CORRIGIR LAYOUT: O HTML gerado anteriormente quebrou o design visual pelos seguintes motivos: {reasons_str}. Reaplique as correções de acessibilidade mantendo o layout intacto."
                                if custom_instruction:
                                    extra_inst += f"\n\nInstrução original: {custom_instruction}"

                                fix_res_2 = await run_fixer(content, file_issues, custom_instruction=extra_inst)
                                fixed_content = fix_res_2.data.get("fixed_html") or fixed_content
                                summary.extend(fix_res_2.data.get("changes_summary", []))
                                summary.append(f"Ajuste visual aplicado para corrigir quebras de design: {reasons_str}")
                                chat_progress.emit({"type": "phase", "text": f"[{path}] Re-correção visual concluída."})
                            else:
                                chat_progress.emit({"type": "phase", "text": f"[{path}] Layout validado com sucesso!"})
                        except Exception as e_vis:
                            logger.error("[a11y_chat] Falha na validação visual de %s: %s", path, e_vis)
                            chat_progress.emit(
                                {
                                    "type": "phase",
                                    "text": f"[{path}] Falha ao executar validação visual. Continuando...",
                                }
                            )

                    # 4. Sanitização semântica de links e rótulos acessíveis (WCAG 2.4.4 / 4.1.2)
                    fixed_content, sanitize_notes = _sanitize_accessible_links_and_labels(fixed_content)
                    summary.extend(sanitize_notes)

                    # 5. Guardrail de seguranca: remove script/iframe/on*/javascript: que o
                    # fixer possa ter introduzido e nao existiam no HTML original (defesa em
                    # profundidade contra prompt injection sobrevivendo ate o ZIP baixado).
                    fixed_content, security_notes = _strip_injected_script_vectors(content, fixed_content)
                    if security_notes:
                        logger.warning(
                            "[a11y_chat] Vetores de execucao removidos do fix de %s (nao existiam no original): %s",
                            path,
                            security_notes,
                        )
                        summary.extend(security_notes)

                    fixed_files.append({"path": path, "content": fixed_content})
                    preview_pages.append({"title": path, "original_html": content, "fixed_html": fixed_content})
                    for s in summary:
                        changes_summary_all.append(f"[{path}] {s}")
                else:
                    fixed_files.append({"path": path, "content": content})
            except Exception as exc:
                logger.error("[a11y_chat] Erro ao corrigir %s: %s", path, exc)
                fixed_files.append({"path": path, "content": content})
        elif is_code and content.strip():
            try:
                code_res = await _run_non_html_code_fixer(path, content, custom_instruction)
                fixed_content = code_res.get("fixed_content") or content
                summary = code_res.get("changes_summary") or []
                fixed_files.append({"path": path, "content": fixed_content})
                for s in summary:
                    changes_summary_all.append(f"[{path}] {s}")
            except Exception as exc:
                logger.error("[a11y_chat] Erro ao processar código %s: %s", path, exc)
                fixed_files.append({"path": path, "content": content})
        elif is_docx and content.strip():
            try:
                docx_bytes = await _fix_docx_document(path, content, custom_instruction)
                fixed_files.append({"path": path, "content": docx_bytes})
                changes_summary_all.append(
                    f"[{path}] Documento DOCX estruturado de forma acessível com estilos de cabeçalho e tabelas semânticas."
                )
            except Exception as exc:
                logger.error("[a11y_chat] Erro ao corrigir DOCX %s: %s", path, exc)
                import docx

                doc = docx.Document()
                doc.add_paragraph(content)
                out_stream = io.BytesIO()
                doc.save(out_stream)
                fixed_files.append({"path": path, "content": out_stream.getvalue()})
        elif is_pdf and content.strip():
            try:
                pdf_html = await _fix_pdf_document_to_html(path, content, custom_instruction)
                html_path = os.path.splitext(path)[0] + ".html"
                fixed_files.append({"path": html_path, "content": pdf_html})
                changes_summary_all.append(
                    f"[{path}] Documento PDF convertido e estruturado para página da web HTML5 totalmente acessível: {html_path}"
                )
            except Exception as exc:
                logger.error("[a11y_chat] Erro ao converter PDF %s: %s", path, exc)
                fixed_files.append(
                    {"path": path + ".html", "content": f"<html><body><pre>{content}</pre></body></html>"}
                )
        elif is_pptx and content.strip():
            try:
                pptx_bytes = await _fix_pptx_document(path, content, custom_instruction)
                fixed_files.append({"path": path, "content": pptx_bytes})
                changes_summary_all.append(
                    f"[{path}] Apresentação PowerPoint PPTX reestruturada de forma acessível com layout de título/conteúdo nativo e tabela posicionada."
                )
            except Exception as exc:
                logger.error("[a11y_chat] Erro ao corrigir PPTX %s: %s", path, exc)
                from pptx import Presentation

                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if slide.shapes.title:
                    slide.shapes.title.text = "Presentation Recovery"
                body = slide.placeholders[1]
                body.text = content
                out_stream = io.BytesIO()
                prs.save(out_stream)
                fixed_files.append({"path": path, "content": out_stream.getvalue()})
        else:
            fixed_files.append({"path": path, "content": content})

    # Atualiza o cache global com todos os issues encontrados no projeto corrigido
    if all_issues_accumulated:
        from backend.src.services.last_analysis_store import set_last_analysis

        set_last_analysis(all_issues_accumulated, "Projeto ZIP corrigido", append=False)

    # Gera o arquivo ZIP
    unique_id = uuid.uuid4().hex[:12]
    zip_filename = f"qa-fixed-{unique_id}.zip"
    temp_dir = tempfile.gettempdir()
    exports_dir = os.path.join(temp_dir, "qa_accessibility_exports")
    os.makedirs(exports_dir, exist_ok=True)
    zip_filepath = os.path.join(exports_dir, zip_filename)

    with zipfile.ZipFile(zip_filepath, "w", zipfile.ZIP_DEFLATED) as z:
        for f in fixed_files:
            z.writestr(f["path"], f["content"])

    from backend.src.config.settings import get_settings

    download_url = f"{get_settings().resolved_public_base_url()}/export/download_zip/{zip_filename}"

    if preview_pages:
        from backend.src.services.last_fix_store import set_last_fix

        logger.info(
            "[a11y_chat] _run_fixes_and_generate_zip salvando %d pagina(s) de preview na sessao %s",
            len(preview_pages),
            session_id,
        )
        set_last_fix(preview_pages)
    else:
        logger.warning(
            "[a11y_chat] _run_fixes_and_generate_zip nao gerou paginas de preview (sessao %s). "
            "Possiveis causas: nenhum arquivo HTML, nenhum issue encontrado, ou erro no fixer.",
            session_id,
        )

    return {
        "download_url": download_url,
        "zip_filename": zip_filename,
        "changes_summary": changes_summary_all,
        "total_files": len(files),
        # Campo interno (não serializado ao modelo -- ver fix_and_zip_files/
        # fix_local_project_files, mesmo padrão de `_raw_violations` em
        # run_remote_test_tool): dá à ferramenta de escrita local acesso ao
        # conteúdo já corrigido sem duplicar o pipeline de correção inteiro.
        "_fixed_files": fixed_files,
    }


_MIN_PLAUSIBLE_FILE_CONTENT_CHARS = 30
_CODE_STRUCTURE_CHARS_RE = re.compile(r"[<>{};]")


def _looks_like_placeholder_file_list(files: list) -> bool:
    """Achado real (validação E2E 2026-08-10): mesmo instruído a OMITIR `files`
    quando não tem o HTML de verdade (ver schema de fix_and_zip_files), o
    modelo às vezes inventa um placeholder em vez de omitir (ex.: um único
    arquivo com content="---"). Heurística estrutural em DOIS sinais, não
    palavra-chave: curto (<30 chars) E sem nenhum caractere típico de
    marcação/código (`<>{};`). Só o tamanho sozinho gerou falso positivo real
    (achado durante a correção deste bug): um snippet de teste genuinamente
    curto tipo "<p>x</p>" (8 chars) tem marcação de verdade e não deve cair
    aqui -- só "---" e afins (curto E sem nenhuma estrutura de código)."""
    if len(files) != 1:
        return False
    content = str(files[0].get("content") or "").strip()
    return len(content) < _MIN_PLAUSIBLE_FILE_CONTENT_CHARS and not _CODE_STRUCTURE_CHARS_RE.search(content)


def fix_and_zip_files(args: dict[str, Any], **_kw: Any) -> str:
    files = args.get("files")
    if files is not None and not isinstance(files, list):
        return json.dumps({"error": "Parâmetro 'files' deve ser uma lista válida."}, ensure_ascii=True)

    if files and _looks_like_placeholder_file_list(files):
        logger.warning(
            "[a11y_chat] fix_and_zip_files recebeu um único arquivo com conteúdo "
            "implausivelmente curto (%r) -- tratando como se 'files' tivesse sido omitido.",
            files[0].get("content"),
        )
        files = None

    files_from_fallback = False
    if not files:
        from backend.src.services.last_analyzed_content_store import get_last_analyzed_content

        html, url = get_last_analyzed_content()
        if not html:
            return json.dumps(
                {
                    "error": (
                        "Nenhum arquivo foi fornecido e não há HTML de uma análise por URL recente "
                        "nesta conversa para usar como fallback. Peça ao usuário para enviar o(s) "
                        "arquivo(s) a corrigir, ou analise uma página por URL primeiro."
                    )
                },
                ensure_ascii=True,
            )
        path = (url.rstrip("/").rsplit("/", 1)[-1] or "index.html") if url else "index.html"
        ext = os.path.splitext(path.lower())[1]
        if "." not in path or ext not in (".html", ".htm"):
            path = "index.html"
        files = [{"path": path, "content": html}]
        files_from_fallback = True

    custom_instruction = args.get("custom_instruction")
    if custom_instruction:
        custom_instruction = str(custom_instruction).strip()

    # Quando nenhum arquivo foi enviado, usamos o HTML/URL da análise anterior.
    # Reaproveitamos os issues já detectados para evitar re-auditoria
    # inconsistente, que frequentemente deixa o live preview sem páginas.
    existing_issues = None
    if files_from_fallback:
        from backend.src.services.last_analysis_store import get_last_analysis

        issues_raw, _ = get_last_analysis()
        existing_issues = [dict(issue) for issue in issues_raw]
        logger.info(
            "[a11y_chat] fix_and_zip_files usando fallback de analise: %d issue(s) carregado(s)",
            len(existing_issues),
        )

    # Checkpoint ANTES de aplicar qualquer correção: a remediação sobrescreve o
    # cache da análise e as páginas de preview, e `undo_last_fix` precisa do
    # estado anterior para ter caminho de volta.
    from backend.src.services import fix_checkpoint_store

    fix_checkpoint_store.create_checkpoint(f"Correção de {len(files)} arquivo(s)")

    try:
        result = _safe_async_run(_run_fixes_and_generate_zip(files, custom_instruction, existing_issues))
        result.pop("_fixed_files", None)
        from backend.src.services.last_fix_store import get_last_fix
        from backend.src.services.session_context import resolve_session

        logger.info(
            "[a11y_chat] fix_and_zip_files concluido (sessao %s): total_files=%d, download_url=%s, "
            "preview_pages_no_cache=%s",
            resolve_session(),
            result.get("total_files", 0),
            result.get("download_url", ""),
            "sim" if get_last_fix() else "nao",
        )
        return json.dumps(result, ensure_ascii=True)
    except Exception as exc:
        logger.error("[a11y_chat] fix_and_zip_files falhou: %s", exc)
        return json.dumps({"error": f"Falha ao corrigir e zipar arquivos: {exc}"}, ensure_ascii=True)


_FIX_LOCAL_PROJECT_FILES_SCHEMA: dict[str, Any] = {
    "description": (
        "Apply accessibility fixes to a real local project directory on the backend machine's disk, "
        "and WRITE the corrected content back to the SAME files at that path (in addition to still "
        "generating a ZIP as a safety-net download). Use this instead of `fix_and_zip_files` when the "
        "user wants the fixes applied directly to their local project files (e.g. 'edita direto no "
        "projeto que está em C:\\meuprojeto'), not just a downloadable ZIP. If `files` is omitted, "
        "this tool re-reads the project directory from disk itself (same files `read_local_project_files` "
        "would return) -- you do not need to have called that tool first. A backup of every original file "
        "is written to a temp folder before overwriting, and its path is returned."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português avisando o usuário que a IA vai sobrescrever arquivos reais no diretório local informado (com backup).",
            },
            "project_dir": {
                "type": "string",
                "description": "Caminho absoluto de um diretório real no disco, exatamente como o usuário informou.",
            },
            "files": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "content": {"type": "string"},
                    },
                    "required": ["path", "content"],
                },
                "description": "Opcional: arquivos e conteúdo já obtidos via read_local_project_files. Se omitido, a ferramenta lê o diretório do disco de novo.",
            },
            "custom_instruction": {
                "type": "string",
                "description": "Optional custom instructions or guidelines for the fixes.",
            },
        },
        "required": ["pre_exec_msg", "project_dir"],
    },
}


def fix_local_project_files(args: dict[str, Any], **_kw: Any) -> str:
    from backend.src.services.local_project_guard import (
        accessibility_scope_denial_message,
        is_accessibility_project_dir,
    )

    project_dir = str(args.get("project_dir") or "").strip()
    if not project_dir:
        return json.dumps({"error": "Parâmetro 'project_dir' é obrigatório."}, ensure_ascii=True)
    if not os.path.isdir(project_dir):
        return json.dumps({"error": f"Diretório não encontrado: {project_dir}"}, ensure_ascii=True)
    if not is_accessibility_project_dir(project_dir):
        return json.dumps({"error": accessibility_scope_denial_message(project_dir)}, ensure_ascii=True)

    files = args.get("files")
    if files is not None and not isinstance(files, list):
        return json.dumps({"error": "Parâmetro 'files' deve ser uma lista válida."}, ensure_ascii=True)
    if not files:
        try:
            files = _walk_local_project_files(project_dir)
        except Exception as exc:
            return json.dumps({"error": f"Falha ao ler o projeto local: {exc}"}, ensure_ascii=True)
        if not files:
            return json.dumps(
                {"error": "Nenhum arquivo com extensão reconhecida foi encontrado nesse diretório."}, ensure_ascii=True
            )

    custom_instruction = args.get("custom_instruction")
    if custom_instruction:
        custom_instruction = str(custom_instruction).strip()

    from backend.src.services import fix_checkpoint_store

    fix_checkpoint_store.create_checkpoint(f"Correção local de {len(files)} arquivo(s) em {project_dir}")

    try:
        result = _safe_async_run(_run_fixes_and_generate_zip(files, custom_instruction))
    except Exception as exc:
        logger.error("[a11y_chat] fix_local_project_files falhou ao corrigir: %s", exc)
        return json.dumps({"error": f"Falha ao corrigir arquivos: {exc}"}, ensure_ascii=True)

    fixed_files = result.pop("_fixed_files", [])

    backup_dir = os.path.join(tempfile.gettempdir(), "qa_accessibility_local_backups", uuid.uuid4().hex[:12])
    written: list[str] = []
    skipped_unsafe: list[str] = []

    for f in fixed_files:
        rel_path = f.get("path", "")
        safe_path = _resolve_safe_project_path(project_dir, rel_path)
        if safe_path is None:
            logger.warning("[a11y_chat] fix_local_project_files recusou caminho fora do projeto: %r", rel_path)
            skipped_unsafe.append(rel_path)
            continue

        content = f.get("content", "")
        try:
            if os.path.exists(safe_path):
                backup_target = os.path.join(backup_dir, rel_path)
                os.makedirs(os.path.dirname(backup_target), exist_ok=True)
                with open(safe_path, "rb") as src, open(backup_target, "wb") as dst:
                    dst.write(src.read())

            os.makedirs(os.path.dirname(safe_path), exist_ok=True)
            if isinstance(content, (bytes, bytearray)):
                with open(safe_path, "wb") as out_bytes:
                    out_bytes.write(content)
            else:
                with open(safe_path, "w", encoding="utf-8") as out_text:
                    out_text.write(content)
            written.append(rel_path)
        except Exception as exc:
            logger.error("[a11y_chat] fix_local_project_files falhou ao escrever %s: %s", safe_path, exc)
            skipped_unsafe.append(rel_path)

    result["project_dir"] = project_dir
    result["written_files"] = written
    result["backup_dir"] = backup_dir if written else None
    if skipped_unsafe:
        result["skipped_files"] = skipped_unsafe

    return json.dumps(result, ensure_ascii=True)


_UNDO_LAST_FIX_SCHEMA: dict[str, Any] = {
    "description": (
        "Undo the last `fix_and_zip_files` run, restoring the accessibility analysis cache "
        "and the live-preview pages to the state they had before the fix. Call this when the "
        "user regrets the remediation, wants to go back, or asks to discard the applied fixes. "
        "Only the most recent fix can be undone."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Desfazendo a última correção e voltando ao estado anterior...').",
            }
        },
        "required": ["pre_exec_msg"],
    },
}


def undo_last_fix(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `undo_last_fix`: repõe o checkpoint anterior à correção."""
    from backend.src.services import fix_checkpoint_store

    checkpoint = fix_checkpoint_store.restore_checkpoint()
    if checkpoint is None:
        return json.dumps(
            {
                "error": "Não há correção para desfazer: nenhum checkpoint foi criado nesta conversa.",
            },
            ensure_ascii=False,
        )

    return json.dumps(
        {
            "restored": True,
            "checkpoint": checkpoint.describe(),
            "issues_restored": len(checkpoint.issues),
            "preview_pages_restored": len(checkpoint.fix_pages),
            "message": "Estado anterior à última correção restaurado.",
        },
        ensure_ascii=False,
    )


_EXPORT_XLSX_SCHEMA = {
    "description": (
        "Get the download link for the Excel spreadsheet containing the accessibility audit results "
        "of the last analyzed page or site. Call this when the user asks for a spreadsheet, "
        "an Excel file, a sheet, or to export the results."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando a planilha Excel com o relatório...').",
            }
        },
        "required": ["pre_exec_msg"],
    },
}


def export_xlsx(args: dict[str, Any], **_kw: Any) -> str:
    """
    Handler da tool `export_xlsx`. Retorna o link de download direto para a planilha.
    """
    from backend.src.config.settings import get_settings

    return json.dumps(
        {
            "download_url": f"{get_settings().resolved_public_base_url()}/export/last_xlsx",
            "message": "Excel spreadsheet download link is ready. The user can download it by clicking the link.",
        },
        ensure_ascii=True,
    )


def _issues_from_last_analysis() -> tuple[list, str]:
    """Loads the last analysis cache and converts raw issue dicts into real
    AccessibilityIssue instances -- vpat_reporter/test_generator access fields like
    `i.severity.value`, which only works on the Pydantic model, not a plain dict."""
    from backend.src.services.last_analysis_store import get_last_analysis
    from backend.src.shared.models import AccessibilityIssue

    issues_raw, url = get_last_analysis()
    issues = [AccessibilityIssue(**item) for item in issues_raw]
    return issues, url


_GENERATE_VPAT_SCHEMA = {
    "description": (
        "Generate a WCAG 2.2 VPAT (Voluntary Product Accessibility Template) -- the "
        "conformance document required for enterprise, government (Section 508), and "
        "procurement processes -- based on the accessibility issues found in the most "
        "recent analysis (analyze_page/analyze_site). Call this when the user asks for "
        "a VPAT, a conformance report, a Section 508 report, or procurement "
        "accessibility documentation."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando o VPAT WCAG 2.2 com base na última análise...').",
            },
            "product_name": {
                "type": "string",
                "description": "Nome do produto avaliado, se o usuário informou. Padrão: 'Produto Avaliado'.",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def generate_vpat(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `generate_vpat`. Usa os issues da última análise em cache."""
    from backend.src.agents.vpat_reporter.vpat_reporter import run_vpat_reporter

    try:
        issues, url = _issues_from_last_analysis()
    except Exception as exc:
        return json.dumps(
            {"error": f"Não foi possível interpretar os issues da última análise: {exc}"}, ensure_ascii=True
        )
    if not issues:
        return json.dumps(
            {"error": "Nenhuma análise recente encontrada. Analise uma página ou site antes de gerar o VPAT."},
            ensure_ascii=True,
        )

    product_name = args.get("product_name") or "Produto Avaliado"
    try:
        chat_progress.emit_tool_progress(None, "generate_vpat", "Montando o VPAT de conformidade...")
        result = _safe_async_run(run_vpat_reporter(issues=issues, target=url, product_name=product_name))
    except Exception as exc:
        logger.error("[a11y_chat] generate_vpat falhou: %s", exc)
        return json.dumps({"error": f"Falha ao gerar VPAT: {exc}"}, ensure_ascii=True)
    if not result.success:
        return json.dumps({"error": f"Falha ao gerar VPAT: {result.error}"}, ensure_ascii=True)
    return json.dumps({"vpat": result.data.get("vpat")}, ensure_ascii=False)


_GENERATE_CHECKLIST_SCHEMA = {
    "description": (
        "Generate a structured accessibility checklist (pass/fail/manual-verification items, "
        "one per WCAG criterion) from the issues found in the most recent analysis "
        "(analyze_page/analyze_site), via the dedicated ChecklistAgent -- NOT free text you write "
        "yourself. Call this whenever the user asks for a checklist. Returns the structured items; "
        "present them to the user as a checklist in your response. If the user also wants it as a "
        "downloadable file, follow up with `export_checklist_pdf`."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando o checklist de acessibilidade estruturado...').",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def generate_checklist(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `generate_checklist`. Usa o ChecklistAgent real (não texto
    solto do modelo) com os issues + HTML da última análise em cache."""
    from backend.src.agents.checklist.checklist import run_checklist
    from backend.src.services.last_analyzed_content_store import get_last_analyzed_content

    try:
        issues, url = _issues_from_last_analysis()
    except Exception as exc:
        return json.dumps(
            {"error": f"Não foi possível interpretar os issues da última análise: {exc}"}, ensure_ascii=True
        )
    if not issues:
        return json.dumps(
            {"error": "Nenhuma análise recente encontrada. Analise uma página ou site antes de gerar o checklist."},
            ensure_ascii=True,
        )

    html_content, _ = get_last_analyzed_content()
    try:
        chat_progress.emit_tool_progress(
            None, "generate_checklist", "Gerando o checklist estruturado de acessibilidade..."
        )
        result = _safe_async_run(run_checklist(issues, html_content=html_content or None))
    except Exception as exc:
        logger.error("[a11y_chat] generate_checklist falhou: %s", exc)
        return json.dumps({"error": f"Falha ao gerar checklist: {exc}"}, ensure_ascii=True)
    if not result.success:
        return json.dumps({"error": f"Falha ao gerar checklist: {result.error}"}, ensure_ascii=True)
    return json.dumps({"checklist": result.data.get("checklist"), "url": url}, ensure_ascii=False)


_EXPORT_CHECKLIST_PDF_SCHEMA = {
    "description": (
        "Get the download link for the accessibility checklist of the last analyzed page as a tagged, "
        "accessible PDF (PDF/UA-1 -- real structure tree, not a flat visual export). Call this when the "
        "user asks for the checklist as a PDF, a file, or a downloadable document."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando o PDF acessível do checklist...').",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def export_checklist_pdf_tool(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `export_checklist_pdf`. Retorna o link de download direto;
    o PDF é gerado sob demanda no GET (mesmo padrão de `export_xlsx`)."""
    from backend.src.config.settings import get_settings

    return json.dumps(
        {
            "download_url": f"{get_settings().resolved_public_base_url()}/export/last_checklist_pdf",
            "message": "Accessible checklist PDF (PDF/UA-1) download link is ready. The user can download it by clicking the link.",
        },
        ensure_ascii=True,
    )


_GENERATE_ACCESSIBILITY_STATEMENT_SCHEMA = {
    "description": (
        "Generate an Accessibility Statement (Declaração de Acessibilidade) -- the public "
        "conformance-status document organizations publish (e.g. at /accessibility), covering "
        "WCAG conformance target, evaluation methodology, known limitations (from the real issues "
        "found in the last analysis), and how users report a barrier. Call this when the user asks "
        "for an accessibility statement, a conformance statement, or wants to scope/document "
        "accessibility for a platform/consultancy engagement. Requires a prior analyze_page/analyze_site. "
        "Never invent the organization name or contact info if the user did not provide them -- pass "
        "only what the user actually said; the tool inserts clearly-marked placeholders otherwise."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando a Declaração de Acessibilidade com base na última análise...').",
            },
            "organization_name": {
                "type": "string",
                "description": "Nome real da organização, SE o usuário informou. Nunca inventar.",
            },
            "product_name": {
                "type": "string",
                "description": "Nome do produto/site avaliado, se informado. Padrão: 'Produto Avaliado'.",
            },
            "contact_email": {
                "type": "string",
                "description": "E-mail real de contato para reportar problemas de acessibilidade, SE o usuário informou. Nunca inventar.",
            },
            "contact_phone": {
                "type": "string",
                "description": "Telefone real de contato, SE o usuário informou. Nunca inventar.",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def generate_accessibility_statement(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `generate_accessibility_statement`. Monta a declaração
    real a partir dos issues da última análise em cache -- nunca texto solto
    inventado pelo modelo (ver accessibility_statement_generator.py)."""
    from backend.src.services.accessibility_statement_generator import build_accessibility_statement
    from backend.src.services.accessibility_statement_store import set_accessibility_statement_options
    from backend.src.services.last_analysis_store import get_last_analysis

    # Usa get_last_analysis() (dicts crus) e não _issues_from_last_analysis() (instâncias
    # AccessibilityIssue) -- build_accessibility_statement usa issue.get(...), que
    # AccessibilityIssue (pydantic) não tem; mesmo padrão da rota de export do PDF.
    issues, url = get_last_analysis()
    if not issues:
        return json.dumps(
            {
                "error": "Nenhuma análise recente encontrada. Analise uma página ou site antes de gerar a declaração de acessibilidade."
            },
            ensure_ascii=True,
        )

    options = {
        "organization_name": args.get("organization_name") or None,
        "product_name": args.get("product_name") or None,
        "contact_email": args.get("contact_email") or None,
        "contact_phone": args.get("contact_phone") or None,
    }
    set_accessibility_statement_options(options)

    chat_progress.emit_tool_progress(
        None, "generate_accessibility_statement", "Montando a declaração pública de acessibilidade..."
    )
    statement = build_accessibility_statement(issues, url, **options)
    return json.dumps(statement, ensure_ascii=False)


_EXPORT_ACCESSIBILITY_STATEMENT_PDF_SCHEMA = {
    "description": (
        "Get the download link for the Accessibility Statement of the last analyzed page as a "
        "tagged, accessible PDF (PDF/UA-1). Call this when the user asks for the accessibility "
        "statement as a PDF, a file, or a downloadable/publishable document. Call "
        "generate_accessibility_statement first if it has not been called yet in this conversation."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando o PDF acessível da Declaração de Acessibilidade...').",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def export_accessibility_statement_pdf_tool(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `export_accessibility_statement_pdf`. Retorna o link de
    download direto; o PDF é gerado sob demanda no GET (mesmo padrão de
    `export_checklist_pdf`)."""
    from backend.src.config.settings import get_settings

    return json.dumps(
        {
            "download_url": f"{get_settings().resolved_public_base_url()}/export/last_accessibility_statement_pdf",
            "message": "Accessible statement PDF (PDF/UA-1) download link is ready. The user can download it by clicking the link.",
        },
        ensure_ascii=True,
    )


_GENERATE_TEST_SUITE_SCHEMA = {
    "description": (
        "Generate a ready-to-run accessibility test suite (Playwright + axe-core) from "
        "the issues found in the most recent analysis (analyze_page/analyze_site), so "
        "the audited team's CI can catch regressions. Call this when the user asks for "
        "automated tests, a test suite, Playwright tests, or CI accessibility gates."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Gerando a suíte de testes Playwright + axe-core...').",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def generate_test_suite(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `generate_test_suite`. Usa os issues da última análise em cache."""
    from backend.src.agents.test_generator.test_generator import run_test_generator

    try:
        issues, url = _issues_from_last_analysis()
    except Exception as exc:
        return json.dumps(
            {"error": f"Não foi possível interpretar os issues da última análise: {exc}"}, ensure_ascii=True
        )
    if not issues:
        return json.dumps(
            {
                "error": "Nenhuma análise recente encontrada. Analise uma página ou site antes de gerar a suíte de testes."
            },
            ensure_ascii=True,
        )

    try:
        chat_progress.emit_tool_progress(
            None, "generate_test_suite", "Escrevendo os testes automatizados de acessibilidade..."
        )
        result = _safe_async_run(run_test_generator(issues=issues, target=url))
    except Exception as exc:
        logger.error("[a11y_chat] generate_test_suite falhou: %s", exc)
        return json.dumps({"error": f"Falha ao gerar suíte de testes: {exc}"}, ensure_ascii=True)
    if not result.success:
        return json.dumps({"error": f"Falha ao gerar suíte de testes: {result.error}"}, ensure_ascii=True)
    return json.dumps({"suite": result.data.get("suite")}, ensure_ascii=False)


_OPEN_LIVE_PREVIEW_SCHEMA = {
    "description": (
        "Open a Live Preview session showing the HTML pages fixed by the most recent "
        "fix_and_zip_files call, side-by-side (original vs. fixed with accessibility "
        "highlights). Call this when the user asks to see/visualize/preview the fixed "
        "page, or to compare before/after. Must be called AFTER fix_and_zip_files has "
        "produced at least one fixed HTML file in this conversation."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Abrindo a visualização ao vivo das páginas corrigidas...').",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def open_live_preview(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da tool `open_live_preview`. Cria uma sessão de preview a partir do cache da última correção."""
    from backend.src.routes.preview import register_preview_session
    from backend.src.services.last_fix_store import get_last_fix
    from backend.src.services.session_context import resolve_session

    chat_progress.emit_tool_progress(None, "open_live_preview", "Preparando a visualização antes e depois...")
    pages = get_last_fix()
    if not pages:
        session_id = resolve_session()
        logger.warning(
            "[a11y_chat] open_live_preview nao encontrou paginas de preview na sessao %s. "
            "Verifique se fix_and_zip_files foi chamado e se o contexto de sessao foi propagado.",
            session_id,
        )
        return json.dumps(
            {
                "error": "Nenhuma correção recente encontrada. Rode fix_and_zip_files antes de abrir a visualização ao vivo."
            },
            ensure_ascii=True,
        )

    session_id = register_preview_session(pages)
    return json.dumps({"session_id": session_id, "total_pages": len(pages)}, ensure_ascii=True)


_TAVILY_SEARCH_SCHEMA = {
    "description": (
        "Search the web using Tavily Search API. Returns a list of matching page titles, descriptions and URLs. "
        "Use this controlled search tool to find target pages for an accessibility audit, web accessibility guidelines, WCAG documentation, or accessibility articles. Do not use it after analyze_page already succeeded for the same URL unless the user asks for external sources."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Buscando informações adicionais na web via Tavily...').",
            },
            "query": {"type": "string", "description": "Search query strictly focused on accessibility"},
            "limit": {"type": "integer", "description": "Max results to return (1-10)", "default": 5},
        },
        "required": ["pre_exec_msg", "query"],
    },
}

_EXA_SEARCH_SCHEMA = {
    "description": (
        "Search the web using Exa.ai Search API. Returns a list of matching page titles, descriptions and URLs. "
        "Use this controlled search tool to find target pages for an accessibility audit, deep technical accessibility specs, WCAG ACT Rules, and accessibility solutions. Do not use it after analyze_page already succeeded for the same URL unless the user asks for external sources."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Buscando especificações técnicas detalhadas no Exa...').",
            },
            "query": {"type": "string", "description": "Search query strictly focused on accessibility"},
            "limit": {"type": "integer", "description": "Max results to return (1-10)", "default": 5},
        },
        "required": ["pre_exec_msg", "query"],
    },
}


def tavily_search(args: dict[str, Any], **_kw: Any) -> str:
    """Runs a web search using the Tavily Search API."""
    query = args.get("query")
    limit = args.get("limit") or 5
    if not query:
        return json.dumps({"error": "Query is required"})
    try:
        import httpx

        from backend.src.config.settings import get_settings

        settings = get_settings()
        api_key = settings.tavily_api_key or os.environ.get("TAVILY_API_KEY", "")
        if not api_key:
            return json.dumps({"error": "Tavily API key is not configured"})

        chat_progress.emit_tool_progress(None, "tavily_search", f"Pesquisando: {query}...")
        response = httpx.post(
            "https://api.tavily.com/search",
            json={"api_key": api_key, "query": query, "max_results": limit},
            timeout=15.0,
        )
        response.raise_for_status()
        data = response.json()

        results = []
        for r in data.get("results", []):
            results.append({"title": r.get("title", ""), "url": r.get("url", ""), "content": r.get("content", "")})
        return json.dumps({"success": True, "data": {"web": results}}, ensure_ascii=False)
    except Exception as exc:
        return json.dumps({"error": f"Tavily search failed: {str(exc)}"})


def exa_search(args: dict[str, Any], **_kw: Any) -> str:
    """Runs a web search using the Exa.ai Search API."""
    query = args.get("query")
    limit = args.get("limit") or 5
    if not query:
        return json.dumps({"error": "Query is required"})
    try:
        import httpx

        from backend.src.config.settings import get_settings

        settings = get_settings()
        api_key = settings.exa_api_key or os.environ.get("EXA_API_KEY", "")
        if not api_key:
            return json.dumps({"error": "Exa API key is not configured"})

        chat_progress.emit_tool_progress(None, "exa_search", f"Pesquisando: {query}...")
        response = httpx.post(
            "https://api.exa.ai/search",
            headers={"x-api-key": api_key, "content-type": "application/json"},
            json={"query": query, "num_results": limit, "use_autoprompt": True},
            timeout=15.0,
        )
        response.raise_for_status()
        data = response.json()

        results = []
        for r in data.get("results", []):
            results.append(
                {
                    "title": r.get("title", ""),
                    "url": r.get("url", ""),
                    "content": r.get("text", "") or r.get("highlights", [""])[0],
                }
            )
        return json.dumps({"success": True, "data": {"web": results}}, ensure_ascii=False)
    except Exception as exc:
        return json.dumps({"error": f"Exa search failed: {str(exc)}"})


_EVALUATE_RESEARCH_SCHEMA: dict[str, Any] = {
    "description": (
        "Use esta ferramenta após pesquisar com tavily_search ou exa_search para registrar "
        "o seu veredito de suficiência. Permite ao agente documentar se a informação recolhida "
        "é suficiente para responder com precisão normativa, ou se é necessária pesquisa adicional. "
        "Faz parte do ciclo Agentic RAG / ReAct 2026."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "verdict": {
                "type": "string",
                "enum": ["sufficient", "insufficient"],
                "description": "'sufficient' se a informação recolhida é suficiente para responder; 'insufficient' se falta informação crítica.",
            },
            "missing_info": {
                "type": "string",
                "description": "Descreve o que está em falta (apenas quando verdict='insufficient').",
            },
            "refined_query": {
                "type": "string",
                "description": "A query refinada para a próxima pesquisa (apenas quando verdict='insufficient').",
            },
            "sources_found": {
                "type": "string",
                "description": "Resumo das fontes normativas encontradas (quando verdict='sufficient').",
            },
        },
        "required": ["verdict"],
    },
}


def evaluate_research(args: dict) -> str:
    """Regista o veredito de suficiência de pesquisa do agente (Agentic RAG / ReAct)."""
    verdict = args.get("verdict", "insufficient")
    missing = args.get("missing_info", "")
    refined = args.get("refined_query", "")
    sources = args.get("sources_found", "")
    return json.dumps(
        {
            "verdict": verdict,
            "missing_info": missing,
            "refined_query": refined,
            "sources_found": sources,
            "action": "continue_search" if verdict == "insufficient" else "synthesize",
        },
        ensure_ascii=False,
    )


_CREATE_GITHUB_ISSUE_SCHEMA: dict[str, Any] = {
    "description": (
        "Cria uma nova Issue no repositório GitHub configurado contendo o diagnóstico "
        "e a sugestão de correção de um problema de acessibilidade encruzado."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Criando uma issue no GitHub com esse problema de acessibilidade...').",
            },
            "title": {
                "type": "string",
                "description": "Título descritivo da Issue (ex.: '[A11Y] WCAG 1.1.1 - Imagem sem alt').",
            },
            "body": {
                "type": "string",
                "description": "Descrição detalhada em Markdown com critério violado e recomendação.",
            },
            "repo_owner": {
                "type": "string",
                "description": "Dono do repositório (opcional, usa GITHUB_OWNER se omitido).",
            },
            "repo_name": {
                "type": "string",
                "description": "Nome do repositório (opcional, usa GITHUB_REPO se omitido).",
            },
        },
        "required": ["pre_exec_msg", "title", "body"],
    },
}


def create_github_issue_tool(args: dict) -> str:
    """Handler da ferramenta de criação de Issue no GitHub."""
    from backend.src.services.github_service import create_github_issue

    chat_progress.emit_tool_progress(None, "create_github_issue", "Abrindo a issue no GitHub...")
    res = create_github_issue(
        title=args.get("title", ""),
        body=args.get("body", ""),
        repo_owner=args.get("repo_owner"),
        repo_name=args.get("repo_name"),
    )
    return json.dumps(res, ensure_ascii=False)


_CREATE_JIRA_ISSUE_SCHEMA: dict[str, Any] = {
    "description": (
        "Cria uma nova Issue no Jira configurado contendo o diagnóstico e a sugestão de "
        "correção de um problema de acessibilidade encontrado."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Criando uma issue no Jira com esse problema de acessibilidade...').",
            },
            "summary": {
                "type": "string",
                "description": "Título/resumo da issue (ex.: '[A11Y] WCAG 4.1.2 - Botão sem nome acessível').",
            },
            "description": {
                "type": "string",
                "description": "Descrição detalhada em texto plano com critério violado e recomendação.",
            },
            "severity": {
                "type": "string",
                "description": "Severidade do issue ('critical'/'high'/'medium'/'low') -- mapeada para a prioridade do Jira.",
            },
            "project_key": {
                "type": "string",
                "description": "Chave do projeto Jira (opcional, usa JIRA_PROJECT_KEY se omitido).",
            },
        },
        "required": ["pre_exec_msg", "summary", "description"],
    },
}


def create_jira_issue_tool(args: dict) -> str:
    """Handler da ferramenta de criação de Issue no Jira."""
    from backend.src.services.ticket_integrations import create_jira_issue

    chat_progress.emit_tool_progress(None, "create_jira_issue", "Abrindo a issue no Jira...")
    res = create_jira_issue(
        summary=args.get("summary", ""),
        description=args.get("description", ""),
        severity=args.get("severity", "medium"),
        project_key=args.get("project_key"),
    )
    return json.dumps(res, ensure_ascii=False)


_CREATE_AZURE_DEVOPS_WORK_ITEM_SCHEMA: dict[str, Any] = {
    "description": (
        "Cria um novo Work Item no Azure DevOps configurado contendo o diagnóstico e a "
        "sugestão de correção de um problema de acessibilidade encontrado."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Criando um work item no Azure DevOps com esse problema de acessibilidade...').",
            },
            "title": {
                "type": "string",
                "description": "Título do work item (ex.: '[A11Y] WCAG 4.1.2 - Botão sem nome acessível').",
            },
            "description": {
                "type": "string",
                "description": "Descrição detalhada com critério violado e recomendação.",
            },
            "severity": {
                "type": "string",
                "description": "Severidade do issue ('critical'/'high'/'medium'/'low') -- mapeada para Microsoft.VSTS.Common.Severity.",
            },
            "project": {
                "type": "string",
                "description": "Projeto do Azure DevOps (opcional, usa AZURE_DEVOPS_PROJECT se omitido).",
            },
        },
        "required": ["pre_exec_msg", "title", "description"],
    },
}


def create_azure_devops_work_item_tool(args: dict) -> str:
    """Handler da ferramenta de criação de Work Item no Azure DevOps."""
    from backend.src.services.ticket_integrations import create_azure_devops_work_item

    chat_progress.emit_tool_progress(None, "create_azure_devops_work_item", "Criando o work item no Azure DevOps...")
    res = create_azure_devops_work_item(
        title=args.get("title", ""),
        description=args.get("description", ""),
        severity=args.get("severity", "medium"),
        project=args.get("project"),
    )
    return json.dumps(res, ensure_ascii=False)


_NVDA_SPEAK_SCHEMA: dict[str, Any] = {
    "description": (
        "Envia um comando de voz diretamente para o leitor de tela NVDA ativo para que ele fale um texto ao usuário."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Enviando o texto para o NVDA ler em voz alta...').",
            },
            "text": {
                "type": "string",
                "description": "Texto a ser falado pelo leitor de tela NVDA.",
            },
        },
        "required": ["pre_exec_msg", "text"],
    },
}


def nvda_speak_tool(args: dict) -> str:
    """Handler da ferramenta de controle de fala do NVDA."""
    from backend.src.services.nvda_service import speak_text

    chat_progress.emit_tool_progress(None, "nvda_speak", "Enviando texto para o NVDA...")
    res = speak_text(args.get("text", ""))
    return json.dumps(res, ensure_ascii=False)


_VERIFY_SCREEN_READER_ANNOUNCEMENTS_SCHEMA: dict[str, Any] = {
    "description": (
        "Verifica os anuncios de leitor de tela de uma URL cruzando a arvore de "
        "acessibilidade REAL computada pelo motor do navegador (Chromium/CDP -- a mesma "
        "API que NVDA/JAWS/Narrator consultam no Windows) contra regras deterministicas "
        "de nome acessivel ausente ou generico. Diferente de `analyze_page` (que estima a "
        "partir do HTML bruto via LLM), este achado e confirmado pelo proprio motor de "
        "acessibilidade do navegador. Se o NVDA real estiver rodando na maquina do usuario "
        "e `speak_via_nvda=true`, os achados sao lidos em voz alta para confirmacao humana. "
        "Requer BROWSERLESS_WS_URL configurado; sem isso, devolve lista vazia (nao falha)."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Verificando os anúncios de leitor de tela desta página...').",
            },
            "url": {
                "type": "string",
                "description": "URL da página a verificar.",
            },
            "speak_via_nvda": {
                "type": "boolean",
                "description": "Se true, e o NVDA real estiver rodando, le os achados em voz alta para confirmação humana.",
            },
        },
        "required": ["pre_exec_msg", "url"],
    },
}


def verify_screen_reader_announcements_tool(args: dict, **_kw: Any) -> str:
    """Handler da ferramenta de verificação de anúncios de leitor de tela."""
    from backend.src.services.screen_reader_verification import verify_screen_reader_announcements

    url = str(args.get("url") or "").strip()
    if not url:
        return json.dumps({"error": "Parâmetro 'url' é obrigatório."}, ensure_ascii=True)

    chat_progress.emit_tool_progress(
        None, "verify_screen_reader_announcements", "Verificando anúncios de leitor de tela..."
    )
    result = _safe_async_run(verify_screen_reader_announcements(url, speak_via_nvda=bool(args.get("speak_via_nvda"))))
    return json.dumps(
        {
            "url": result.url,
            "total_interactive_nodes": result.total_interactive_nodes,
            "findings": [
                {
                    "role": f.role,
                    "path": f.path,
                    "problem": f.problem,
                    "severity": f.severity,
                    "announcement_preview": f.announcement_preview,
                }
                for f in result.findings
            ],
            "nvda_running": result.nvda_running,
            "spoken_findings": result.spoken_findings,
        },
        ensure_ascii=False,
    )


_DESIGN_REVIEW_SCHEMA: dict[str, Any] = {
    "description": (
        "Anticipates accessibility risks from a requirement, user story, PRD excerpt, or "
        "component/flow description written in free text -- BEFORE any code exists. "
        "Different from every other analysis tool here, which audits HTML/code that already "
        "exists: this one reads the plan and flags concrete risks (with WCAG 2.2 criteria, "
        "severity, rationale tied to the actual requirement, and an actionable design "
        "recommendation) so the team can decide before building, not after. Call this when the "
        "user pastes a requirement/user story/PRD excerpt and asks for an accessibility review, "
        "or asks 'what should I watch out for' before building something."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português a ser exibida para o usuário antes de executar a ferramenta, justificando o que será feito (ex: 'Revisando o requisito em busca de riscos de acessibilidade...').",
            },
            "requirement_text": {
                "type": "string",
                "description": "O texto do requisito, user story, trecho de PRD ou descrição do componente/fluxo a revisar.",
            },
            "component_type": {
                "type": "string",
                "description": "Opcional: tipo de componente/fluxo já identificado (ex.: 'formulário multi-etapa', 'drag-and-drop', 'modal').",
            },
        },
        "required": ["pre_exec_msg", "requirement_text"],
    },
}


def design_review_tool(args: dict[str, Any], **_kw: Any) -> str:
    """Handler da ferramenta de revisão de acessibilidade pré-desenvolvimento (shift-left)."""
    from backend.src.agents.design_review.design_review import run_design_review

    requirement_text = str(args.get("requirement_text") or "").strip()
    if not requirement_text:
        return json.dumps({"error": "Parâmetro 'requirement_text' é obrigatório."}, ensure_ascii=True)

    chat_progress.emit_tool_progress(
        None, "design_review", "Revisando o requisito em busca de riscos de acessibilidade..."
    )
    result = _safe_async_run(run_design_review(requirement_text, args.get("component_type")))
    if not result.success:
        return json.dumps({"error": result.error or "Falha na revisão de design."}, ensure_ascii=True)
    return json.dumps(result.data, ensure_ascii=False)


_GENERATE_AUTOMATION_SCRIPT_SCHEMA: dict[str, Any] = {
    "description": (
        "Gera um script de teste de acessibilidade automatizado pronto para uso em "
        "Cypress (cypress-axe), Postman (Coleção JSON / Newman) ou Selenium (axe-selenium-python)."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "framework": {
                "type": "string",
                "enum": ["cypress", "postman", "selenium", "playwright"],
                "description": "O framework de automação desejado (cypress, postman, selenium ou playwright).",
            },
            "target_url": {
                "type": "string",
                "description": "A URL alvo a ser testada no script.",
            },
            "scope_selector": {
                "type": "string",
                "description": "Seletor CSS para restringir o escopo da auditoria (opcional).",
            },
        },
        "required": ["framework", "target_url"],
    },
}


def generate_automation_script_tool(args: dict) -> str:
    """Handler para gerar scripts de automação de testes de acessibilidade sob demanda."""
    framework = args.get("framework", "cypress").lower()
    url = args.get("target_url", "https://example.com")
    selector = args.get("scope_selector", "")

    if framework == "cypress":
        scope_arg = f", '{selector}'" if selector else ""
        script = f"""// cypress/e2e/a11y_spec.cy.js
import 'cypress-axe';

describe('Auditoria de Acessibilidade Cypress', () => {{
  beforeEach(() => {{
    cy.visit('{url}');
    cy.injectAxe();
  }});

  it('Deve passar na auditoria WCAG 2.2', () => {{
    cy.checkA11y({scope_arg if selector else "null"}, {{
      includedImpacts: ['critical', 'serious']
    }});
  }});
}});"""
    elif framework == "selenium":
        script = f"""# test_a11y.py (Selenium Python + axe-selenium-python)
import pytest
from selenium import webdriver
from axe_selenium_python import Axe

def test_accessibility_wcag():
    driver = webdriver.Chrome()
    try:
        driver.get('{url}')
        axe = Axe(driver)
        axe.inject()
        results = axe.run()
        violations = results.get('violations', [])
        assert len(violations) == 0, axe.report(violations)
    finally:
        driver.quit()
"""
    elif framework == "postman":
        script = """// Postman Test Script (Insira na aba 'Tests' da sua requisição)
pm.test("Status code é 200 OK", function () {
    pm.response.to.have.status(200);
});

pm.test("Score de Acessibilidade >= 80", function () {
    var jsonData = pm.response.json();
    pm.expect(jsonData.score).to.be.at.least(80);
});

pm.test("Zero violações críticas", function () {
    var jsonData = pm.response.json();
    var critical = jsonData.issues_by_severity ? jsonData.issues_by_severity.critical : 0;
    pm.expect(critical || 0).to.eql(0);
});"""
    else:  # playwright
        script = f"""// tests/a11y.spec.ts (Playwright)
import {{ test, expect }} from '@playwright/test';
import AxeBuilder from '@axe-core/playwright';

test('Auditoria WCAG 2.2 com Playwright', async ({{ page }}) => {{
  await page.goto('{url}');
  const accessibilityScanResults = await new AxeBuilder({{ page }})
    .withTags(['wcag2a', 'wcag2aa', 'wcag22aa'])
    .analyze();

  expect(accessibilityScanResults.violations).toEqual([]);
}});"""

    return json.dumps(
        {
            "status": "ok",
            "framework": framework,
            "script": script,
        },
        ensure_ascii=False,
    )


_RUN_REMOTE_TEST_SCHEMA: dict[str, Any] = {
    "description": (
        "Executa um teste de acessibilidade real. 'cypress' e 'selenium' rodam o axe-core de verdade "
        "(Deque Systems, o mesmo motor que cypress-axe/axe-selenium-python usam) contra a página "
        "renderizada, com violações WCAG reais. 'postman' roda uma collection Postman real via Newman "
        "quando 'npx' está disponível no servidor (fallback honesto para uma verificação leve de "
        "contrato de API quando não está). "
        "IMPORTANT for runner='cypress' or runner='selenium': there are real, independently valid "
        "places to run either -- 'local' (the actual Cypress/Selenium binary on the machine where this "
        "backend runs, checked for real -- it may already be installed there) or 'cloud' (axe-core, "
        "same detection engine, via remote Playwright/Browserless, no local install needed). This is "
        "NOT a fallback the model picks silently -- it is the USER's decision. If the user has not "
        "already stated a preference, you MUST ask them via `clarify` BEFORE calling this tool, and "
        "pass their choice as `location`. Local execution runs real commands on the backend machine, "
        "so ALWAYS ask for confirmation first with options like ['Executar localmente uma vez', "
        "'Sempre aprovar execução local nesta conversa', 'Rodar na nuvem']; if they pick 'sempre', pass "
        "`remember_choice: true` on this call so you don't have to ask again this conversation for the "
        "same runner. If 'local' is requested but the tool reports it is not actually installed there, "
        "it returns a clear error offering to install it for real (retry with location='install_local') "
        "-- ask the user first whether they want that (it can take a few minutes) or prefer cloud "
        "instead; never install anything without that explicit go-ahead."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "runner": {
                "type": "string",
                "enum": ["cypress", "postman", "selenium"],
                "description": "O executor desejado (cypress, postman ou selenium).",
            },
            "target_url": {
                "type": "string",
                "description": "A URL do site ou da API a ser testada.",
            },
            "scope_selector": {
                "type": "string",
                "description": "Seletor CSS para escopo (opcional para Cypress).",
            },
            "location": {
                "type": "string",
                "enum": ["local", "cloud", "install_local"],
                "description": (
                    "Onde executar (relevante/obrigatório para runner='cypress' ou runner='selenium'): "
                    "'local' roda o binário de verdade na máquina do backend, se já instalado; 'cloud' "
                    "roda via axe-core remoto; 'install_local' instala de verdade (só depois do usuário "
                    "confirmar explicitamente que quer isso, pode levar minutos). Pergunte ao usuário "
                    "antes de decidir -- não escolha sozinho."
                ),
            },
            "remember_choice": {
                "type": "boolean",
                "description": "True se o usuário escolheu 'sempre aprovar execução local nesta conversa' -- evita perguntar de novo pro mesmo runner neste chat.",
            },
            "local_project_dir": {
                "type": "string",
                "description": (
                    "Para runner='cypress' com location='local': caminho de qual instalação local "
                    "usar, SE a busca automática já encontrou mais de uma no computador do usuário e "
                    "retornou status='needs_selection' com a lista de candidatos. Pergunte ao usuário "
                    "qual das opções reais encontradas ele quer usar (via `clarify`, apresentando os "
                    "caminhos exatos retornados) e passe a escolha aqui. Não invente um caminho -- só "
                    "use um dos candidatos que a ferramenta já retornou."
                ),
            },
        },
        "required": ["runner", "target_url"],
    },
}


def run_remote_test_tool(args: dict) -> str:
    """Handler síncrono para despachar a execução remota de testes."""
    runner = args.get("runner", "cypress").lower()
    target_url = args.get("target_url", "")
    scope = args.get("scope_selector", "")
    location = args.get("location")
    remember_choice = bool(args.get("remember_choice"))
    local_project_dir = args.get("local_project_dir") or None

    from backend.src.services import local_exec_consent_store
    from backend.src.services.remote_runners import (
        run_remote_cypress_simulation,
        run_remote_postman_contract,
        run_remote_selenium,
    )

    if location in ("local", "install_local") and runner in ("cypress", "selenium") and remember_choice:
        local_exec_consent_store.set_consent(runner, "always")

    # Sem fallback silencioso pra cypress/selenium: é decisão do usuário, não
    # do modelo. Checagem estrutural (parâmetro ausente), não de palavra-chave.
    _needs_location_decision = runner in ("cypress", "selenium") and location not in ("local", "cloud", "install_local")

    if runner == "postman":
        chat_progress.emit_tool_progress(None, "run_remote_test", "Executando contrato Postman/Newman...")
        res = _safe_async_run(run_remote_postman_contract(target_url))
    elif _needs_location_decision:
        runner_label = "Cypress" if runner == "cypress" else "Selenium"
        res = {
            "error": (
                f"É preciso saber onde rodar o {runner_label} antes de executar: localmente (na "
                "máquina do backend, se já tiver o necessário instalado) ou na nuvem (axe-core "
                "remoto, sem instalação). Pergunte ao usuário qual prefere via `clarify` -- se ele "
                "escolher local, confirme se é 'só essa vez' ou 'sempre aprovar' -- e chame esta "
                "ferramenta de novo com `location` (e `remember_choice` se for o caso)."
            )
        }
    elif runner == "selenium":
        chat_progress.emit_tool_progress(None, "run_remote_test", f"Executando Selenium ({location})...")
        res = _safe_async_run(run_remote_selenium(target_url, location=location))
    else:
        chat_progress.emit_tool_progress(None, "run_remote_test", f"Executando Cypress ({location})...")
        res = _safe_async_run(
            run_remote_cypress_simulation(target_url, scope, location=location, project_dir_override=local_project_dir)
        )

    raw_violations = res.pop("_raw_violations", None)
    if raw_violations is not None and res.get("status") == "ok" and target_url:
        _cache_remote_test_result_for_deliverables(runner, target_url, raw_violations)

    return json.dumps(res, ensure_ascii=False)


_RUN_CROSS_BROWSER_TEST_SCHEMA: dict[str, Any] = {
    "description": (
        "Run a REAL axe-core accessibility audit against a URL across the 3 real browser "
        "rendering engines Playwright ships (Chromium, Firefox/Gecko, WebKit) -- not just one. "
        "Screen reader and accessibility-tree behavior genuinely differs between engines (e.g. "
        "WebKit is the actual engine VoiceOver on macOS/iOS relies on; Firefox/Gecko has its own "
        "accessibility tree implementation with its own quirks). Use this when the user wants to "
        "know whether an issue is specific to one browser engine, or wants broader real coverage "
        "than a single-engine test. Returns a per-engine summary plus a diff of which violations "
        "appear only in some engines, not all -- real cross-browser behavior differences, not "
        "guessed. Feeds the same deliverable pipeline (spreadsheet/checklist/PDF) as analyze_page."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português avisando que a IA vai rodar o teste real nos 3 motores de navegador.",
            },
            "target_url": {
                "type": "string",
                "description": "URL real a testar.",
            },
        },
        "required": ["pre_exec_msg", "target_url"],
    },
}


def run_cross_browser_test_tool(args: dict) -> str:
    """Handler síncrono: roda axe-core real nos 3 motores (Chromium/Firefox/WebKit)
    e devolve resumo por motor + diff das diferenças reais entre eles."""

    from backend.src.services.browser import diff_cross_browser_violations, run_axe_core_cross_browser_audit
    from backend.src.services.remote_runners import _summarize_axe_results

    target_url = args.get("target_url", "")
    if not target_url:
        return json.dumps({"error": "Parâmetro 'target_url' é obrigatório."}, ensure_ascii=True)

    try:
        chat_progress.emit_tool_progress(
            None, "run_cross_browser_test", f"Testando {target_url} nos 3 motores de navegador..."
        )
        raw = _safe_async_run(run_axe_core_cross_browser_audit(target_url))
    except Exception as exc:
        logger.error("[a11y_chat] run_cross_browser_test_tool falhou: %s", exc)
        return json.dumps({"error": f"Falha ao rodar teste cross-browser: {exc}"}, ensure_ascii=True)

    per_engine = raw.get("per_engine", {})
    diff = diff_cross_browser_violations(per_engine)

    per_engine_summary: dict[str, Any] = {}
    all_violations_union: dict[str, dict] = {}
    for engine, data in per_engine.items():
        if data.get("success"):
            summary = _summarize_axe_results(data["results"])
            for v in summary.pop("_raw_violations", []):
                all_violations_union.setdefault(v.get("id", ""), v)
            per_engine_summary[engine] = summary
        else:
            per_engine_summary[engine] = {"error": data.get("error")}

    result = {
        "status": "ok" if diff["engines_succeeded"] else "error",
        "url": target_url,
        "engines_succeeded": diff["engines_succeeded"],
        "engines_failed": diff["engines_failed"],
        "per_engine_summary": per_engine_summary,
        "cross_browser_only_differences": diff["cross_browser_only_differences"],
    }

    # Achado real (2026-08-11, pedido do usuário): quando um motor falha por
    # binário ausente ("playwright install" na mensagem de erro do próprio
    # Playwright), avisa a IA pra oferecer a instalação em vez de só reportar
    # o erro cru -- o usuário decide via `install_playwright_browsers` (tool
    # separada, requer confirmação explícita antes de baixar qualquer coisa).
    missing_engines = [
        engine
        for engine, data in per_engine_summary.items()
        if isinstance(data.get("error"), str) and "playwright install" in data["error"].lower()
    ]
    if missing_engines:
        result["install_suggestion"] = (
            f"Os motores {missing_engines} não estão instalados nesta máquina. "
            "Pergunte ao usuário se quer que você instale agora (via "
            "`install_playwright_browsers`) para rodar o teste real nesses motores também."
        )

    union_violations = list(all_violations_union.values())
    if union_violations and diff["engines_succeeded"]:
        _cache_remote_test_result_for_deliverables("cross_browser_axe", target_url, union_violations)

    return json.dumps(result, ensure_ascii=True)


_INSTALL_PLAYWRIGHT_BROWSERS_SCHEMA: dict[str, Any] = {
    "description": (
        "Install the real Playwright browser binaries (Chromium, Firefox, WebKit) on the machine "
        "running this backend, via `python -m playwright install`. ONLY call this after the user has "
        "explicitly confirmed in chat that they want it installed (e.g. after run_cross_browser_test_tool "
        "returned an `install_suggestion`) -- this downloads real software onto the user's/server's disk, "
        "so it must never run without an explicit yes from the user in THIS conversation."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "pre_exec_msg": {
                "type": "string",
                "description": "Uma mensagem curta em português avisando que a IA vai instalar os navegadores do Playwright agora, confirmado pelo usuário.",
            },
        },
        "required": ["pre_exec_msg"],
    },
}


def install_playwright_browsers_tool(args: dict, **_kw: Any) -> str:
    """Handler síncrono: instala os binários reais do Playwright (Chromium/
    Firefox/WebKit) via subprocess. Só deve ser chamada após confirmação
    explícita do usuário no chat (ver schema) -- nunca roda sozinha."""
    import subprocess
    import sys

    chat_progress.emit_tool_progress(None, "install_playwright_browsers", "Instalando os navegadores do Playwright...")
    try:
        proc = subprocess.run(
            [sys.executable, "-m", "playwright", "install", "chromium", "firefox", "webkit"],
            capture_output=True,
            text=True,
            timeout=600,
        )
        if proc.returncode != 0:
            logger.error("[a11y_chat] install_playwright_browsers falhou: %s", proc.stderr[-2000:])
            return json.dumps(
                {
                    "status": "error",
                    "error": f"Falha ao instalar os navegadores do Playwright: {proc.stderr[-1000:]}",
                },
                ensure_ascii=True,
            )
        # Achado real (checado ao vivo nesta máquina, 2026-08-12): quando os
        # navegadores já estão instalados, `playwright install` roda em
        # segundos e não imprime NADA (stdout vazio, sem linhas "Downloading
        # ..."). Só quando baixa algo de fato é que aparece "Downloading" no
        # stdout. Usa isso pra avisar corretamente "já estava instalado" em
        # vez de sempre dizer "instalado com sucesso" mesmo sem baixar nada.
        already_installed = "Downloading" not in proc.stdout
        message = (
            "Chromium, Firefox e WebKit já estavam instalados nesta máquina -- nada precisou ser baixado. "
            "O teste cross-browser real já pode usar os 3 motores."
            if already_installed
            else "Chromium, Firefox e WebKit instalados com sucesso agora. O teste cross-browser real já pode usar os 3 motores."
        )
        logger.info("[a11y_chat] install_playwright_browsers concluído (already_installed=%s).", already_installed)
        return json.dumps(
            {
                "status": "ok",
                "already_installed": already_installed,
                "message": message,
            },
            ensure_ascii=True,
        )
    except subprocess.TimeoutExpired:
        return json.dumps({"status": "error", "error": "Instalação excedeu o tempo limite de 600s."}, ensure_ascii=True)
    except Exception as exc:
        logger.error("[a11y_chat] install_playwright_browsers falhou: %s", exc)
        return json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=True)


def _cache_remote_test_result_for_deliverables(runner: str, target_url: str, raw_violations: list[dict]) -> None:
    """Achado real (pedido do usuário, 2026-08-11): sem isso, um teste real de
    Cypress/Selenium detectava problemas de verdade mas a IA não conseguia
    depois gerar planilha/checklist/PDF/VPAT a partir DESSES resultados --
    generate_checklist/export_xlsx etc. só liam de last_analysis_store, que só
    era populado por analyze_page/analyze_site. Agora um teste remoto real
    alimenta o mesmo cache, e também busca+cacheia o HTML renderizado da
    página testada (last_analyzed_content_store) para fix_and_zip_files/
    open_live_preview funcionarem a partir de um teste remoto também, não só
    de uma análise via analyze_page. Nunca bloqueia o resultado do teste em
    si -- falha aqui só significa "sem entrega automática depois", registrada
    em log, não repassada como erro pro usuário.
    """
    from backend.src.services.self_healing import axe_violations_to_accessibility_issues

    try:
        issues = axe_violations_to_accessibility_issues(raw_violations)
        issues_dicts = [i.model_dump() for i in issues]
        for issue in issues_dicts:
            issue["url"] = target_url

        from backend.src.services.last_analysis_store import set_last_analysis

        set_last_analysis(issues_dicts, target_url)

        from backend.src.routes.analyze import _extract_semantic_html
        from backend.src.services.browser import fetch_rendered_html_and_screenshot

        raw_html, _screenshot = _safe_async_run(fetch_rendered_html_and_screenshot(target_url))
        resolved_content = _extract_semantic_html(raw_html)
        from backend.src.services.last_analyzed_content_store import set_last_analyzed_content

        set_last_analyzed_content(resolved_content, target_url)

        logger.info(
            "[a11y_chat] Resultado do teste remoto (%s) cacheado para entregas (planilha/checklist/PDF/correção): "
            "%d issues.",
            runner,
            len(issues_dicts),
        )
    except Exception as exc:
        logger.warning(
            "[a11y_chat] Não foi possível cachear o resultado do teste remoto (%s) para entregas automáticas: %s",
            runner,
            exc,
        )


_RUN_DEEP_RESEARCH_SCHEMA: dict[str, Any] = {
    "description": (
        "Executa uma pesquisa normativa profunda de acessibilidade (WCAG 2.2, WAI-ARIA APG, "
        "Section 508, EN 301 549, PDF/UA) investigando fontes primárias e trazendo citações completas."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "question": {
                "type": "string",
                "description": "A pergunta ou dúvida técnica sobre a norma de acessibilidade a ser investigada.",
            },
        },
        "required": ["question"],
    },
}


def run_deep_research_tool(args: dict) -> str:
    """Handler síncrono para despachar o agente de Deep Research."""
    question = args.get("question", "")
    from backend.src.agents.deep_research.deep_research import run_deep_research

    chat_progress.emit_tool_progress(None, "run_deep_research", "Investigando fontes normativas...")
    res = _safe_async_run(run_deep_research(question))
    return json.dumps(res, ensure_ascii=False)


_CLARIFY_SCHEMA: dict[str, Any] = {
    "description": (
        "Pergunta algo ao usuário e aguarda a resposta antes de continuar. Use para apresentar "
        "um plano de análise e pedir aprovação, ou para confirmar correções antes de alterar código. "
        "Devolve JSON com a resposta escolhida pelo usuário."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "question": {
                "type": "string",
                "description": "A pergunta a fazer ao usuário, em linguagem natural.",
            },
            "options": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Opções de resposta sugeridas (ex.: ['Aprovar Plano', 'Alterar Foco', 'Cancelar']).",
            },
        },
        "required": ["question"],
    },
}


def clarify(args: dict, clarify_callback: Callable[[str, list[str]], str]) -> str:
    """Faz a pergunta ao usuário através do canal de clarify do chat e devolve a resposta."""
    question = str(args.get("question") or "").strip()
    if not question:
        return json.dumps({"error": "A pergunta não pode ser vazia."}, ensure_ascii=False)
    options = [str(option) for option in (args.get("options") or [])]
    answer = str(clarify_callback(question, options))
    logger.info("[clarify] pergunta respondida pelo usuário (%d opções)", len(options))
    return json.dumps(
        {"question": question, "options": options, "answer": answer},
        ensure_ascii=False,
    )


def register_chat_tools() -> None:
    """Registra o toolset 'a11y_chat' no registry local de tools (idempotente)."""
    try:
        from tools.registry import registry
    except Exception as exc:  # pragma: no cover - registry indisponivel
        logger.warning("[a11y_chat] tools.registry indisponivel: %s", exc)
        return

    registry.register(
        name="analyze_page",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_ANALYZE_PAGE_SCHEMA,
        handler=analyze_page,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (analyze_page)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="analyze_site",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_ANALYZE_SITE_SCHEMA,
        handler=analyze_site,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (analyze_site)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="unzip_and_list_files",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_UNZIP_AND_LIST_FILES_SCHEMA,
        handler=unzip_and_list_files,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (unzip_and_list_files)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="read_local_project_files",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_READ_LOCAL_PROJECT_FILES_SCHEMA,
        handler=read_local_project_files,
        is_async=False,
        emoji="",
        # Lê arquivos reais de um caminho no disco que o usuário informou: efeito
        # real (acesso a filesystem fora da sandbox de HTML), exige confirmação.
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (read_local_project_files)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="fix_local_project_files",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_FIX_LOCAL_PROJECT_FILES_SCHEMA,
        handler=fix_local_project_files,
        is_async=False,
        emoji="",
        # Sobrescreve arquivos reais no disco do usuário: efeito destrutivo, exige
        # confirmação explícita (mesmo com backup automático).
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (fix_local_project_files)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="analyze_document",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_ANALYZE_DOCUMENT_SCHEMA,
        handler=analyze_document,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (analyze_document)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="fix_and_zip_files",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_FIX_AND_ZIP_FILES_SCHEMA,
        handler=fix_and_zip_files,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (fix_and_zip_files)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="undo_last_fix",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_UNDO_LAST_FIX_SCHEMA,
        handler=undo_last_fix,
        is_async=False,
        emoji="",
        # Descarta o resultado da remediação e repõe o estado anterior: efeito real.
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (undo_last_fix)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="export_xlsx",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_EXPORT_XLSX_SCHEMA,
        handler=export_xlsx,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (export_xlsx)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="generate_checklist",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_GENERATE_CHECKLIST_SCHEMA,
        handler=generate_checklist,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (generate_checklist)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="export_checklist_pdf",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_EXPORT_CHECKLIST_PDF_SCHEMA,
        handler=export_checklist_pdf_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (export_checklist_pdf)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="generate_accessibility_statement",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_GENERATE_ACCESSIBILITY_STATEMENT_SCHEMA,
        handler=generate_accessibility_statement,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (generate_accessibility_statement)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="export_accessibility_statement_pdf",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_EXPORT_ACCESSIBILITY_STATEMENT_PDF_SCHEMA,
        handler=export_accessibility_statement_pdf_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (export_accessibility_statement_pdf)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="generate_vpat",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_GENERATE_VPAT_SCHEMA,
        handler=generate_vpat,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (generate_vpat)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="generate_test_suite",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_GENERATE_TEST_SUITE_SCHEMA,
        handler=generate_test_suite,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (generate_test_suite)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="open_live_preview",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_OPEN_LIVE_PREVIEW_SCHEMA,
        handler=open_live_preview,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (open_live_preview)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="tavily_search",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_TAVILY_SEARCH_SCHEMA,
        handler=tavily_search,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (tavily_search)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="exa_search",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_EXA_SEARCH_SCHEMA,
        handler=exa_search,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (exa_search)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="evaluate_research",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_EVALUATE_RESEARCH_SCHEMA,
        handler=evaluate_research,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (evaluate_research)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="create_github_issue",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_CREATE_GITHUB_ISSUE_SCHEMA,
        handler=create_github_issue_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (create_github_issue)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="create_jira_issue",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_CREATE_JIRA_ISSUE_SCHEMA,
        handler=create_jira_issue_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (create_jira_issue)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="create_azure_devops_work_item",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_CREATE_AZURE_DEVOPS_WORK_ITEM_SCHEMA,
        handler=create_azure_devops_work_item_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (create_azure_devops_work_item)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="nvda_speak",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_NVDA_SPEAK_SCHEMA,
        handler=nvda_speak_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (nvda_speak)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="verify_screen_reader_announcements",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_VERIFY_SCREEN_READER_ANNOUNCEMENTS_SCHEMA,
        handler=verify_screen_reader_announcements_tool,
        is_async=False,
        emoji="",
        # speak_via_nvda=true produz efeito audivel real na maquina do usuario --
        # mesma classe de efeito colateral do nvda_speak acima.
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (verify_screen_reader_announcements)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="design_review",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_DESIGN_REVIEW_SCHEMA,
        handler=design_review_tool,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (design_review)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="generate_automation_script",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_GENERATE_AUTOMATION_SCRIPT_SCHEMA,
        handler=generate_automation_script_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (generate_automation_script)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="run_remote_test",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_RUN_REMOTE_TEST_SCHEMA,
        handler=run_remote_test_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (run_remote_test)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="run_cross_browser_test",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_RUN_CROSS_BROWSER_TEST_SCHEMA,
        handler=run_cross_browser_test_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (run_cross_browser_test)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="install_playwright_browsers",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_INSTALL_PLAYWRIGHT_BROWSERS_SCHEMA,
        handler=install_playwright_browsers_tool,
        is_async=False,
        emoji="",
        requires_approval=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (install_playwright_browsers)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="run_deep_research",
        toolset=A11Y_CHAT_TOOLSET,
        schema=_RUN_DEEP_RESEARCH_SCHEMA,
        handler=run_deep_research_tool,
        is_async=False,
        emoji="",
    )
    logger.info("[a11y_chat] toolset '%s' registrado (run_deep_research)", A11Y_CHAT_TOOLSET)

    registry.register(
        name="clarify",
        toolset=CLARIFY_TOOLSET,
        schema=_CLARIFY_SCHEMA,
        handler=clarify,
        is_async=False,
        emoji="",
        needs_clarify_callback=True,
    )
    logger.info("[a11y_chat] toolset '%s' registrado (clarify)", CLARIFY_TOOLSET)


register_chat_tools()
