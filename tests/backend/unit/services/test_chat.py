import json
import threading
from typing import Any
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

from backend.src.services.chat_runtime import (
    _DEFAULT_CLARIFY_QUESTION,
    format_human_friendly_error,
    stream_chat,
)
from backend.src.services.chat_tools import (
    A11Y_CHAT_TOOLSET,
    CLARIFY_TOOLSET,
    analyze_page,
    analyze_site,
    export_checklist_pdf_tool,
    fix_and_zip_files,
    generate_checklist,
    register_chat_tools,
    run_remote_test_tool,
    unzip_and_list_files,
)
from backend.src.shared.models import AgentResult


def test_format_human_friendly_error_ollama_402():
    raw = "Error code: 402 - {'error': 'this model uses extra usage only (not included plan usage) and your extra usage balance is empty, add extra usage or turn on auto reload at https://ollama.com/settings'}"
    formatted = format_human_friendly_error(raw)
    assert "saldo de créditos" in formatted
    assert "Configurações" in formatted
    assert "ollama.com/settings" in formatted


class _CapturingAgent:
    last_kwargs: dict[str, Any] = {}
    last_user_message = ""

    def __init__(self, **kwargs):
        type(self).last_kwargs = kwargs

    def run_conversation(self, user_message):
        type(self).last_user_message = user_message
        return {"final_response": "ok", "failed": False}


# ── chat_tools.analyze_page ───────────────────────────────────────────────────


class TestAnalyzePageTool:
    def test_no_input_returns_error(self):
        out = json.loads(analyze_page({}))
        assert "error" in out

    def test_url_path_fetches_and_analyzes(self):
        issues = [{"id": "p-1", "criterion": "1.1.1", "severity": "high",
                   "element": "img", "description": "x"}]
        fake = AgentResult(agent="orchestrator", success=True, data={"issues": issues})
        with patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(return_value=("<html><img src=x></html>", "fake_screenshot")),
        ) as fetch_mock, patch(
            "backend.src.services.browser.fetch_accessibility_tree_snapshot",
            new=AsyncMock(return_value=""),
        ), patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<img src=x>"
        ), patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate",
            new=AsyncMock(return_value=fake),
        ):
            out = json.loads(analyze_page({"url": "https://example.com"}))

        assert out["total_issues"] == 1
        assert out["counts_by_severity"]["high"] == 1
        assert out["score"] == 90
        fetch_mock.assert_awaited_once_with("https://example.com")

    def test_counts_by_severity_matches_top_issues_when_severity_is_enum_member(self):
        """Achado real (147 issues, counts_by_severity todos 0, mas top_issues[0]
        com severity="low"): model_dump() em modo python preserva o membro do
        enum Severity (str, Enum), e str(Severity.LOW) retorna "Severity.LOW"
        (Enum.__str__ sobrescrito), nao "low" -- quebrando o match contra
        `counts`. Ver backend/src/services/chat_tools.py::_summarize_issues.
        """
        from backend.src.shared.models import Severity

        issues = [
            {"id": "p-1", "criterion": "1.1.1", "severity": Severity.LOW, "element": "img", "description": "x"},
            {"id": "p-2", "criterion": "2.1.1", "severity": Severity.CRITICAL, "element": "btn", "description": "y"},
        ]
        fake = AgentResult(agent="orchestrator", success=True, data={"issues": issues})
        with patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(return_value=("<html></html>", None)),
        ), patch(
            "backend.src.services.browser.fetch_accessibility_tree_snapshot",
            new=AsyncMock(return_value=""),
        ), patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<html></html>"
        ), patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate",
            new=AsyncMock(return_value=fake),
        ):
            out = json.loads(analyze_page({"url": "https://example.com"}))

        assert out["total_issues"] == 2
        assert out["counts_by_severity"]["low"] == 1
        assert out["counts_by_severity"]["critical"] == 1
        assert out["top_issues"][0]["severity"] == "low"
        assert out["score"] == 78


class TestAnalyzePageRegressionVsPreviousScan:
    """Achado real (2026-08-11, "shift-right" sob demanda, pedido do usuário):
    reanalisar a MESMA URL depois de já ter sido analisada antes deve
    detectar regressão real (issues novos) e melhoria real (issues
    resolvidos), sem criar nada sozinho -- só sinaliza pro modelo oferecer
    `create_github_issue`."""

    def setup_method(self):
        import shutil

        from backend.src.services import url_scan_history_store
        shutil.rmtree(url_scan_history_store._HISTORY_DIR, ignore_errors=True)

    def _run_analyze(self, url: str, issues: list[dict]):
        fake = AgentResult(agent="orchestrator", success=True, data={"issues": issues})
        with patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(return_value=("<html></html>", None)),
        ), patch(
            "backend.src.services.browser.fetch_accessibility_tree_snapshot",
            new=AsyncMock(return_value=""),
        ), patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<html></html>"
        ), patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate",
            new=AsyncMock(return_value=fake),
        ):
            return json.loads(analyze_page({"url": url}))

    def test_first_scan_of_a_url_has_no_regression_field(self):
        out = self._run_analyze(
            "https://example.com/regressao-1",
            [{"id": "p-1", "criterion": "1.1.1", "severity": "high", "element": "img.a", "description": "x"}],
        )
        assert "regression_vs_previous_scan" not in out

    def test_second_scan_with_new_issue_reports_real_regression(self):
        url = "https://example.com/regressao-2"
        self._run_analyze(url, [
            {"id": "p-1", "criterion": "1.1.1", "severity": "high", "element": "img.a", "description": "x"},
        ])
        out = self._run_analyze(url, [
            {"id": "p-1", "criterion": "1.1.1", "severity": "high", "element": "img.a", "description": "x"},
            {"id": "p-2", "criterion": "2.4.4", "severity": "medium", "element": "a.link", "description": "y"},
        ])
        assert out["regression_vs_previous_scan"]["new_issues_count"] == 1
        assert out["regression_vs_previous_scan"]["new_issues_preview"][0]["criterion"] == "2.4.4"
        assert out["regression_vs_previous_scan"]["resolved_issues_count"] == 0

    def test_second_scan_with_fixed_issue_reports_real_improvement(self):
        url = "https://example.com/regressao-3"
        self._run_analyze(url, [
            {"id": "p-1", "criterion": "1.1.1", "severity": "high", "element": "img.a", "description": "x"},
            {"id": "p-2", "criterion": "2.4.4", "severity": "medium", "element": "a.link", "description": "y"},
        ])
        out = self._run_analyze(url, [
            {"id": "p-1", "criterion": "1.1.1", "severity": "high", "element": "img.a", "description": "x"},
        ])
        assert out["regression_vs_previous_scan"]["resolved_issues_count"] == 1
        assert out["regression_vs_previous_scan"]["new_issues_count"] == 0

    def test_identical_second_scan_has_no_regression_field(self):
        url = "https://example.com/regressao-4"
        issues = [{"id": "p-1", "criterion": "1.1.1", "severity": "high", "element": "img.a", "description": "x"}]
        self._run_analyze(url, issues)
        out = self._run_analyze(url, issues)
        assert "regression_vs_previous_scan" not in out


class TestAnalyzeSiteTool:
    def test_no_input_returns_error(self):
        out = json.loads(analyze_site({}))
        assert "error" in out

    def test_crawl_site_path_crawls_and_analyzes(self):
        from backend.src.services.crawler import CrawlPageResult
        pages = [CrawlPageResult(url="https://example.com/a", html="<html></html>")]
        issues = [{"id": "p-1", "criterion": "1.1.1", "severity": "high",
                   "element": "img", "description": "x"}]
        fake = AgentResult(agent="orchestrator", success=True, data={"issues": issues})

        with patch(
            "backend.src.services.crawler.crawl_site",
            new=AsyncMock(return_value=pages),
        ) as crawl_mock, patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<html></html>"
        ), patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate",
            new=AsyncMock(return_value=fake),
        ) as orch_mock:
            out = json.loads(analyze_site({"url": "https://example.com", "max_pages": 5}))

        crawl_mock.assert_awaited_once_with("https://example.com", max_pages=5)
        orch_mock.assert_awaited_once()
        assert out["total_pages"] == 1
        assert out["total_issues"] == 1
        assert out["score"] == 90

    def test_list_urls_path_analyzes_each(self):
        issues = [{"id": "p-1", "criterion": "1.1.1", "severity": "low",
                   "element": "p", "description": "alt"}]
        fake = AgentResult(agent="orchestrator", success=True, data={"issues": issues})

        with patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(return_value=("<html></html>", None)),
        ) as fetch_mock, patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<html></html>"
        ), patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate",
            new=AsyncMock(return_value=fake),
        ) as orch_mock, patch(
            "backend.src.services.last_analysis_store.set_last_analysis"
        ) as store_mock:
            out = json.loads(analyze_site({"urls": ["https://example.com/a", "https://example.com/b"]}))

        assert fetch_mock.call_count == 2
        assert orch_mock.call_count == 2
        assert out["total_pages"] == 2
        assert out["pages_ok"] == 2
        assert out["total_issues"] == 2
        assert out["score"] == 96
        store_mock.assert_called_once()

    def test_registered_in_hermes(self):
        register_chat_tools()
        from tools.registry import registry
        assert "analyze_page" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)
        assert "analyze_site" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)
        # Regression test: generate_vpat/generate_test_suite are real backend routes
        # (/analyze/vpat, /analyze/tests) that were never wrapped as chat tools -- the
        # README advertises VPAT/test-suite generation as headline features, but they
        # were unreachable through any real user action (chat or UI).
        assert "generate_vpat" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)
        assert "generate_test_suite" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)
        # Regression test: open_live_preview is what lets the model actually trigger the
        # [LIVE_PREVIEW:id:total] marker the frontend (ChatScreen.tsx) watches for -- without
        # this tool the Live Preview feature has no way to ever fire.
        assert "open_live_preview" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)
        assert "generate_accessibility_statement" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)
        assert "export_accessibility_statement_pdf" in registry.get_tool_names_for_toolset(A11Y_CHAT_TOOLSET)


class TestAccessibilityStatementTools:
    def test_without_prior_analysis_returns_clear_error(self):
        from backend.src.services.chat_tools import generate_accessibility_statement
        with patch("backend.src.services.last_analysis_store.get_last_analysis", return_value=([], "")):
            out = json.loads(generate_accessibility_statement({"pre_exec_msg": "Gerando declaração..."}))
        assert "error" in out
        assert "análise" in out["error"].lower()

    def test_uses_cached_issues_and_never_fabricates_org_when_not_given(self):
        from backend.src.services.chat_tools import generate_accessibility_statement

        cached_issue = {
            "id": "issue-1", "criterion": "1.1.1 Non-text Content", "severity": "critical",
            "element": "img", "description": "Imagem sem texto alternativo",
        }
        with patch("backend.src.services.last_analysis_store.get_last_analysis",
                   return_value=([cached_issue], "https://example.com")):
            out = json.loads(generate_accessibility_statement({"pre_exec_msg": "Gerando declaração..."}))

        assert out["total_issues"] == 1
        assert out["conformance_level"] == "Parcialmente conforme"
        assert out["organization_name"] == "[Nome da Organização]"
        assert out["known_limitations"][0]["criterion"] == "1.1.1 Non-text Content"

    def test_uses_real_org_and_contact_when_provided_never_a_default_fabrication(self):
        from backend.src.services.chat_tools import generate_accessibility_statement

        cached_issue = {"id": "issue-1", "criterion": "1.1.1 Non-text Content", "severity": "high", "description": "x"}
        with patch("backend.src.services.last_analysis_store.get_last_analysis",
                   return_value=([cached_issue], "https://example.com")):
            out = json.loads(generate_accessibility_statement({
                "pre_exec_msg": "Gerando declaração...",
                "organization_name": "Acme Ltda",
                "contact_email": "a11y@acme.com",
            }))

        assert out["organization_name"] == "Acme Ltda"
        assert out["contact_email"] == "a11y@acme.com"

    def test_export_pdf_returns_download_link(self):
        from backend.src.services.chat_tools import export_accessibility_statement_pdf_tool
        out = json.loads(export_accessibility_statement_pdf_tool({}))
        assert "last_accessibility_statement_pdf" in out["download_url"]


class TestGenerateVpatAndTestSuiteTools:
    def test_generate_vpat_without_prior_analysis_returns_clear_error(self):
        from backend.src.services.chat_tools import generate_vpat
        with patch("backend.src.services.last_analysis_store.get_last_analysis", return_value=([], "")):
            out = json.loads(generate_vpat({"pre_exec_msg": "Gerando VPAT..."}))
        assert "error" in out
        assert "análise" in out["error"].lower()

    def test_generate_test_suite_without_prior_analysis_returns_clear_error(self):
        from backend.src.services.chat_tools import generate_test_suite
        with patch("backend.src.services.last_analysis_store.get_last_analysis", return_value=([], "")):
            out = json.loads(generate_test_suite({"pre_exec_msg": "Gerando testes..."}))
        assert "error" in out

    def test_generate_vpat_uses_cached_issues_and_returns_report(self):
        from backend.src.services.chat_tools import generate_vpat
        from backend.src.shared.models import AgentResult

        cached_issue = {
            "id": "issue-1",
            "guideline": "wcag_2_2",
            "criterion": "1.1.1 Non-text Content",
            "severity": "critical",
            "element": "img",
            "description": "Imagem sem texto alternativo",
            "suggestion": "Adicione um atributo alt descritivo",
            "wcag_url": "https://www.w3.org/WAI/WCAG22/Understanding/non-text-content",
        }
        fake_result = AgentResult(agent="vpat_reporter", success=True, data={"vpat": {"total_criteria_evaluated": 50}})

        with patch("backend.src.services.last_analysis_store.get_last_analysis",
                   return_value=([cached_issue], "https://example.com")), \
             patch("backend.src.agents.vpat_reporter.vpat_reporter.run_vpat_reporter",
                   new=AsyncMock(return_value=fake_result)) as mock_run:
            out = json.loads(generate_vpat({"pre_exec_msg": "Gerando VPAT...", "product_name": "Meu Site"}))

        assert out["vpat"]["total_criteria_evaluated"] == 50
        mock_run.assert_awaited_once()
        _, call_kwargs = mock_run.call_args
        assert call_kwargs["product_name"] == "Meu Site"
        assert call_kwargs["target"] == "https://example.com"
        assert len(call_kwargs["issues"]) == 1

    def test_generate_test_suite_uses_cached_issues_and_returns_report(self):
        from backend.src.services.chat_tools import generate_test_suite
        from backend.src.shared.models import AgentResult

        cached_issue = {
            "id": "issue-1",
            "guideline": "wcag_2_2",
            "criterion": "1.1.1 Non-text Content",
            "severity": "critical",
            "element": "img",
            "description": "Imagem sem texto alternativo",
            "suggestion": "Adicione um atributo alt descritivo",
            "wcag_url": "https://www.w3.org/WAI/WCAG22/Understanding/non-text-content",
        }
        fake_result = AgentResult(agent="test_generator", success=True, data={"suite": {"total_tests": 3}})

        with patch("backend.src.services.last_analysis_store.get_last_analysis",
                   return_value=([cached_issue], "https://example.com")), \
             patch("backend.src.agents.test_generator.test_generator.run_test_generator",
                   new=AsyncMock(return_value=fake_result)) as mock_run:
            out = json.loads(generate_test_suite({"pre_exec_msg": "Gerando testes..."}))

        assert out["suite"]["total_tests"] == 3
        mock_run.assert_awaited_once()


class TestOpenLivePreviewTool:
    def test_without_prior_fix_returns_clear_error(self):
        from backend.src.services.chat_tools import open_live_preview
        with patch("backend.src.services.last_fix_store.get_last_fix", return_value=[]):
            out = json.loads(open_live_preview({"pre_exec_msg": "Abrindo preview..."}))
        assert "error" in out
        assert "fix_and_zip_files" in out["error"]

    def test_creates_session_from_cached_fix_pages(self):
        from backend.src.services.chat_tools import open_live_preview

        cached_pages = [
            {"title": "index.html", "original_html": "<p>old</p>", "fixed_html": "<p>new</p>"},
            {"title": "about.html", "original_html": "<p>old2</p>", "fixed_html": "<p>new2</p>"},
        ]

        with patch("backend.src.services.last_fix_store.get_last_fix", return_value=cached_pages), \
             patch("backend.src.routes.preview.register_preview_session", return_value="abc12345") as mock_register:
            out = json.loads(open_live_preview({"pre_exec_msg": "Abrindo preview..."}))

        assert out["session_id"] == "abc12345"
        assert out["total_pages"] == 2
        mock_register.assert_called_once_with(cached_pages)

    def test_session_created_is_actually_renderable(self):
        """Integration-style: a session opened via the tool must be servable by the real /preview/render route."""
        from backend.src.services import last_fix_store
        from backend.src.services.chat_tools import open_live_preview

        cached_pages = [{"title": "index.html", "original_html": "<p>old</p>", "fixed_html": "<p>new</p>"}]
        last_fix_store.set_last_fix(cached_pages)
        try:
            out = json.loads(open_live_preview({"pre_exec_msg": "Abrindo preview..."}))
            session_id = out["session_id"]

            from fastapi.testclient import TestClient

            from backend.src.main import app
            client = TestClient(app)
            resp = client.get(f"/preview/render/{session_id}/0")
            assert resp.status_code == 200
            assert "new" in resp.text
        finally:
            last_fix_store.set_last_fix([])


class _FakeAgent:
    """Simula o AIAgent: ao rodar, dispara os callbacks e retorna a resposta."""

    def __init__(self, **kwargs):
        self._cb = kwargs

    def run_conversation(self, user_message, **_k):
        assert self._cb["thinking_callback"] is not None
        assert self._cb["reasoning_callback"] is not None
        self._cb["thinking_callback"]("Analisando o problema...")
        self._cb["reasoning_callback"]("Raciocinando sobre WCAG 2.2...")
        self._cb["tool_start_callback"]("id1", "analyze_page", {})
        self._cb["tool_complete_callback"]("id1", "analyze_page", {}, "{}")
        self._cb["stream_delta_callback"]("Ola")
        self._cb["stream_delta_callback"](" mundo")
        return {"final_response": "Ola mundo", "failed": False}


class _FailingAgent:
    def __init__(self, **kwargs):
        pass

    def run_conversation(self, user_message, **_k):
        return {"failed": True, "error": "provider 400"}


@pytest.mark.asyncio
async def test_stream_chat_emits_event_sequence():
    with patch("backend.src.services.chat_runtime.AIAgent", new=_FakeAgent):
        events = [ev async for ev in stream_chat("audita isto", history=[])]

    types = [e["type"] for e in events]
    assert types[0] == "stream_id"
    assert "thinking" in types
    assert "reasoning" in types
    assert "tool_start" in types
    assert "tool_result" in types
    assert types.count("token") == 2
    assert types[-1] == "done"
    assert events[-1]["final"] == "Ola mundo"


@pytest.mark.asyncio
async def test_stream_chat_persists_user_and_assistant_messages_to_history(tmp_path, monkeypatch):
    """chat_history_store.py: histórico de sessão persistido no backend --
    antes só existia em memória do React (perdido em cada reload)."""
    monkeypatch.setattr("tempfile.gettempdir", lambda: str(tmp_path))
    from backend.src.services import chat_history_store
    chat_history_store._sessions.clear()

    with patch("backend.src.services.chat_runtime.AIAgent", new=_FakeAgent):
        _ = [ev async for ev in stream_chat("audita isto", history=[], conversation_id="conv-persist-teste")]

    history = chat_history_store.get_history(session_id="conv-persist-teste")
    assert [m["role"] for m in history] == ["user", "assistant"]
    assert history[0]["content"] == "audita isto"
    assert history[1]["content"] == "Ola mundo"


class _SlowAgent:
    """Simula um turno lento: bloqueia num Event até o teste liberar,
    permitindo cancelar o turno enquanto ele ainda "está rodando" na thread."""

    release: threading.Event

    def __init__(self, **kwargs):
        self._cb = kwargs

    def run_conversation(self, user_message, **_k):
        _SlowAgent.release.wait(timeout=5.0)
        return {"final_response": "resposta tardia", "failed": False}


@pytest.mark.asyncio
async def test_stream_chat_cancel_stops_the_stream_before_agent_finishes():
    from backend.src.services import chat_progress

    _SlowAgent.release = threading.Event()
    with patch("backend.src.services.chat_runtime.AIAgent", new=_SlowAgent):
        gen = stream_chat("audita isto", history=[])
        first = await gen.__anext__()
        assert first["type"] == "stream_id"
        stream_id = first["id"]

        assert chat_progress.request_cancel(stream_id) is True

        events = [ev async for ev in gen]

    assert events[-1]["type"] == "cancelled"
    # Libera a thread bloqueada só depois da asserção, para não deixar o
    # agente "preso" além do teste (limpeza, não parte do comportamento testado).
    _SlowAgent.release.set()


@pytest.mark.asyncio
async def test_chat_cancel_route_delivers_to_the_progress_registry():
    from backend.src.routes.chat import CancelRequest, chat_cancel
    from backend.src.services import chat_progress

    token = chat_progress.new_cancel_token()
    try:
        result = await chat_cancel(CancelRequest(stream_id=token))
        assert result == {"cancelled": True}
        ev = chat_progress.cancel_event(token)
        assert ev is not None and ev.is_set()
    finally:
        chat_progress.clear_cancel_token(token)


@pytest.mark.asyncio
async def test_chat_cancel_route_unknown_stream_id_returns_false():
    from backend.src.routes.chat import CancelRequest, chat_cancel

    result = await chat_cancel(CancelRequest(stream_id="naoexiste"))
    assert result == {"cancelled": False}


@pytest.mark.asyncio
async def test_stream_chat_emits_error_event():
    with patch("backend.src.services.chat_runtime.AIAgent", new=_FailingAgent):
        events = [ev async for ev in stream_chat("oi")]
    assert events[-1]["type"] == "error"
    assert "provedor de IA" in events[-1]["error"]


@pytest.mark.asyncio
async def test_stream_chat_applies_provider_model_override():
    with patch("backend.src.services.chat_runtime.AIAgent", new=_CapturingAgent):
        _ = [ev async for ev in stream_chat("oi", provider="anthropic", model="claude-opus-4-5")]
    assert _CapturingAgent.last_kwargs["provider"] == "anthropic"
    assert _CapturingAgent.last_kwargs["model"] == "claude-opus-4-5"


@pytest.mark.asyncio
async def test_stream_chat_uses_controlled_toolsets_only():
    with patch("backend.src.services.chat_runtime.AIAgent", new=_CapturingAgent):
        _ = [ev async for ev in stream_chat("análise a página")]
    assert _CapturingAgent.last_kwargs["enabled_toolsets"] == [A11Y_CHAT_TOOLSET, CLARIFY_TOOLSET]
    assert "web" not in _CapturingAgent.last_kwargs["enabled_toolsets"]
    # O toolset pedido tem de existir de facto no registry, senão o modelo fica
    # sem a tool `clarify` que as regras 12 e 13 do system prompt exigem.
    from tools.registry import registry
    assert registry.get_tool_names_for_toolset(CLARIFY_TOOLSET) == ["clarify"]


@pytest.mark.asyncio
async def test_stream_chat_surfaces_the_clarifier_question(monkeypatch):
    """Bug real de auditoria: o intent `needs_clarification` do clarifier e a
    pergunta gerada por ele eram calculados e descartados -- só `out_of_scope`
    tinha branch, e o turno seguia como se o pedido estivesse completo."""
    monkeypatch.delenv("PYTEST_CURRENT_TEST", raising=False)
    clarifier_result = AgentResult(
        agent="clarifier",
        success=True,
        data={
            "intent": "needs_clarification",
            "needs_clarification": True,
            "question": "Qual página você quer que eu audite?",
            "explanation": "Pedido sem alvo.",
        },
    )
    with patch(
        "backend.src.agents.clarifier.run_clarifier",
        new=AsyncMock(return_value=clarifier_result),
    ), patch("backend.src.services.chat_runtime.AIAgent", new=_CapturingAgent):
        events = [ev async for ev in stream_chat("analisa isto")]

    assert any(
        ev["type"] == "token" and "Qual página você quer que eu audite?" in ev["text"]
        for ev in events
    )
    assert events[-1]["type"] == "done"
    assert events[-1]["final"] == "Qual página você quer que eu audite?"


@pytest.mark.asyncio
async def test_stream_chat_falls_back_to_a_default_clarify_question(monkeypatch):
    """Se o clarifier classificar como ambíguo sem devolver pergunta, o turno
    ainda tem de perguntar algo ao usuário em vez de seguir em frente."""
    monkeypatch.delenv("PYTEST_CURRENT_TEST", raising=False)
    clarifier_result = AgentResult(
        agent="clarifier",
        success=True,
        data={"intent": "needs_clarification", "needs_clarification": True, "question": ""},
    )
    with patch(
        "backend.src.agents.clarifier.run_clarifier",
        new=AsyncMock(return_value=clarifier_result),
    ), patch("backend.src.services.chat_runtime.AIAgent", new=_CapturingAgent):
        events = [ev async for ev in stream_chat("analisa isto")]

    assert events[-1]["type"] == "done"
    assert events[-1]["final"] == _DEFAULT_CLARIFY_QUESTION


@pytest.mark.asyncio
async def test_stream_chat_merges_consecutive_messages_and_prevents_duplicate_roles():
    history = [
        {"role": "user", "content": "olá"},
        {"role": "user", "content": "tudo bem?"},
        {"role": "assistant", "content": "sim, e com você?"},
        {"role": "assistant", "content": "como posso ajudar?"},
        {"role": "user", "content": "preciso de ajuda"},
    ]
    with patch("backend.src.services.chat_runtime.AIAgent", new=_CapturingAgent):
        _ = [ev async for ev in stream_chat("com acessibilidade", history=history)]
    prefill = _CapturingAgent.last_kwargs["prefill_messages"]
    assert len(prefill) == 2
    assert prefill[0]["role"] == "user"
    assert prefill[0]["content"] == "olá\n\ntudo bem?"
    assert prefill[1]["role"] == "assistant"
    assert prefill[1]["content"] == "sim, e com você?\n\ncomo posso ajudar?"
    assert _CapturingAgent.last_user_message == "preciso de ajuda\n\ncom acessibilidade"


class TestUnzipAndListFilesTool:
    def test_unzip_no_input_returns_error(self):
        out = json.loads(unzip_and_list_files({}))
        assert "error" in out

    def test_unzip_valid_zip_extracts_files(self):
        import base64
        import io
        import zipfile

        # Cria um zip simples na memória
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w") as z:
            z.writestr("index.html", "<html></html>")
            z.writestr("styles.css", "body {}")
            z.writestr("invalid.exe", "binary")

        zip_base64 = base64.b64encode(zip_buffer.getvalue()).decode("utf-8")
        out = json.loads(unzip_and_list_files({"zip_base64": zip_base64}))

        assert isinstance(out, list)
        assert len(out) == 2  # Apenas index.html e styles.css
        paths = [item["path"] for item in out]
        assert "index.html" in paths
        assert "styles.css" in paths
        assert "invalid.exe" not in paths


class TestFixAndZipFilesTool:
    @pytest.fixture(autouse=True)
    def _clean_last_analyzed_content(self):
        """Sem isso, um teste anterior que chama analyze_page() nesta mesma
        sessão (default) deixa HTML cacheado em last_analyzed_content_store,
        e o fallback de fix_and_zip_files passa a achar conteúdo onde os
        testes abaixo esperam vazio."""
        from backend.src.services import last_analyzed_content_store
        last_analyzed_content_store._sessions.clear()
        yield
        last_analyzed_content_store._sessions.clear()

    def test_fix_no_files_returns_error(self):
        out = json.loads(fix_and_zip_files({}))
        assert "error" in out

    def test_fix_and_zip_html_files(self):
        fake_audit = AgentResult(agent="orchestrator", success=True, data={"issues": [{"id": "1", "criterion": "1.1.1", "severity": "high", "element": "img", "description": "x"}]})
        fake_fix = AgentResult(agent="fixer", success=True, data={"fixed_html": "<html>fixed</html>", "changes_summary": ["fixed img alt"]})
        fake_non_html_fix = {"path": "styles.css", "content": "body {}", "changes_summary": []}

        files = [
            {"path": "index.html", "content": "<html><img src=x></html>"},
            {"path": "styles.css", "content": "body {}"}
        ]

        with patch("backend.src.agents.orchestrator.orchestrator.orchestrate", new=AsyncMock(return_value=fake_audit)), \
             patch("backend.src.agents.fixer.fixer.run_fixer", new=AsyncMock(return_value=fake_fix)), \
             patch("backend.src.services.chat_tools._run_non_html_code_fixer", new=AsyncMock(return_value=fake_non_html_fix)), \
             patch("backend.src.services.chat_tools._render_html_to_screenshot", new=AsyncMock(return_value="fake_b64")), \
             patch("backend.src.services.chat_tools._verify_layout_visually", new=AsyncMock(return_value={"layout_ok": True, "reasons": []})):
            out = json.loads(fix_and_zip_files({"files": files}))

        assert "download_url" in out
        assert "zip_filename" in out
        assert len(out["changes_summary"]) == 1
        assert out["total_files"] == 2

    def test_fix_no_files_no_cached_analysis_returns_error(self):
        with patch(
            "backend.src.services.last_analyzed_content_store.get_last_analyzed_content",
            return_value=("", ""),
        ):
            out = json.loads(fix_and_zip_files({}))
        assert "error" in out

    def test_fix_no_files_falls_back_to_last_analyzed_content(self):
        fake_audit = AgentResult(agent="orchestrator", success=True, data={"issues": []})
        fake_fix = AgentResult(agent="fixer", success=True, data={"fixed_html": "<html>fixed</html>", "changes_summary": ["fixed img alt"]})

        with patch(
            "backend.src.services.last_analyzed_content_store.get_last_analyzed_content",
            return_value=("<html><img src=x></html>", "https://example.com/page"),
        ), patch("backend.src.agents.orchestrator.orchestrator.orchestrate", new=AsyncMock(return_value=fake_audit)), \
           patch("backend.src.agents.fixer.fixer.run_fixer", new=AsyncMock(return_value=fake_fix)), \
           patch("backend.src.services.chat_tools._render_html_to_screenshot", new=AsyncMock(return_value="fake_b64")), \
           patch("backend.src.services.chat_tools._verify_layout_visually", new=AsyncMock(return_value={"layout_ok": True, "reasons": []})):
            out = json.loads(fix_and_zip_files({}))

        assert "download_url" in out
        assert out["total_files"] == 1

    def test_fix_placeholder_file_content_falls_back_to_cache_instead_of_using_it(self):
        """Achado real (validação E2E completa, 2026-08-10): mesmo instruído a
        OMITIR `files` quando não tem o HTML, o modelo às vezes inventa um
        único arquivo com conteúdo implausivelmente curto (ex.: "---") em vez
        de omitir. Isso precisa cair no mesmo fallback do cache, não tentar
        "corrigir" o placeholder como se fosse HTML de verdade."""
        fake_audit = AgentResult(agent="orchestrator", success=True, data={"issues": []})
        fake_fix = AgentResult(agent="fixer", success=True, data={"fixed_html": "<html>fixed</html>", "changes_summary": []})

        with patch(
            "backend.src.services.last_analyzed_content_store.get_last_analyzed_content",
            return_value=("<html><img src=x></html>", "https://example.com/page"),
        ) as mock_cache, patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate", new=AsyncMock(return_value=fake_audit)
        ), patch("backend.src.agents.fixer.fixer.run_fixer", new=AsyncMock(return_value=fake_fix)), \
           patch("backend.src.services.chat_tools._render_html_to_screenshot", new=AsyncMock(return_value="fake_b64")), \
           patch("backend.src.services.chat_tools._verify_layout_visually", new=AsyncMock(return_value={"layout_ok": True, "reasons": []})):
            out = json.loads(fix_and_zip_files({"files": [{"path": "mars.html", "content": "---"}]}))

        assert "download_url" in out
        assert out["total_files"] == 1
        # o cache foi consultado -- prova de que caiu no fallback em vez de
        # tentar "corrigir" o placeholder "---" como se fosse HTML de verdade
        mock_cache.assert_called_once()

    def test_fix_multiple_files_with_short_content_is_not_treated_as_placeholder(self):
        """A heurística só se aplica quando é UM único arquivo curto (o padrão
        real observado) -- um projeto de verdade com vários arquivos pequenos
        (ex.: configs) não deve ser descartado como placeholder."""
        fake_audit = AgentResult(agent="orchestrator", success=True, data={"issues": []})
        fake_fix = AgentResult(agent="fixer", success=True, data={"fixed_html": "<html>x</html>", "changes_summary": []})
        fake_non_html_fix = {"path": "a.css", "content": "a{}", "changes_summary": []}

        files = [
            {"path": "index.html", "content": "<html>x</html>"},
            {"path": "a.css", "content": "a{}"},
        ]
        with patch("backend.src.agents.orchestrator.orchestrator.orchestrate", new=AsyncMock(return_value=fake_audit)), \
             patch("backend.src.agents.fixer.fixer.run_fixer", new=AsyncMock(return_value=fake_fix)), \
             patch("backend.src.services.chat_tools._run_non_html_code_fixer", new=AsyncMock(return_value=fake_non_html_fix)), \
             patch("backend.src.services.chat_tools._render_html_to_screenshot", new=AsyncMock(return_value="fake_b64")), \
             patch("backend.src.services.chat_tools._verify_layout_visually", new=AsyncMock(return_value={"layout_ok": True, "reasons": []})):
            out = json.loads(fix_and_zip_files({"files": files}))

        assert out["total_files"] == 2

    def test_fix_short_but_real_markup_is_never_treated_as_placeholder(self):
        """Achado real corrigido na própria correção acima: a heurística
        (tamanho sozinho) gerou falso positivo contra um teste existente que
        passa um único arquivo genuinamente curto mas com marcação real
        ('<p>x</p>', 8 chars). Só deve cair no fallback quando é curto E não
        tem NENHUM caractere de marcação/código (<>{};) -- '<p>x</p>' tem."""
        fake_audit = AgentResult(agent="orchestrator", success=True, data={"issues": []})
        fake_fix = AgentResult(agent="fixer", success=True, data={"fixed_html": "<p>x</p>", "changes_summary": []})

        with patch(
            "backend.src.services.last_analyzed_content_store.get_last_analyzed_content",
        ) as mock_cache, patch(
            "backend.src.agents.orchestrator.orchestrator.orchestrate", new=AsyncMock(return_value=fake_audit)
        ), patch("backend.src.agents.fixer.fixer.run_fixer", new=AsyncMock(return_value=fake_fix)), \
           patch("backend.src.services.chat_tools._render_html_to_screenshot", new=AsyncMock(return_value="fake_b64")), \
           patch("backend.src.services.chat_tools._verify_layout_visually", new=AsyncMock(return_value={"layout_ok": True, "reasons": []})):
            out = json.loads(fix_and_zip_files({"files": [{"path": "snippet.html", "content": "<p>x</p>"}]}))

        assert out["total_files"] == 1
        mock_cache.assert_not_called()  # NÃO deve ter caído no fallback


class TestGenerateChecklistTool:
    """Achado real corrigido: o chat gerava o checklist como texto solto lido
    do JSON cru da última análise, reimplementando o que o ChecklistAgent
    dedicado (backend/src/agents/checklist/checklist.py) já faz de forma
    estruturada. generate_checklist agora chama o agente real."""

    def test_generate_checklist_no_analysis_returns_error(self):
        with patch(
            "backend.src.services.last_analysis_store.get_last_analysis",
            return_value=([], ""),
        ):
            out = json.loads(generate_checklist({}))
        assert "error" in out

    def test_generate_checklist_calls_real_checklist_agent(self):
        fake_issues = [{
            "id": "1", "guideline": "WCAG 2.2", "criterion": "1.1.1 Non-text Content",
            "severity": "critical", "element": "img", "description": "x", "suggestion": "y",
        }]
        fake_checklist_result = AgentResult(
            agent="checklist",
            success=True,
            data={"checklist": [{
                "id": "chk-1", "criterion": "1.1.1 Non-text Content", "guideline": "WCAG 2.2",
                "status": "fail", "priority": "critical", "notes": "Adicione alt.",
            }]},
        )
        with patch(
            "backend.src.services.last_analysis_store.get_last_analysis",
            return_value=(fake_issues, "https://example.com"),
        ), patch(
            "backend.src.services.last_analyzed_content_store.get_last_analyzed_content",
            return_value=("<html></html>", "https://example.com"),
        ), patch(
            "backend.src.agents.checklist.checklist.run_checklist",
            new=AsyncMock(return_value=fake_checklist_result),
        ) as mock_run_checklist:
            out = json.loads(generate_checklist({}))

        assert out["url"] == "https://example.com"
        assert len(out["checklist"]) == 1
        assert out["checklist"][0]["status"] == "fail"
        # HTML da última página analisada (Fix #1) foi passado pro agente --
        # habilita as perguntas de verificação manual (alt text, labels, etc).
        _, kwargs = mock_run_checklist.call_args
        assert kwargs.get("html_content") == "<html></html>"

    def test_generate_checklist_agent_failure_returns_error(self):
        fake_issues = [{
            "id": "1", "guideline": "WCAG 2.2", "criterion": "1.1.1 Non-text Content",
            "severity": "critical", "element": "img", "description": "x", "suggestion": "y",
        }]
        with patch(
            "backend.src.services.last_analysis_store.get_last_analysis",
            return_value=(fake_issues, "https://example.com"),
        ), patch(
            "backend.src.services.last_analyzed_content_store.get_last_analyzed_content",
            return_value=("", ""),
        ), patch(
            "backend.src.agents.checklist.checklist.run_checklist",
            new=AsyncMock(return_value=AgentResult(agent="checklist", success=False, data={}, error="boom")),
        ):
            out = json.loads(generate_checklist({}))

        assert "error" in out
        assert "boom" in out["error"]


class TestExportChecklistPdfTool:
    def test_export_checklist_pdf_returns_download_link(self):
        out = json.loads(export_checklist_pdf_tool({}))
        assert out["download_url"].endswith("/export/last_checklist_pdf")


class TestRunRemoteTestToolLocationDecision:
    """Achado real: rodar Cypress/Selenium local ou na nuvem é decisão do
    USUÁRIO, não do modelo -- sem fallback silencioso de um pro outro. O
    handler bloqueia a chamada até `location` vir preenchido."""

    def test_cypress_without_location_asks_to_clarify_instead_of_running(self):
        with patch("backend.src.services.remote_runners.run_remote_cypress_simulation") as mock_run:
            out = json.loads(run_remote_test_tool({"runner": "cypress", "target_url": "https://example.com"}))

        assert "error" in out
        assert "clarify" in out["error"]
        mock_run.assert_not_called()

    def test_selenium_without_location_asks_to_clarify_instead_of_running(self):
        with patch("backend.src.services.remote_runners.run_remote_selenium") as mock_run:
            out = json.loads(run_remote_test_tool({"runner": "selenium", "target_url": "https://example.com"}))

        assert "error" in out
        assert "clarify" in out["error"]
        mock_run.assert_not_called()

    def test_cypress_with_location_runs_immediately(self):
        with patch(
            "backend.src.services.remote_runners.run_remote_cypress_simulation",
            new=AsyncMock(return_value={"status": "ok", "runner": "cypress_remote"}),
        ) as mock_run:
            out = json.loads(run_remote_test_tool({
                "runner": "cypress", "target_url": "https://example.com", "location": "cloud",
            }))

        assert out["status"] == "ok"
        mock_run.assert_called_once()

    def test_postman_never_needs_location(self):
        with patch(
            "backend.src.services.remote_runners.run_remote_postman_contract",
            new=AsyncMock(return_value={"status": "ok", "runner": "postman_remote"}),
        ) as mock_run:
            out = json.loads(run_remote_test_tool({"runner": "postman", "target_url": "https://api.example.com"}))

        assert out["status"] == "ok"
        mock_run.assert_called_once()

    def test_install_local_is_accepted_as_a_valid_location(self):
        """Achado real (pedido do usuário): quando o Cypress/Selenium local não
        está instalado, a IA deve poder oferecer instalar de verdade -- não só
        recusar. `location='install_local'` passa direto, sem pedir clarify de novo."""
        with patch(
            "backend.src.services.remote_runners.run_remote_cypress_simulation",
            new=AsyncMock(return_value={"status": "ok", "runner": "cypress_local"}),
        ) as mock_run:
            out = json.loads(run_remote_test_tool({
                "runner": "cypress", "target_url": "https://example.com", "location": "install_local",
            }))

        assert out["status"] == "ok"
        mock_run.assert_called_once()
        _, kwargs = mock_run.call_args
        assert kwargs.get("location") == "install_local"

    def test_remember_choice_saves_standing_consent_for_the_runner(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store._sessions.clear()
        with patch(
            "backend.src.services.remote_runners.run_remote_selenium",
            new=AsyncMock(return_value={"status": "ok", "runner": "selenium_local"}),
        ):
            run_remote_test_tool({
                "runner": "selenium", "target_url": "https://example.com",
                "location": "local", "remember_choice": True,
            })

        assert local_exec_consent_store.has_standing_consent("selenium") is True
        local_exec_consent_store._sessions.clear()

    def test_remember_choice_false_does_not_save_consent(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store._sessions.clear()
        with patch(
            "backend.src.services.remote_runners.run_remote_cypress_simulation",
            new=AsyncMock(return_value={"status": "ok", "runner": "cypress_local"}),
        ):
            run_remote_test_tool({
                "runner": "cypress", "target_url": "https://example.com", "location": "local",
            })

        assert local_exec_consent_store.has_standing_consent("cypress") is False


class TestRunRemoteTestToolFeedsDeliverables:
    """Achado real (pedido do usuário, 2026-08-11): sem isso, um teste real de
    Cypress/Selenium detectava problemas de verdade mas a IA não conseguia
    depois gerar planilha/checklist/PDF/VPAT a partir DESSES resultados --
    esses geradores só liam de last_analysis_store, que só era populado por
    analyze_page/analyze_site. Agora um teste remoto real com sucesso também
    alimenta o mesmo cache (issues + HTML renderizado da página testada)."""

    _RAW_VIOLATIONS = [
        {
            "id": "image-alt", "impact": "critical", "description": "Images must have alt text",
            "help": "Images must have alt text", "helpUrl": "https://dequeuniversity.com/rules/axe/4.13/image-alt",
            "tags": ["wcag2a", "wcag111"],
            "nodes": [{"target": ["img.hero"], "html": "<img class=hero>", "failureSummary": "Fix: add alt"}],
        },
    ]

    def _clear_last_analysis_store(self):
        import os

        from backend.src.services import last_analysis_store
        last_analysis_store._sessions.clear()
        # set_last_analysis também persiste em disco (get_last_analysis
        # recarrega de lá quando a memória está vazia) -- sem apagar o
        # arquivo, um teste anterior nesta classe vazava pro próximo mesmo
        # com _sessions.clear(), porque a sessão default é a mesma em todos.
        cache_path = last_analysis_store.get_cache_filepath()
        if os.path.exists(cache_path):
            os.remove(cache_path)

    def setup_method(self):
        self._clear_last_analysis_store()

    def teardown_method(self):
        self._clear_last_analysis_store()

    def test_successful_cypress_run_populates_last_analysis_store(self):
        from backend.src.services import last_analysis_store
        with patch(
            "backend.src.services.remote_runners.run_remote_cypress_simulation",
            new=AsyncMock(return_value={
                "status": "ok", "runner": "cypress_remote", "total_violations": 1,
                "_raw_violations": self._RAW_VIOLATIONS,
            }),
        ), patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(return_value=("<html><img class='hero'></html>", None)),
        ), patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<img class='hero'>",
        ):
            run_remote_test_tool({
                "runner": "cypress", "target_url": "https://example.com", "location": "cloud",
            })

        issues, url = last_analysis_store.get_last_analysis()
        assert url == "https://example.com"
        assert len(issues) == 1
        assert issues[0]["criterion"] == "image-alt"
        assert issues[0]["severity"] == "critical"

    def test_raw_violations_field_never_leaks_into_the_response_sent_to_the_model(self):
        with patch(
            "backend.src.services.remote_runners.run_remote_cypress_simulation",
            new=AsyncMock(return_value={
                "status": "ok", "runner": "cypress_remote", "total_violations": 1,
                "_raw_violations": self._RAW_VIOLATIONS,
            }),
        ), patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(return_value=("<html></html>", None)),
        ), patch(
            "backend.src.routes.analyze._extract_semantic_html", return_value="<html></html>",
        ):
            out = json.loads(run_remote_test_tool({
                "runner": "cypress", "target_url": "https://example.com", "location": "cloud",
            }))

        assert "_raw_violations" not in out

    def test_postman_result_without_raw_violations_does_not_touch_the_store(self):
        from backend.src.services import last_analysis_store
        last_analysis_store._sessions.clear()
        with patch(
            "backend.src.services.remote_runners.run_remote_postman_contract",
            new=AsyncMock(return_value={"status": "ok", "runner": "postman_remote"}),
        ):
            run_remote_test_tool({"runner": "postman", "target_url": "https://api.example.com"})

        issues, _url = last_analysis_store.get_last_analysis()
        assert issues == []

    def test_caching_failure_never_breaks_the_test_result_itself(self):
        """Se buscar/cachear a página falhar (rede, timeout), o resultado real
        do teste (que já rodou com sucesso) ainda deve chegar ao usuário --
        essa é só uma conveniência extra, nunca deve derrubar o turno."""
        with patch(
            "backend.src.services.remote_runners.run_remote_cypress_simulation",
            new=AsyncMock(return_value={
                "status": "ok", "runner": "cypress_remote", "total_violations": 1,
                "_raw_violations": self._RAW_VIOLATIONS,
            }),
        ), patch(
            "backend.src.services.browser.fetch_rendered_html_and_screenshot",
            new=AsyncMock(side_effect=RuntimeError("network down")),
        ):
            out = json.loads(run_remote_test_tool({
                "runner": "cypress", "target_url": "https://example.com", "location": "cloud",
            }))

        assert out["status"] == "ok"
        assert out["total_violations"] == 1


class TestLocalExecConsentStore:
    def setup_method(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store._sessions.clear()

    def teardown_method(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store._sessions.clear()

    def test_once_level_never_persists(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store.set_consent("cypress", "once", session_id="s1")
        assert local_exec_consent_store.has_standing_consent("cypress", session_id="s1") is False

    def test_always_level_persists_per_session(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store.set_consent("cypress", "always", session_id="s1")
        assert local_exec_consent_store.has_standing_consent("cypress", session_id="s1") is True
        assert local_exec_consent_store.has_standing_consent("cypress", session_id="s2") is False

    def test_consent_is_scoped_per_runner(self):
        from backend.src.services import local_exec_consent_store
        local_exec_consent_store.set_consent("cypress", "always", session_id="s1")
        assert local_exec_consent_store.has_standing_consent("selenium", session_id="s1") is False


class TestReadXlsxStructure:
    """`_read_xlsx_structure` alimenta run_excel_accessibility -- extrai
    estrutura real (openpyxl), não conteúdo célula a célula."""

    def _make_xlsx_bytes(self):
        import io as _io

        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Vendas"
        ws.append(["Mês", "Total"])
        ws.append(["Jan", 120])
        ws.merge_cells("A1:B1")
        buf = _io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        return buf

    def test_extracts_real_sheet_structure(self):
        from backend.src.services.chat_tools import _read_xlsx_structure
        summary = _read_xlsx_structure(self._make_xlsx_bytes())
        assert "Vendas" in summary
        assert "A1:B1" in summary  # célula mesclada detectada
        assert "cabeçalho" in summary.lower()

    def test_read_error_returns_diagnostic_string_not_exception(self):
        import io as _io

        from backend.src.services.chat_tools import _read_xlsx_structure
        summary = _read_xlsx_structure(_io.BytesIO(b"not a real xlsx"))
        assert "Erro" in summary


class TestAnalyzeDocumentTool:
    """Achado real: PDF/XLSX enviados já chegam como TEXTO extraído no
    contexto do modelo (chat_runtime.py pré-processa antes do turno) -- por
    isso analyze_document recebe `document_text`, nunca bytes/base64 (mesma
    lição do Fix #1 em fix_and_zip_files)."""

    def test_missing_params_returns_error(self):
        from backend.src.services.chat_tools import analyze_document
        out = json.loads(analyze_document({"filename": "relatorio.pdf"}))
        assert "error" in out

    def test_unsupported_extension_returns_error(self):
        from backend.src.services.chat_tools import analyze_document
        out = json.loads(analyze_document({"filename": "foo.docx", "document_text": "algum texto"}))
        assert "error" in out

    def test_pdf_dispatches_to_pdf_accessibility_agent_and_caches_result(self):
        from backend.src.services.chat_tools import analyze_document
        fake_result = AgentResult(
            agent="pdf_accessibility",
            success=True,
            data={"issues": [{
                "id": "pdf-1", "guideline": "WCAG 2.2", "criterion": "1.3.1 Info and Relationships",
                "severity": "critical", "element": "document", "description": "sem tags", "suggestion": "tagear",
            }]},
        )
        with patch(
            "backend.src.agents.pdf_accessibility.pdf_accessibility.run_pdf_accessibility",
            new=AsyncMock(return_value=fake_result),
        ), patch("backend.src.services.last_analysis_store.set_last_analysis") as mock_cache:
            out = json.loads(analyze_document({
                "filename": "relatorio.pdf", "document_text": "[Texto extraído do PDF]\n...",
            }))

        assert out["filename"] == "relatorio.pdf"
        assert out["document_type"] == "pdf"
        mock_cache.assert_called_once()

    def test_xlsx_dispatches_to_excel_accessibility_agent(self):
        from backend.src.services.chat_tools import analyze_document
        fake_result = AgentResult(agent="excel_accessibility", success=True, data={"issues": []})
        with patch(
            "backend.src.agents.excel_accessibility.excel_accessibility.run_excel_accessibility",
            new=AsyncMock(return_value=fake_result),
        ) as mock_run, patch("backend.src.services.last_analysis_store.set_last_analysis"):
            out = json.loads(analyze_document({
                "filename": "planilha.xlsx", "document_text": "[Estrutura extraída do XLSX]\n...",
            }))

        assert out["document_type"] == "xlsx"
        mock_run.assert_called_once()

    def test_agent_failure_returns_clear_error(self):
        from backend.src.services.chat_tools import analyze_document
        fake_result = AgentResult(agent="pdf_accessibility", success=False, data={}, error="boom")
        with patch(
            "backend.src.agents.pdf_accessibility.pdf_accessibility.run_pdf_accessibility",
            new=AsyncMock(return_value=fake_result),
        ):
            out = json.loads(analyze_document({"filename": "x.pdf", "document_text": "texto"}))
        assert "error" in out
        assert "boom" in out["error"]


class TestSearchTools:
    def test_tavily_search_no_query_returns_error(self):
        from backend.src.services.chat_tools import tavily_search
        out = json.loads(tavily_search({}))
        assert "error" in out

    def test_exa_search_no_query_returns_error(self):
        from backend.src.services.chat_tools import exa_search
        out = json.loads(exa_search({}))
        assert "error" in out

    def test_tavily_search_calls_provider(self):
        from unittest.mock import MagicMock

        from backend.src.services.chat_tools import tavily_search

        mock_settings = MagicMock()
        mock_settings.tavily_api_key = "test-tavily-key"

        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.json.return_value = {
            "results": [{"title": "Tavily Result", "url": "http://tavily", "content": "Some info"}]
        }

        with patch("backend.src.config.settings.get_settings", return_value=mock_settings), \
             patch("httpx.post", return_value=mock_response) as mock_post:
            out = json.loads(tavily_search({"query": "wcag"}))

        assert out["success"] is True
        assert out["data"]["web"][0]["title"] == "Tavily Result"
        mock_post.assert_called_once()

    def test_exa_search_calls_provider(self):
        from unittest.mock import MagicMock

        from backend.src.services.chat_tools import exa_search

        mock_settings = MagicMock()
        mock_settings.exa_api_key = "test-exa-key"

        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.json.return_value = {
            "results": [{"title": "Exa Result", "url": "http://exa", "text": "Some text"}]
        }

        with patch("backend.src.config.settings.get_settings", return_value=mock_settings), \
             patch("httpx.post", return_value=mock_response) as mock_post:
            out = json.loads(exa_search({"query": "wcag"}))

        assert out["success"] is True
        assert out["data"]["web"][0]["title"] == "Exa Result"
        mock_post.assert_called_once()

    def test_tavily_search_missing_key_reports_not_configured_not_attribute_error(self, monkeypatch):
        # Regression test: Settings used to have no `tavily_api_key` field at all, so
        # `settings.tavily_api_key` raised AttributeError before the os.environ fallback
        # was ever reached -- broken regardless of whether TAVILY_API_KEY was set. Using
        # a real Settings instance (not a MagicMock, which fabricates any attribute and
        # would hide this exact bug) proves the field really exists now.
        from backend.src.config.settings import Settings
        from backend.src.services.chat_tools import tavily_search

        monkeypatch.delenv("TAVILY_API_KEY", raising=False)
        real_settings = Settings(secret_key="test", tavily_api_key=None)
        with patch("backend.src.config.settings.get_settings", return_value=real_settings):
            out = json.loads(tavily_search({"query": "wcag"}))

        assert out.get("error") == "Tavily API key is not configured"

    def test_exa_search_missing_key_reports_not_configured_not_attribute_error(self, monkeypatch):
        from backend.src.config.settings import Settings
        from backend.src.services.chat_tools import exa_search

        monkeypatch.delenv("EXA_API_KEY", raising=False)
        real_settings = Settings(secret_key="test", exa_api_key=None)
        with patch("backend.src.config.settings.get_settings", return_value=real_settings):
            out = json.loads(exa_search({"query": "wcag"}))

        assert out.get("error") == "Exa API key is not configured"

    def test_tavily_search_reads_real_settings_field(self):
        from unittest.mock import MagicMock

        from backend.src.config.settings import Settings
        from backend.src.services.chat_tools import tavily_search

        real_settings = Settings(secret_key="test", tavily_api_key="real-tavily-key")
        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.json.return_value = {"results": []}

        with patch("backend.src.config.settings.get_settings", return_value=real_settings), \
             patch("httpx.post", return_value=mock_response) as mock_post:
            out = json.loads(tavily_search({"query": "wcag"}))

        assert out["success"] is True
        assert mock_post.call_args.kwargs["json"]["api_key"] == "real-tavily-key"

    def test_read_pptx_text(self):
        import io

        from pptx import Presentation

        from backend.src.services.chat_tools import _read_pptx_text

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        if slide.shapes.title:
            slide.shapes.title.text = "Slide 1 Title"
        body = slide.placeholders[1]
        body.text = "This is bullet point 1"

        # Add a table
        from pptx.util import Inches
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Val 1"
        table.cell(1, 1).text = "Val 2"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        text = _read_pptx_text(buf)
        assert "Slide 1 Title" in text
        assert "This is bullet point 1" in text
        assert "Header 1 | Header 2" in text
        assert "Val 1 | Val 2" in text

    @pytest.mark.asyncio
    @patch("backend.src.services.llm_client.call_llm")
    async def test_fix_pptx_document_success(self, mock_call_llm):
        import io

        from pptx import Presentation

        from backend.src.services.chat_tools import _fix_pptx_document

        llm_response = {
            "metadata": {
                "title": "Acessibilidade de Apresentacoes",
                "language": "pt-BR"
            },
            "slides": [
                {
                    "title": "Introducao",
                    "layout": "title_content",
                    "bullets": ["Marcador 1", "Marcador 2"],
                    "table": {
                        "headers": ["A", "B"],
                        "rows": [["1", "2"]]
                    }
                }
            ]
        }
        mock_call_llm.return_value = json.dumps(llm_response)

        fixed_bytes = await _fix_pptx_document("presentation.pptx", "old content", None)
        assert isinstance(fixed_bytes, bytes)

        prs = Presentation(io.BytesIO(fixed_bytes))
        assert len(prs.slides) == 1
        assert prs.core_properties.title == "Acessibilidade de Apresentacoes"
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Introducao"

        tables = [shape.table for shape in slide.shapes if shape.has_table]
        assert len(tables) == 1
        assert tables[0].cell(0, 0).text == "A"
        assert tables[0].cell(1, 0).text == "1"


# ── usage de tokens no evento `done` ─────────────────────────────────────────


class _UsageAgent:
    """Agente falso que devolve a contagem de tokens, como os 4 providers reais."""

    usage: dict[str, int] | None = None

    def __init__(self, **kwargs):
        pass

    def run_conversation(self, user_message):
        result: dict[str, Any] = {"final_response": "ok", "failed": False}
        if type(self).usage is not None:
            result["usage"] = type(self).usage
        return result


@pytest.mark.asyncio
async def test_done_event_carries_the_turn_token_usage():
    """O backend já somava o usage dos 4 caminhos de provider, mas não o enviava
    para o frontend -- o evento `done` só levava o texto final, e a UI não tinha
    como mostrar o custo do turno."""
    _UsageAgent.usage = {"input_tokens": 800, "output_tokens": 434, "total_tokens": 1234}
    with patch("backend.src.services.chat_runtime.AIAgent", new=_UsageAgent):
        events = [ev async for ev in stream_chat("olá")]

    done = events[-1]
    assert done["type"] == "done"
    assert done["usage"] == {"input_tokens": 800, "output_tokens": 434, "total_tokens": 1234}


@pytest.mark.asyncio
async def test_done_event_omits_usage_when_the_provider_did_not_count():
    """Sem contagem não se inventa um zero: o campo fica ausente e a UI esconde
    o indicador em vez de mostrar '0 tokens'."""
    _UsageAgent.usage = None
    with patch("backend.src.services.chat_runtime.AIAgent", new=_UsageAgent):
        events = [ev async for ev in stream_chat("olá")]

    done = events[-1]
    assert done["type"] == "done"
    assert "usage" not in done


class TestCrossBrowserAndPlaywrightInstallFlow:
    """Fluxo real pedido pelo usuário (2026-08-12): analisar/testar uma
    página, ver quais motores de navegador rodaram, e -- se algum estiver
    faltando -- a IA sugerir instalar o Playwright local; se já estiver
    instalado, avisar isso em vez de fingir que instalou algo."""

    def test_missing_engine_produces_install_suggestion(self):
        from backend.src.services.chat_tools import run_cross_browser_test_tool

        fake_audit = {
            "url": "https://example.com",
            "per_engine": {
                "chromium": {"success": True, "results": {"violations": [], "incomplete": []}},
                "firefox": {
                    "success": False,
                    "error": "Executable doesn't exist... Please run the following command to download new browsers:\n    playwright install",
                },
                "webkit": {"success": True, "results": {"violations": [], "incomplete": []}},
            },
        }
        with patch(
            "backend.src.services.browser.run_axe_core_cross_browser_audit",
            new=AsyncMock(return_value=fake_audit),
        ):
            raw = run_cross_browser_test_tool({"target_url": "https://example.com"})
        result = json.loads(raw)
        assert "firefox" in result["install_suggestion"]
        assert "install_playwright_browsers" in result["install_suggestion"]

    def test_all_engines_present_has_no_install_suggestion(self):
        from backend.src.services.chat_tools import run_cross_browser_test_tool

        fake_audit = {
            "url": "https://example.com",
            "per_engine": {
                engine: {"success": True, "results": {"violations": [], "incomplete": []}}
                for engine in ("chromium", "firefox", "webkit")
            },
        }
        with patch(
            "backend.src.services.browser.run_axe_core_cross_browser_audit",
            new=AsyncMock(return_value=fake_audit),
        ):
            raw = run_cross_browser_test_tool({"target_url": "https://example.com"})
        result = json.loads(raw)
        assert "install_suggestion" not in result
        assert result["status"] == "ok"
        assert result["engines_succeeded"] == ["chromium", "firefox", "webkit"]

    def test_install_playwright_browsers_reports_already_installed(self):
        from backend.src.services.chat_tools import install_playwright_browsers_tool

        fake_proc = MagicMock(returncode=0, stdout="", stderr="")
        with patch("subprocess.run", return_value=fake_proc):
            raw = install_playwright_browsers_tool({"pre_exec_msg": "instalando..."})
        result = json.loads(raw)
        assert result["status"] == "ok"
        assert result["already_installed"] is True
        assert "já estavam instalados" in result["message"]

    def test_install_playwright_browsers_reports_freshly_installed(self):
        from backend.src.services.chat_tools import install_playwright_browsers_tool

        fake_proc = MagicMock(returncode=0, stdout="Downloading Firefox 132.0 (playwright build v1466)...\n", stderr="")
        with patch("subprocess.run", return_value=fake_proc):
            raw = install_playwright_browsers_tool({"pre_exec_msg": "instalando..."})
        result = json.loads(raw)
        assert result["status"] == "ok"
        assert result["already_installed"] is False
        assert "instalados com sucesso agora" in result["message"]

    def test_install_playwright_browsers_reports_failure(self):
        from backend.src.services.chat_tools import install_playwright_browsers_tool

        fake_proc = MagicMock(returncode=1, stdout="", stderr="network error")
        with patch("subprocess.run", return_value=fake_proc):
            raw = install_playwright_browsers_tool({"pre_exec_msg": "instalando..."})
        result = json.loads(raw)
        assert result["status"] == "error"
        assert "network error" in result["error"]

    def test_install_playwright_browsers_requires_approval(self):
        """Baixa software real no disco -- precisa estar marcado como ação
        que exige aprovação explícita do usuário no registry de tools."""
        from tools.registry import registry
        assert registry.tools["install_playwright_browsers"]["requires_approval"] is True
